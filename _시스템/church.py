# -*- coding: utf-8 -*-
"""○○교회 종합행정시스템.
교인DB(심방이력 누적) · 증명서 자동발급 · 출석/재정 · 주보 · 교안 · Excel 동기화.
명령으로 운용: python church.py <명령> ...  (도움말: python church.py help)"""
import os, sys, json, argparse, datetime, re

BASE=os.path.dirname(os.path.abspath(__file__))
ROOT=os.path.dirname(BASE)                          # 프로그램 루트
DB=os.path.join(BASE,"church_db.json")
def CAT(n):
    m={"01":"01 교회 기본·조직","02":"02 예배·성례·예식","03":"03 교육·훈련","04":"04 목양·셀·심방·새가족",
       "05":"05 선교","06":"06 재정·증명","07":"07 행사·홍보·자료","08":"08 대외·인사·시설","09":"09 설교"}
    ar=_C.get("아카이브루트")   # 설정 시: 실제 교회 원본 아카이브(0X. 이름)에 저장·연결
    if ar and os.path.isdir(ar):
        real=os.path.join(ar, n+". "+m[n][3:])   # 원본 폴더는 번호 뒤 점: "07. 행사·홍보·자료"
        try: os.makedirs(real,exist_ok=True); return real
        except Exception: pass
    d=os.path.join(ROOT,m[n]); os.makedirs(d,exist_ok=True); return d
def today(): return datetime.date.today().isoformat()
def _skip_scan(dp):
    """파일 검색에서 제외할 폴더 — 프로그램·파이썬 라이브러리 내부. 목사님 문서(01~09)만 뜨게."""
    segs=set(s.lower() for s in dp.replace("/","\\").split("\\"))
    return bool(segs & {"_시스템","_피닉스","python","__pycache__","site-packages","_관리대장","scripts"})
CONFIG=os.path.join(BASE,"church_config.json")
def _load_cfg():
    if os.path.exists(CONFIG):
        try: return json.load(open(CONFIG,encoding='utf-8'))
        except Exception: pass
    return {"교회명":"○○교회","담임":"담임목사"}
_C=_load_cfg(); CHURCH=_C.get("교회명","○○교회"); PASTOR=_C.get("담임","담임목사")
COPYRIGHT_NOTICE=("본 프로그램과 모든 참고자료의 저작권은 세움교회 김용원 목사에게 있습니다. "
    "교회 목회 목적의 사용은 자유이나 무단 판매·상업적 이용을 금합니다. "
    "성경 인용은 공개역본(KJV·WEB)·성구 주소만 사용했으며 저작권 있는 타 저작물을 포함하지 않습니다. "
    "본 자료는 목회를 돕는 참고 자료이며 교단 규정·법률·세무 등 세부 사항은 각자 확인·전문가 상담을 권합니다. "
    "본 프로그램은 '있는 그대로' 제공되며 사용에 따른 책임은 사용자에게 있습니다.")

# ───────── DB ─────────
def load():
    if os.path.exists(DB):
        return json.load(open(DB,encoding='utf-8'))
    return {"교회":{"이름":CHURCH,"담임":PASTOR},"교인":[],"출석":[],"재정":[],"_seq":0}
PHOENIX=os.path.join(BASE,"_피닉스")
def _phoenix_snapshot(db):
    """🔥 저장 때마다 시점 스냅샷 보관(최근 60개). 자료가 삭제돼도 되살릴 수 있음."""
    try:
        import datetime
        os.makedirs(PHOENIX,exist_ok=True)
        stamp=datetime.datetime.now().strftime("%Y%m%d_%H%M%S_%f")  # 마이크로초까지 → 1초내 다중저장도 유실 없음
        json.dump(db,open(os.path.join(PHOENIX,f"db_{stamp}.json"),'w',encoding='utf-8'),ensure_ascii=False,indent=2)
        snaps=sorted(f for f in os.listdir(PHOENIX) if f.startswith("db_") and f.endswith(".json"))
        for old in snaps[:-60]:
            try: os.remove(os.path.join(PHOENIX,old))
            except Exception: pass
    except Exception: pass
def save(db):
    import shutil, tempfile
    try:  # 롤링 백업(직전 상태 보존) — 업데이트·오작동 대비
        if os.path.exists(DB): shutil.copy2(DB, DB+".bak")
    except Exception: pass
    # 원자적 저장: 임시파일에 먼저 쓰고 os.replace로 교체 → 쓰기 도중 중단·동시쓰기에도 DB 손상/유실 없음
    d=os.path.dirname(DB) or "."
    fd,tmp=tempfile.mkstemp(dir=d,suffix=".tmp"); os.close(fd)
    try:
        with open(tmp,'w',encoding='utf-8') as f:
            json.dump(db,f,ensure_ascii=False,indent=2)
        os.replace(tmp,DB)
    finally:
        try:
            if os.path.exists(tmp): os.remove(tmp)
        except Exception: pass
    _phoenix_snapshot(db)   # 🔥 피닉스 스냅샷
    try: _ext_copy(DB)      # D·USB 자동 이중저장(설정 '백업폴더')
    except Exception: pass
def find(db,name):
    if not name: return []
    hit=[m for m in db["교인"] if (m.get("이름") or "")==name]
    if not hit: hit=[m for m in db["교인"] if name in (m.get("이름") or "")]
    return hit
def _intarg(val, default, name="숫자"):
    """--days/--ahead 등 정수 인자 안전 파싱: 숫자가 아니면 크래시 없이 기본값 사용하고 안내."""
    if val is None or val=="": return default
    try:
        return int(val)
    except (ValueError, TypeError):
        print(f"⚠ '{val}'은(는) 숫자가 아니어서 {name} 기본값({default})으로 진행합니다.")
        return default
def _pint(s):
    """argparse용 정수 파서 — 숫자 아니면 영문 대신 한글 안내로 오류."""
    try:
        return int(s)
    except (ValueError, TypeError):
        raise argparse.ArgumentTypeError(f"'{s}'은(는) 숫자가 아닙니다. 금액·인원은 숫자로 입력해 주세요(예: 50000).")

# ───────── docx 스타일 ─────────
from docx import Document
from docx.shared import Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH as AL, WD_LINE_SPACING
from docx.enum.table import WD_TABLE_ALIGNMENT as TA
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
NAVY=RGBColor(0x1B,0x2A,0x4A); GOLD=RGBColor(0xA6,0x7C,0x1E); MAROON=RGBColor(0x8B,0x2E,0x2E)
GRAY=RGBColor(0x40,0x40,0x40); LG=RGBColor(0x88,0x88,0x88); WHITE=RGBColor(0xFF,0xFF,0xFF)
SANS="맑은 고딕"; SERIF="바탕"
def setf(r,n):
    r.font.name=n; rpr=r.element.get_or_add_rPr()
    for k in ('w:eastAsia','w:ascii','w:hAnsi'): rpr.get_or_add_rFonts().set(qn(k),n)
def para(d,text="",size=11,color=None,bold=False,align=None,before=0,after=6,font=SERIF,italic=False):
    # 문서 폰트에서 □로 깨지는 이모지·결합기호(⃞ 등) 자동 제거 — 모든 생성 문서 안전
    text=re.sub(r'[\U0001F000-\U0001FAFF️⃐-⃿]\s?','',str(text)) if text else text
    p=d.add_paragraph(); r=p.add_run(text); r.font.size=Pt(size); r.bold=bold; r.italic=italic
    if color: r.font.color.rgb=color
    setf(r,font); pf=p.paragraph_format; pf.space_before=Pt(before); pf.space_after=Pt(after)
    if align: p.alignment=align
    return p
def hr(d,color="A67C1E",sz=8,before=2,after=8):
    p=d.add_paragraph(); p.paragraph_format.space_before=Pt(before); p.paragraph_format.space_after=Pt(after)
    pPr=p._p.get_or_add_pPr(); bd=OxmlElement('w:pBdr'); e=OxmlElement('w:bottom')
    e.set(qn('w:val'),'single'); e.set(qn('w:sz'),str(sz)); e.set(qn('w:space'),'1'); e.set(qn('w:color'),color); bd.append(e); pPr.append(bd)
def shade(c,fill):
    tcPr=c._tc.get_or_add_tcPr(); sh=OxmlElement('w:shd'); sh.set(qn('w:val'),'clear'); sh.set(qn('w:fill'),fill); tcPr.append(sh)
def ctext(c,text,size=10,color=None,bold=False,font=SANS,align=None):
    c.text=""; p=c.paragraphs[0]; r=p.add_run(str(text)); r.font.size=Pt(size); r.bold=bold
    if color: r.font.color.rgb=color
    setf(r,font)
    if align: p.alignment=align
def tborders(t,color="BBBBBB",sz=4):
    b=OxmlElement('w:tblBorders')
    for e in ('top','left','bottom','right','insideH','insideV'):
        x=OxmlElement(f'w:{e}'); x.set(qn('w:val'),'single'); x.set(qn('w:sz'),str(sz)); x.set(qn('w:color'),color); b.append(x)
    t._tbl.tblPr.append(b)
def newdoc(sermon=False):
    d=Document()
    if sermon:  # 설교문 표준 서식(기준 hwpx): 바탕·줄간격160%·여백 좌우30/위20/아래15mm
        for s in d.sections: s.top_margin=Cm(2.0); s.bottom_margin=Cm(1.5); s.left_margin=Cm(3.0); s.right_margin=Cm(3.0)
    else:
        for s in d.sections: s.top_margin=Cm(2.0); s.bottom_margin=Cm(2.0); s.left_margin=Cm(2.2); s.right_margin=Cm(2.2)
    st=d.styles['Normal']; st.font.size=Pt(11); st.font.name=SERIF
    st.element.get_or_add_rPr().get_or_add_rFonts().set(qn('w:eastAsia'),SERIF)
    if sermon:
        st.paragraph_format.line_spacing_rule=WD_LINE_SPACING.MULTIPLE; st.paragraph_format.line_spacing=1.6
    _orig=d.save
    def _save2(path,*ar,**kw):   # 저장 시 D·USB 자동 이중저장(설정 '백업폴더')
        r=_orig(path,*ar,**kw); _ext_copy(path); return r
    d.save=_save2
    return d
def sanit(n): return re.sub(r'[:\\/*?"<>|]','·',(n or ''))[:50]  # 파일명 길이 50자 컷 — 긴 제목→경로초과·저장실패(파일 소실) 방지
def _openfile(p):
    """생성한 파일을 기본 프로그램(한글·워드·엑셀)으로 바로 열기 — 열린 문서에서 즉시 인쇄(Ctrl+P·쪽 지정).
       웹 실행 시엔 부모(웹서버·바탕화면 세션)가 확실히 열도록 표식만 출력(자식 subprocess 열기는 불안정)."""
    q=os.path.normpath(p)
    if os.environ.get("CHURCH_WEB")=="1":
        print("__OPENFILE__"+q); return
    try: os.startfile(q)
    except Exception: pass
def _savedoc(d, out):
    """docx 저장 + 존재·크기 검증 + 생성 즉시 열기. 파일이 열려 있으면(PermissionError) 시각 붙여 새 이름으로 재시도(덮어쓰기 실패 방지)."""
    try:
        d.save(out)
    except PermissionError:
        base,ext=os.path.splitext(out); out=f"{base}_{datetime.datetime.now().strftime('%H%M%S')}{ext}"
        try: d.save(out)
        except Exception as e:
            print(f"✗ 문서 저장 실패({os.path.basename(out)}): {e}"); return False
    except Exception as e:
        print(f"✗ 문서 저장 실패({os.path.basename(out)}): {e}"); return False
    if not (os.path.exists(out) and os.path.getsize(out)>0):
        print(f"✗ 문서가 실제로 저장되지 않았습니다(제목이 너무 길거나 경로 문제): {out}"); return False
    _openfile(out)
    return True
def _also_hwp(docx_path):
    """docx를 한글(.hwp)로도 저장 — 한글(한컴오피스) 설치 시. 없으면 조용히 건너뜀."""
    try:
        import win32com.client as w, subprocess
        try: subprocess.run(["taskkill","/F","/IM","Hwp.exe"],capture_output=True)
        except Exception: pass
        h=w.gencache.EnsureDispatch("HWPFrame.HwpObject")
        try: h.RegisterModule("FilePathCheckDLL","FilePathCheckerModule")
        except Exception: pass
        h.SetMessageBoxMode(0x20000); h.Open(docx_path)
        out=os.path.splitext(docx_path)[0]+".hwp"; h.SaveAs(out,"HWP"); h.Quit()
        _ext_copy(out); return out
    except Exception: return None
def _ext_copy(path):
    """생성물(문서·악보·곡·자료)을 외부 백업폴더(D·USB)에도 폴더구조 보존해 복사. 설정 '백업폴더' 시 자동."""
    dest=_C.get("백업폴더")
    if not dest or not path or not os.path.exists(path): return
    try:
        import shutil; ap=os.path.abspath(path); rel=None
        for base in (ROOT, BASE):
            try:
                if os.path.commonpath([ap,base])==base: rel=os.path.relpath(ap,base); break
            except Exception: pass
        tgt=os.path.join(dest,rel) if rel else os.path.join(dest,os.path.basename(path))
        os.makedirs(os.path.dirname(tgt),exist_ok=True); shutil.copy2(path,tgt)
    except Exception: pass

# ───────── 명령: 교인 ─────────
def member_add(a):
    if not (a.name or "").strip(): print("⚠ 이름을 입력해 주세요."); return
    db=load()
    if find(db,a.name):   # ★이미 등록된 교인이면 조회로 보여줌(등록 & 조회 통합·중복 방지)
        print(f"ℹ '{a.name}' 님은 이미 등록되어 있어 아래에 정보를 보여드립니다.")
        print("  (같은 이름의 다른 분을 새로 등록하시려면 이름 뒤에 구분을 붙여 주세요 — 예: 김철수B)")
        print("─"*34)
        return member_show(a)
    db["_seq"]+=1
    fam=[]   # '이름:관계[:생년월일]' 쌍을 ;로 구분 (교우 등록 시 가족 함께 입력)
    for item in re.split(r'[;\n]', getattr(a,"family",None) or ""):
        parts=[p.strip() for p in item.split(":")]
        if not parts or not parts[0]: continue
        fam.append({"이름":parts[0],"관계":(parts[1] if len(parts)>1 else ""),"생년월일":(parts[2] if len(parts)>2 else ""),"직분":"","연락처":""})
    m={"id":db["_seq"],"이름":a.name,"성별":a.sex or "","생년월일":a.birth or "","연락처":a.tel or "",
       "주소":a.addr or "","직분":a.role or "","세례":a.baptism or "","세례일":a.baptismdate or "","결혼기념일":a.wedding or "","소속셀":a.cell or "","인도자":a.leader or "",
       "등록일":a.date or today(),"심방주기":a.cycle or "","가족":fam,"직분이력":([{"날짜":a.date or today(),"직분":a.role}] if a.role else []),
       "상태":"재적","심방이력":[],"기도제목누적":[],"메모":a.memo or ""}
    db["교인"].append(m); save(db)
    print(f"✅ 교우 등록: {a.name} (id={m['id']}) · 소속셀 {m['소속셀'] or '-'} · 등록일 {m['등록일']}"+(f" · 가족 {len(fam)}명" if fam else ""))

def member_list(a):
    db=load(); ms=db["교인"]
    if a.cell: ms=[m for m in ms if m["소속셀"]==a.cell]
    print(f"■ 교인 {len(ms)}명"+(f" (셀:{a.cell})" if a.cell else ""))
    for m in ms:
        print(f"  #{m['id']:>3} {m['이름']:<8} {m['직분'] or '':<6} {m['소속셀'] or '':<8} 심방 {len(m['심방이력'])}회 {m['연락처']}")

def member_show(a):
    db=load(); hit=find(db,a.name)
    if not hit: print(f"✗ '{a.name}' 없음"); return
    m=hit[0]
    print(f"═══ {m['이름']}  (#{m['id']}) ═══")
    for k in ["성별","생년월일","연락처","주소","직분","세례","소속셀","인도자","등록일","상태"]:
        if m.get(k): print(f"  {k}: {m[k]}")
    if m["가족"]: print("  가족: "+", ".join(f"{g.get('이름','')}({g.get('관계','')})" for g in m["가족"]))
    if m["메모"]: print(f"  메모: {m['메모']}")
    print(f"  심방 {len(m['심방이력'])}회 · 누적 기도제목 {len(m['기도제목누적'])}건")
def family_add(a):
    """가족 등록 — 교인에게 가족 구성원 추가(관계·생일·직분·연락처)."""
    if not a.name or not a.member: print("⚠ 교인 이름과 가족 이름을 입력해 주세요. (가족은 '교우 등록' 카드에서 함께 입력할 수 있습니다)"); return
    db=load(); hit=find(db,a.name)
    if not hit: print(f"✗ 교인 '{a.name}' 없음 — 먼저 '교우 등록'에서 등록해 주세요"); return
    m=hit[0]; m.setdefault("가족",[]).append({"이름":a.member or "","관계":a.relation or "","생년월일":a.birth or "","직분":a.role or "","연락처":a.tel or ""})
    save(db); print(f"✅ [{m['이름']}] 가족 추가: {a.member or ''}({a.relation or ''}) → 총 {len(m['가족'])}명")
def family_list(a):
    """가족 조회 — 한 교인의 가족 구성원."""
    db=load(); hit=find(db,a.name)
    if not hit: print(f"✗ 교인 '{a.name}' 없음"); return
    m=hit[0]; fam=m.get("가족",[])
    print(f"■ {m['이름']} 가족 {len(fam)}명")
    for g in fam: print(f"  · {g.get('이름','')} ({g.get('관계','')}) {g.get('생년월일','')} {g.get('직분','')} {g.get('연락처','')}")
    if not fam: print("  (가족 미등록 — '교우 등록' 카드에서 가족을 함께 입력해 주세요)")
def office_add(a):
    """직분 이력 — 임직·취임 등 직분 변동 기록(날짜·직분·비고). 현재 직분도 갱신."""
    if not a.name: print("⚠ 교인 이름을 입력해 주세요."); return
    db=load(); hit=find(db,a.name)
    if not hit: print(f"✗ 교인 '{a.name}' 없음"); return
    m=hit[0]; m.setdefault("직분이력",[]).append({"날짜":a.date or today(),"직분":a.role or "","비고":a.memo or ""})
    if a.role: m["직분"]=a.role
    save(db); print(f"✅ [{m['이름']}] 직분 이력: {a.date or today()} {a.role or ''} {a.memo or ''} · 현재 직분 {m['직분']}")
def office_list(a):
    """직분 이력·세례일 조회."""
    db=load(); hit=find(db,a.name)
    if not hit: print(f"✗ 교인 '{a.name}' 없음"); return
    m=hit[0]; oh=m.get("직분이력",[])
    print(f"■ {m['이름']} 직분 이력 {len(oh)}건 · 현재 {m.get('직분','')}")
    for o in oh: print(f"  · {o.get('날짜','')} {o.get('직분','')} {o.get('비고','')}")
    if m.get("세례일"): print(f"  세례일: {m['세례일']}")
    if m.get("심방주기"): print(f"  심방주기: {m['심방주기']}")
def group_add(a):
    """소그룹(셀/구역/목장) 등록 — 리더·모임요일·장소."""
    if not a.name: print("⚠ 소그룹 이름을 입력해 주세요."); return
    db=load(); db.setdefault("소그룹",[])
    db["소그룹"].append({"이름":a.name or "","리더":a.leader or "","부리더":a.subleader or "","요일":a.day or "","장소":a.place or "","메모":a.memo or "","보고":[]})
    save(db); print(f"✅ 소그룹 등록: {a.name or ''}(리더 {a.leader or '-'}) · {a.day or ''} {a.place or ''}")
def group_list(a):
    db=load(); gs=db.get("소그룹",[])
    print(f"■ 소그룹 {len(gs)}개")
    for g in gs:
        n=len([m for m in db["교인"] if m.get("소속셀")==g["이름"]])
        print(f"  · {g['이름']} — 리더 {g.get('리더','')} · 조원 {n}명 · {g.get('요일','')} {g.get('장소','')}")
    if not gs: print("  (없음 — '소그룹(셀) 등록'에서 등록해 주세요)")
def group_assign(a):
    """교인을 소그룹에 배정."""
    db=load(); hit=find(db,a.name)
    if not hit: print(f"✗ 교인 '{a.name}' 없음"); return
    m=hit[0]; m["소속셀"]=a.group or ""; save(db); print(f"✅ {m['이름']} → 소그룹 '{a.group or ''}' 배정")
def group_roster(a):
    """소그룹 명단(docx) — 조원·직분·연락처."""
    db=load(); members=[m for m in db["교인"] if m.get("소속셀")==a.name]
    g=next((x for x in db.get("소그룹",[]) if x["이름"]==a.name), {"이름":a.name})
    d=newdoc()
    para(d,f"{a.name} 명단",20,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,f"리더 {g.get('리더','')} · {g.get('요일','')} {g.get('장소','')} · {CHURCH}",10,GRAY,AL.CENTER,after=6,font=SANS); hr(d,"3C6E3C",12)
    t=d.add_table(rows=len(members)+1,cols=4)
    for j,h in enumerate(["이름","직분","연락처","비고"]): ctext(t.rows[0].cells[j],h,10,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],"3C6E3C")
    for i,m in enumerate(members,1):
        for j,v in enumerate([m["이름"],m.get("직분",""),m.get("연락처",""),""]):
            ctext(t.rows[i].cells[j],v,10,GRAY,font=SERIF)
    tborders(t,"BBBBBB",4)
    out=os.path.join(CAT("04"),f"[소그룹명단] {sanit(a.name)}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 소그룹 명단: {out} ({len(members)}명)")
def group_report(a):
    """주간 소그룹 보고 — 모임·참석·나눔·기도제목·결석자(누적)."""
    db=load(); g=next((x for x in db.get("소그룹",[]) if x["이름"]==a.name), None)
    if not g: print(f"✗ 소그룹 '{a.name}' 없음 — 먼저 '소그룹(셀) 등록'에서 등록해 주세요"); return
    g.setdefault("보고",[]).append({"날짜":a.date or today(),"참석":a.attend or "","나눔":a.note or "","기도제목":a.prayer or "","결석":a.absent or ""})
    save(db); print(f"✅ [{a.name}] 주간보고 기록: {a.date or today()} · 참석 {a.attend or ''} · 결석 {a.absent or ''}")
def train_add(a):
    """양육·제자훈련 기록 — 교인의 교육과정 등록·진도·수료(새가족반·제자훈련·성경공부 등)."""
    if not a.name: print("⚠ 교인 이름을 입력해 주세요. (양육·제자훈련 기록)"); return
    db=load(); hit=find(db,a.name)
    if not hit: print(f"✗ 교인 '{a.name}' 없음"); return
    m=hit[0]; m.setdefault("양육이력",[]).append({"과정":a.course or "","단계":a.stage or "","상태":a.status or "진행","날짜":a.date or today(),"비고":a.memo or ""})
    save(db); print(f"✅ [{m['이름']}] 양육: {a.course or ''} {a.stage or ''} [{a.status or '진행'}]")
def train_list(a):
    """양육 이력 — 교인별(--name) 또는 과정별 수강·수료 현황."""
    db=load()
    if a.name:
        hit=find(db,a.name)
        if not hit: print(f"✗ '{a.name}' 없음"); return
        m=hit[0]; h=m.get("양육이력",[])
        print(f"■ {m['이름']} 양육 이력 {len(h)}건")
        for t in h: print(f"  · {t.get('날짜','')} {t.get('과정','')} {t.get('단계','')} [{t.get('상태','')}]")
        if not h: print("  (없음 — '양육·제자훈련 기록'에서 기록해 주세요)")
    else:
        from collections import defaultdict
        by=defaultdict(list)
        for m in db["교인"]:
            for t in m.get("양육이력",[]):
                if not a.course or a.course in t.get("과정",""): by[t.get("과정","")].append((m["이름"],t.get("상태","")))
        print(f"■ 교육 과정 {len(by)}개")
        for c,people in by.items():
            done=sum(1 for _,s in people if s=="수료")
            print(f"  · {c}: 수강 {len(people)}명 · 수료 {done}명")
            for name,s in people: print(f"      - {name} [{s}]")
        if not by: print("  (없음 — '양육·제자훈련 기록'에서 기록해 주세요)")
def group_form(a):
    """셀·구역 주간 보고서 양식(인쇄용 빈 양식) — 셀리더이 손으로 작성용."""
    d=newdoc()
    para(d,f"{CHURCH} 셀·구역 주간 보고서",19,NAVY,True,AL.CENTER,after=2,font=SANS)
    if a.name: para(d,a.name,12,MAROON,True,AL.CENTER,after=4,font=SANS)
    hr(d,"3C6E3C",12)
    def frow(labels):
        t=d.add_table(rows=1,cols=len(labels)*2)
        for j,lab in enumerate(labels):
            lc=t.rows[0].cells[j*2]; vc=t.rows[0].cells[j*2+1]
            ctext(lc,lab,10,WHITE,True,SANS,AL.CENTER); shade(lc,"3C6E3C"); lc.width=Cm(3); vc.width=Cm(5.5)
        tborders(t,"BBBBBB",4)
    for labs in [["셀/구역","셀리더"],["모임일시","장소"],["참석 인원","결석 인원"]]: frow(labs)
    def area(title,n):
        para(d,title,11,RGBColor(0x3C,0x6E,0x3C),True,before=8,after=3,font=SANS)
        t=d.add_table(rows=n,cols=1)
        for r in t.rows: r.cells[0].width=Cm(17)
        tborders(t,"BBBBBB",4)
    area("나눈 말씀·나눔 내용",4); area("기도 제목",3); area("심방 요청",2); area("결석자 명단",2); area("특이사항·건의",2)
    para(d,f"셀리더 서명 :               · {CHURCH}",9,LG,align=AL.RIGHT,before=8,font=SANS)
    out=os.path.join(CAT("04"),f"[셀보고서양식] {sanit(a.name or '셀')}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 셀 주간보고서 양식: {out}")

# ───────── 명령: 심방(핵심) ─────────
def _visit_briefing(m):
    """심방 브리핑 출력(반복 방지) — 지난 말씀·기도제목·후속을 최근순으로."""
    H=m.get("심방이력",[])
    print(f"╔══ 심방 브리핑: {m['이름']} ({m.get('소속셀') or '-'}) ══╗")
    print(f"  직분 {m.get('직분') or '-'} · 연락처 {m.get('연락처') or '-'} · 주소 {m.get('주소') or '-'}")
    if not H: print("  ▸ 첫 심방입니다 (이력 없음)."); return
    print(f"  ▸ 지난 심방 {len(H)}회. 최근순:")
    for v in list(reversed(H))[:5]:
        print(f"   · {v.get('날짜','')} [{v.get('구분','')}] 말씀: {v.get('말씀') or '-'}")
        if v.get('내용'): print(f"       내용: {str(v['내용'])[:60]}")
        if v.get('기도제목'): print(f"       기도제목: {', '.join(v['기도제목'])}")
        if v.get('후속'): print(f"       후속: {v['후속']}")
    used=[v.get('말씀') for v in H if v.get('말씀')]
    print(f"  ⚠ 이미 나눈 말씀(중복 주의): {', '.join(used) if used else '없음'}")
def visit_add(a):
    db=load(); hit=find(db,a.name)
    if not hit: print(f"✗ '{a.name}' 없음 — 먼저 '교우 등록'에서 등록해 주세요"); return
    m=hit[0]
    _visit_briefing(m)   # ★ 기록 전에 항상 브리핑(지난 말씀·기도제목 반복 방지)
    has=any([(getattr(a,'word','') or '').strip(),(getattr(a,'note','') or '').strip(),
             (getattr(a,'prayer','') or '').strip(),(getattr(a,'followup','') or '').strip()])
    if not has:
        print("  ─────────────────────────")
        print("  ※ 브리핑만 보셨습니다. 심방 다녀오신 뒤 '전한 말씀·기도제목·나눈 내용'을 채워 다시 실행하면 기록됩니다.")
        return
    prayers=[p.strip() for p in (a.prayer or "").split(";") if p.strip()]
    v={"날짜":a.date or today(),"심방자":a.by or PASTOR,"구분":a.kind or "정기",
       "말씀":a.word or "","내용":a.note or "","기도제목":prayers,"후속":a.followup or ""}
    m["심방이력"].append(v)
    for p in prayers:
        if p not in m["기도제목누적"]: m["기도제목누적"].append(p)
    save(db)
    print("  ─────────────────────────")
    print(f"✅ 심방 기록 완료: {m['이름']} · {v['날짜']} · 말씀 '{v['말씀'] or '-'}' · 기도제목 {len(prayers)}건 누적")

def visit_daesim(a):
    """🏘️ 춘계·추계 대심방 현황 — 올해 대심방(구분에 '대심방' 포함) 완료 세대와 미심방 세대·진행률(전교인 심방 점검)."""
    db=load(); mem=db.get("교인",[])
    yr=(getattr(a,'year','') or str(datetime.date.today().year)).strip()
    done=[]; notyet=[]
    for m in mem:
        if (m.get("상태") or "재적") not in ("재적","재적중",""): continue
        hit=[v for v in m.get("심방이력",[]) if "대심방" in str(v.get("구분","")) and str(v.get("날짜","")).startswith(yr)]
        (done if hit else notyet).append((m,hit[-1] if hit else None))
    total=len(done)+len(notyet); pct=round(len(done)/total*100) if total else 0
    d=newdoc()
    para(d,f"{CHURCH}",11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,f"{yr}년 대심방 현황",20,NAVY,True,AL.CENTER,before=4,after=2,font=SANS); hr(d,"1E6E6A",12)
    para(d,f"완료 {len(done)} / 전체 {total}세대  ·  진행률 {pct}%  ·  {today()}",10,GRAY,align=AL.CENTER,after=8,font=SANS)
    para(d,f"■ 심방 완료 ({len(done)}세대)",13,RGBColor(0x2F,0x7D,0x5B),True,before=6,after=3,font=SANS)
    for m,v in sorted(done,key=lambda x:str((x[1] or {}).get("날짜",""))):
        para(d,f"  · {m['이름']} {m.get('직분','') or ''} ({m.get('소속셀','') or '-'}) — {v.get('날짜','') if v else ''}",10.5,after=1)
    if not done: para(d,"  (아직 없음)",10,LG,after=1)
    para(d,f"■ 미심방 ({len(notyet)}세대)",13,MAROON,True,before=8,after=3,font=SANS)
    for m,_ in notyet:
        para(d,f"  · {m['이름']} {m.get('직분','') or ''} ({m.get('소속셀','') or '-'}) — {m.get('연락처','') or ''}",10.5,after=1)
    if not notyet: para(d,"  (전 세대 완료! 🎉)",10,LG,after=1)
    para(d,"※ 심방 기록 시 '구분'에 '춘계대심방' 또는 '추계대심방'을 넣으면 이 현황에 집계됩니다.",9,LG,before=8,font=SANS)
    out=os.path.join(CAT("04"),f"[대심방현황] {sanit(yr)}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 대심방 현황: {out}  (완료 {len(done)}/{total}세대 · {pct}%)")
def visit_brief(a):
    """다음 심방 전 브리핑 — 반복 방지(심방 기록 카드에 통합됨, CLI 호환용 유지)."""
    db=load(); hit=find(db,a.name)
    if not hit: print(f"✗ '{a.name}' 없음"); return
    _visit_briefing(hit[0])
    print(f"  ⚠ 이미 받은 기도제목(재질문 주의): {', '.join(m['기도제목누적']) if m['기도제목누적'] else '없음'}")
    print("  ▶ 이번 심방은 위와 겹치지 않는 새 말씀·새 안부로 진행하세요.")

# ───────── 명령: 증명서 자동발급 ─────────
CERTS={"교인":("교인 증명원","위의 사람은 본 교회 교인임을 증명합니다."),
       "세례":("세례교인 증명원","위의 사람은 본 교회 세례교인임을 증명합니다."),
       "헌금":("헌금사실 증명원","위의 사람이 본 교회에 헌금하였음을 증명합니다."),
       "재직":("재직증명서","위의 사람은 본 교회의 교인으로 재직하고 있음을 증명합니다."),
       "수료":("수료증","위의 사람은 위 과정을 성실히 이수하였기에 이 증서를 수여합니다."),
       "재적":("재적 증명원","위의 사람은 본 교회에 재적하고 있음을 증명합니다."),
       "출석":("출석 증명원","위의 사람은 본 교회에 성실히 출석하고 있음을 증명합니다."),
       "이명":("이명 증서","위의 사람은 본 교회 교인으로서 다른 교회로 이명함을 증명합니다.")}
AWARDS={"임명":"임명장","위촉":"위촉장"}   # 상장형(명조 격식)
def _cert_ledger(db,no,kind,name,purpose,out):
    """발급대장 자동기록 — 문서번호·발급일·종류·대상·용도·파일."""
    db.setdefault("발급대장",[]).append({"문서번호":no,"발급일":today(),"종류":kind,"대상":name,"용도":purpose or "","파일":os.path.basename(out)})
def cert(a):
    """📜 증명서·상장 발급 & 대장(통합) — 이름을 넣으면 발급(교인·세례·헌금·재직·수료/임명·위촉), 이름을 비우면 발급대장을 본다."""
    if not (getattr(a,'name','') or '').strip():
        return cert_ledger(a)   # 이름 비우면 발급대장 보기(발급 & 대장 통합)
    kind=a.kind or "교인"
    db=load(); hit=find(db,a.name); m=hit[0] if hit else {"이름":a.name,"생년월일":"","주소":""}
    no=a.no or _docno(db,"발급")   # 증명서·상장 공통 단일 연속번호(발급대장 고유번호)
    if kind in AWARDS:   # 상장형(임명장/위촉장)
        title=AWARDS[kind]
        d=newdoc()
        para(d,CHURCH,13,GOLD,True,AL.CENTER,after=0,font=SERIF)
        para(d,title,34,NAVY,True,AL.CENTER,before=10,after=4,font=SERIF); hr(d,"A67C1E",14)
        para(d,f"제 {no} 호",10,LG,align=AL.RIGHT,after=16,font=SANS)
        if a.role: para(d,a.role,16,MAROON,True,AL.CENTER,after=8,font=SERIF)
        para(d,f"{a.name}  귀하",20,NAVY,True,AL.CENTER,after=16,font=SERIF)
        para(d,f"위 사람을 {CHURCH} {a.role or ''}(으)로 {kind}합니다.",14,GRAY,AL.CENTER,before=6,after=4,font=SERIF)
        if a.term: para(d,f"임기: {a.term}",12,GRAY,AL.CENTER,after=4,font=SANS)
        if a.purpose: para(d,f"근거: {a.purpose}",11,LG,AL.CENTER,after=6,font=SANS)
        para(d,f"이에 {title}을(를) 수여합니다.",13,GRAY,AL.CENTER,before=14,after=22,font=SERIF)
        para(d,today().replace('-','년 ',1).replace('-','월 ')+"일",14,NAVY,True,AL.CENTER,after=10,font=SERIF)
        para(d,f"{CHURCH}  담임목사  {PASTOR}   (직인)",15,NAVY,True,AL.CENTER,before=6,font=SERIF)
        out=os.path.join(CAT("01"),f"[{title}] {sanit(a.name)}_{today()}.docx")
    else:               # 증명서형(교인·세례·헌금·재직·수료)
        title,body=CERTS.get(kind,CERTS["교인"])
        d=newdoc()
        para(d,CHURCH,12,GOLD,True,AL.CENTER,after=0,font=SANS)
        para(d,title,22,NAVY,True,AL.CENTER,before=6,after=2,font=SANS); hr(d,"A67C1E",12)
        para(d,f"제 {no} 호",10,LG,after=14,font=SANS)
        for k in ["이름","생년월일","주소"]:
            para(d,f"{k} :  {m.get(k,'')}",12,GRAY,before=2,after=6)
        if kind=="세례": para(d,f"세례일자 :  {m.get('세례','')}",12,GRAY,after=6)
        if kind=="재직": para(d,f"직분 :  {a.role or m.get('직분','')}     재직기간 :  {a.term or ((m.get('등록일','')+' ~ 현재') if m.get('등록일') else '현재 재직')}",12,GRAY,after=6)
        if kind=="수료" and a.role: para(d,f"과정 :  {a.role}",12,GRAY,after=6)
        if a.purpose: para(d,f"용도 :  {a.purpose}",12,GRAY,after=6)
        para(d,"",after=16); para(d,body,12,NAVY,True,AL.CENTER,after=18)
        para(d,today().replace('-','년 ',1).replace('-','월 ')+"일",11,GRAY,AL.CENTER,after=16)
        para(d,f"{CHURCH}  당회장  {PASTOR}   (인)",13,NAVY,True,AL.CENTER,font=SANS)
        out=os.path.join(CAT("06"),f"[발급] {title}_{sanit(m['이름'])}_{today()}.docx")
    _savedoc(d,out); _cert_ledger(db,no,kind,a.name,a.purpose,out); save(db)
    print(f"✅ {title} 발급(제 {no} 호): {out}")
def cert_ledger(a):
    """🗂️ 발급대장 — 지금까지 발급한 증명서·상장 목록(문서번호·발급일·종류·대상·용도)."""
    db=load(); rows=db.get("발급대장",[])
    if a.name: rows=[r for r in rows if a.name in (r.get("대상") or "")]
    print(f"■ 발급대장 · 총 {len(rows)}건" + (f" (대상 '{a.name}')" if a.name else ""))
    for r in rows[-50:]:
        print(f"  {r.get('문서번호','')}  {r.get('발급일','')}  [{r.get('종류','')}]  {r.get('대상','')}  {('· '+r.get('용도')) if r.get('용도') else ''}")
    if not rows: print("  (아직 발급 내역이 없습니다)")

# ───────── 명령: 출석 / 재정 ─────────
def attend_add(a):
    db=load(); db["출석"].append({"날짜":a.date or today(),"예배":a.service or "주일낮","남":a.men or 0,"여":a.women or 0,
        "새신자":a.new or 0,"계":(a.men or 0)+(a.women or 0)}); save(db)
    print(f"✅ 출석 기록: {a.date or today()} {a.service or '주일낮'} 계 {(a.men or 0)+(a.women or 0)}명")
def finance_add(a):
    if not a.amount:
        return finance_sum(a)   # 금액 없으면 재정 요약 보기(기록 & 요약 통합)
    db=load(); db["재정"].append({"날짜":a.date or today(),"구분":a.kind or "수입","항목":a.item or "","금액":a.amount or 0,
        "교인":a.name or "","부서":(getattr(a,'dept','') or "").strip(),"메모":a.memo or ""}); save(db)
    print(f"✅ 재정 기록: {a.date or today()} [{a.kind or '수입'}] {a.item or ''} {a.amount or 0:,}원"+(f" · 부서 {a.dept}" if getattr(a,'dept','') else ""))
def finance_sum(a):
    db=load(); rec=[r for r in db["재정"] if (not a.month or r["날짜"].startswith(a.month))]
    inc=sum(r["금액"] for r in rec if r["구분"]=="수입"); exp=sum(r["금액"] for r in rec if r["구분"]=="지출")
    print(f"■ 재정 요약{f' ({a.month})' if a.month else ''}: 수입 {inc:,}원 · 지출 {exp:,}원 · 잔액 {inc-exp:,}원 ({len(rec)}건)")
def finance_stats(a):
    """재정 통계(JSON) — 월별 수입·지출·헌금종류별·지출항목별·부서별·예산집행. 웹 그래프(시각화)용."""
    import json as _j, datetime as _dt
    db=load(); recs=db.get("재정",[])
    yr=(getattr(a,'year','') or "").strip() or str(_dt.date.today().year)
    yrecs=[r for r in recs if (r.get("날짜") or "").startswith(yr)]
    def _amt(r):
        try: return int(r.get("금액") or 0)
        except Exception: return 0
    mon=[0]*12; mexp=[0]*12; inc={}; exp={}; dept={}
    for r in yrecs:
        v=_amt(r); k=(r.get("항목") or "기타").strip() or "기타"; dp=(r.get("부서") or "미지정").strip() or "미지정"
        try: mi=int(str(r.get("날짜"))[5:7])-1
        except Exception: mi=-1
        if r.get("구분")=="지출":
            if 0<=mi<12: mexp[mi]+=v
            exp[k]=exp.get(k,0)+v; dept.setdefault(dp,{"수입":0,"지출":0})["지출"]+=v
        else:
            if 0<=mi<12: mon[mi]+=v
            inc[k]=inc.get(k,0)+v; dept.setdefault(dp,{"수입":0,"지출":0})["수입"]+=v
    yb=db.get("예산",{}).get(yr,{}); act_in=sum(inc.values()); act_out=sum(exp.values())
    out={"year":yr,"months":mon,"mexp":mexp,
         "inc":sorted(inc.items(),key=lambda x:-x[1])[:8],
         "exp":sorted(exp.items(),key=lambda x:-x[1])[:8],
         "dept":sorted(dept.items(),key=lambda x:-(x[1]["수입"]+x[1]["지출"]))[:10],
         "budget":{"in_bud":sum(yb.get("수입",{}).values()),"in_act":act_in,
                   "out_bud":sum(yb.get("지출",{}).values()),"out_act":act_out},
         "total":{"수입":act_in,"지출":act_out,"잔액":act_in-act_out,"건수":len(yrecs)}}
    print("__FSTATS__"+_j.dumps(out,ensure_ascii=False))
_HEONGEUM=["십일조","주일헌금","주정헌금","감사헌금","작정헌금","서원헌금","일천번제","선교헌금","전도헌금","건축헌금",
 "단기선교헌금","임직헌금","맥추감사헌금","추수감사헌금","성탄감사헌금","부활절헌금","절기헌금","구제헌금","장학헌금","교육헌금","주일학교헌금",
 "청년부헌금","찬양헌금","새벽기도헌금","금요기도헌금","심방헌금","입교·세례헌금","경조헌금","차량헌금","목적헌금","특별헌금","기타헌금"]
_JICHUL=["인건비(사례비)","목회활동비","예배비","강사사례비","교육비","주일학교비","전도비","선교비","구제비","장학비",
 "시설관리비","건축·수리비","비품구입비","소모품비","공과금(전기·수도·가스)","통신비","보험료","행사비","수련회비","친교·식당비",
 "차량유지비","도서·인쇄비","노회비·상회비","경조비","예비비","기타지출"]
def _heongeum(): return _C.get("헌금종류") or _HEONGEUM   # 교회별 설정(config '헌금종류') 우선
def _jichul(): return _C.get("지출항목") or _JICHUL       # 교회별 설정(config '지출항목') 우선
def _finance_annual(db,recs,yr):
    """연간 재정 총결산 — 월별 수입/지출·헌금종류별·지출항목별·부서별 연간 집계(DOCX)."""
    def _amt(r):
        try: return int(r.get("금액") or 0)
        except Exception: return 0
    yrecs=[r for r in recs if (r.get("날짜") or "").startswith(yr)]
    mon=[0]*12; mexp=[0]*12; inc={}; exp={}; dept={}
    for r in yrecs:
        v=_amt(r); k=(r.get("항목") or "기타").strip() or "기타"; dp=(r.get("부서") or "미지정").strip() or "미지정"
        try: mi=int(str(r.get("날짜"))[5:7])-1
        except Exception: mi=-1
        dept.setdefault(dp,[0,0])
        if r.get("구분")=="지출":
            if 0<=mi<12: mexp[mi]+=v
            exp[k]=exp.get(k,0)+v; dept[dp][1]+=v
        else:
            if 0<=mi<12: mon[mi]+=v
            inc[k]=inc.get(k,0)+v; dept[dp][0]+=v
    ti=sum(inc.values()); te=sum(exp.values())
    d=newdoc()
    para(d,f"{CHURCH}  {yr}년 재정 총결산",22,NAVY,bold=True,align=AL.CENTER,after=2,font=SANS)
    para(d,f"작성일: {today()} · 총 {len(yrecs)}건",12,MAROON,bold=True,align=AL.CENTER,after=8,font=SANS); hr(d,"1F5C9E",12)
    para(d,"■ 월별 수입·지출",13,NAVY,bold=True,before=6,after=3,font=SANS)
    t=d.add_table(rows=14,cols=4)
    for j,h in enumerate(["월","수입","지출","잔액"]):
        ctext(t.rows[0].cells[j],h,10.5,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],"1F5C9E")
    for m in range(12):
        ctext(t.rows[m+1].cells[0],f"{m+1}월",10.5,NAVY,font=SANS,align=AL.CENTER)
        ctext(t.rows[m+1].cells[1],f"{mon[m]:,}",10.5,GRAY,font=SANS,align=AL.RIGHT)
        ctext(t.rows[m+1].cells[2],f"{mexp[m]:,}",10.5,GRAY,font=SANS,align=AL.RIGHT)
        ctext(t.rows[m+1].cells[3],f"{mon[m]-mexp[m]:,}",10.5,(MAROON if mon[m]-mexp[m]<0 else NAVY),font=SANS,align=AL.RIGHT)
    for j,val in enumerate(["합계",f"{ti:,}",f"{te:,}",f"{ti-te:,}"]):
        ctext(t.rows[13].cells[j],val,11,MAROON,True,SANS,(AL.CENTER if j==0 else AL.RIGHT)); shade(t.rows[13].cells[j],"EAF0F7")
    for i in range(14):
        for j in range(4): t.rows[i].cells[j].width=Cm(2.6 if j==0 else 4.3)
    tborders(t,"BBBBBB",4)
    def _tbl(title,data,hexc,canon):
        para(d,title,13,NAVY,bold=True,before=10,after=3,font=SANS)
        keys=[k for k in canon if k in data]+[k for k in data if k not in canon]
        if not keys: keys=["(없음)"]
        tt=d.add_table(rows=len(keys)+1,cols=2)
        ctext(tt.rows[0].cells[0],"항목",10.5,WHITE,True,SANS,AL.CENTER); shade(tt.rows[0].cells[0],hexc)
        ctext(tt.rows[0].cells[1],"연간 합계(원)",10.5,WHITE,True,SANS,AL.CENTER); shade(tt.rows[0].cells[1],hexc)
        for i,k in enumerate(keys,1):
            ctext(tt.rows[i].cells[0],k,10.5,NAVY,font=SANS); ctext(tt.rows[i].cells[1],f"{data.get(k,0):,}",10.5,GRAY,font=SANS,align=AL.RIGHT)
        for i in range(len(keys)+1): tt.rows[i].cells[0].width=Cm(9.5); tt.rows[i].cells[1].width=Cm(5.5)
        tborders(tt,"BBBBBB",4)
    _tbl("■ 헌금 종류별 (수입 연간)",inc,"1F5C9E",_heongeum())
    _tbl("■ 지출 항목별 (연간)",exp,"8B2E2E",_jichul())
    if dept:
        para(d,"■ 부서별 수입·지출 (연간)",13,NAVY,bold=True,before=10,after=3,font=SANS)
        dt=d.add_table(rows=len(dept)+1,cols=3)
        for j,h in enumerate(["부서","수입","지출"]):
            ctext(dt.rows[0].cells[j],h,10.5,WHITE,True,SANS,AL.CENTER); shade(dt.rows[0].cells[j],"2F7D5B")
        for i,(dp,iv) in enumerate(sorted(dept.items(),key=lambda x:-(x[1][0]+x[1][1])),1):
            ctext(dt.rows[i].cells[0],dp,10.5,NAVY,font=SANS)
            ctext(dt.rows[i].cells[1],f"{iv[0]:,}",10.5,GRAY,font=SANS,align=AL.RIGHT)
            ctext(dt.rows[i].cells[2],f"{iv[1]:,}",10.5,GRAY,font=SANS,align=AL.RIGHT)
        for i in range(len(dept)+1):
            dt.rows[i].cells[0].width=Cm(6); dt.rows[i].cells[1].width=Cm(4.5); dt.rows[i].cells[2].width=Cm(4.5)
        tborders(dt,"BBBBBB",4)
    para(d,f"■ 연간 총계 — 수입 {ti:,}원 · 지출 {te:,}원 · 잔액 {ti-te:,}원",13,MAROON,bold=True,before=12,font=SANS)
    para(d,"재정부장 (인)              회      계 (인)              담임목사 (인)",11,GRAY,align=AL.CENTER,before=20,font=SANS)
    out=os.path.join(CAT("06"),f"[연간재정총결산] {sanit(yr)}년_{today()}.docx"); _savedoc(d,out); _ext_copy(out)
    print(f"✅ {yr}년 재정 총결산: {out}")
    print(f"   수입 {ti:,}원 · 지출 {te:,}원 · 잔액 {ti-te:,}원 · {len(yrecs)}건 · 부서 {len(dept)}개")
def finance_report(a):
    """주일 재정 결산서 — 헌금 종류별 수입 집계 + 지출 항목별 집계 + 전주이월·당주잔액·차주이월. '연도'를 넣으면 연간 총결산."""
    db=load(); recs=db.get("재정",[])
    yr=(getattr(a,'year','') or "").strip()
    if yr and not (getattr(a,'date',None) or getattr(a,'dfrom',None) or getattr(a,'dto',None)):
        return _finance_annual(db,recs,yr)   # 연간 총결산(월별·종류별·부서별)
    d1=a.date or getattr(a,"dfrom",None) or today()
    d2=a.date or getattr(a,"dto",None) or (a.date or today())
    if a.date: d1=d2=a.date
    if getattr(a,"dfrom",None): d1=a.dfrom
    if getattr(a,"dto",None): d2=a.dto
    def _amt(r):
        try: return int(r.get("금액") or 0)
        except Exception: return 0
    period=[r for r in recs if d1<=(r.get("날짜") or "")<=d2]
    prev=[r for r in recs if (r.get("날짜") or "")<d1]
    prev_bal=sum(_amt(r) for r in prev if r.get("구분")=="수입")-sum(_amt(r) for r in prev if r.get("구분")=="지출")
    inc={}; exp={}
    for r in period:
        key=(r.get("항목") or "").strip() or "기타"
        tgt=exp if r.get("구분")=="지출" else inc
        tgt[key]=tgt.get(key,0)+_amt(r)
    inc_tot=sum(inc.values()); exp_tot=sum(exp.values()); net=inc_tot-exp_tot; carry=prev_bal+net
    d=newdoc(); label=d1 if d1==d2 else f"{d1} ~ {d2}"
    para(d,f"{CHURCH}  주일 재정 결산서",22,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,f"결산일: {label}",12,MAROON,True,AL.CENTER,after=8,font=SANS); hr(d,"1F5C9E",12)
    def _amt_table(title,data,hexc,total_label,total,canon=()):
        para(d,title,13,NAVY,True,before=8,after=3,font=SANS)
        ordered=[(k,data[k]) for k in canon if k in data]+[(k,v) for k,v in data.items() if k not in canon]
        items=ordered or [("(기록 없음)",0)]
        t=d.add_table(rows=len(items)+2,cols=2)
        ctext(t.rows[0].cells[0],"항목",11,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[0],hexc)
        ctext(t.rows[0].cells[1],"금액(원)",11,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[1],hexc)
        for i,(k,v) in enumerate(items,1):
            ctext(t.rows[i].cells[0],k,11,NAVY,font=SANS); ctext(t.rows[i].cells[1],f"{v:,}",11,GRAY,font=SANS,align=AL.RIGHT)
        lr=len(items)+1
        ctext(t.rows[lr].cells[0],total_label,12,MAROON,True,SANS,AL.CENTER); shade(t.rows[lr].cells[0],"EAF0F7")
        ctext(t.rows[lr].cells[1],f"{total:,}",12,MAROON,True,SANS,AL.RIGHT); shade(t.rows[lr].cells[1],"EAF0F7")
        for i in range(len(items)+2): t.rows[i].cells[0].width=Cm(9.5); t.rows[i].cells[1].width=Cm(5.5)
        tborders(t,"BBBBBB",4)
    _amt_table("■ 수입 (헌금)",inc,"1F5C9E","수입 합계",inc_tot,_heongeum())
    _amt_table("■ 지출",exp,"8B2E2E","지출 합계",exp_tot,_jichul())
    para(d,"■ 결산",13,NAVY,True,before=10,after=3,font=SANS)
    srows=[("전주 이월",prev_bal),("당주 수입",inc_tot),("당주 지출",exp_tot),("당주 잔액 (수입−지출)",net),("차주 이월 (전주이월＋당주잔액)",carry)]
    t=d.add_table(rows=len(srows),cols=2)
    for i,(k,v) in enumerate(srows):
        bold=i>=3
        ctext(t.rows[i].cells[0],k,11,(MAROON if bold else NAVY),bold,SANS)
        ctext(t.rows[i].cells[1],f"{v:,} 원",11,(MAROON if bold else GRAY),bold,SANS,AL.RIGHT)
        t.rows[i].cells[0].width=Cm(9.5); t.rows[i].cells[1].width=Cm(5.5)
        if bold: shade(t.rows[i].cells[0],"FBF3E4"); shade(t.rows[i].cells[1],"FBF3E4")
    tborders(t,"BBBBBB",4)
    para(d,"재정부장 (인)              회      계 (인)              담임목사 (인)",11,GRAY,AL.CENTER,before=20,font=SANS)
    out=os.path.join(CAT("06"),f"[주일재정결산] {d1}{'' if d1==d2 else '_'+d2}.docx"); _savedoc(d,out); _ext_copy(out)
    print(f"✅ 주일 재정 결산서: {out}")
    print(f"   수입 {inc_tot:,}원 · 지출 {exp_tot:,}원 · 당주잔액 {net:,}원 · 차주이월 {carry:,}원 (기록 {len(period)}건)")
    if not period: print("   ※ 이 기간 재정 기록이 없습니다 — 먼저 '재정 기록(수입/지출)'에서 헌금·지출을 입력해 주세요.")
def finance_items(a):
    """헌금 종류·지출 항목 표준표 — 재정부가 참고·인쇄할 표준 항목 목록(DOCX). 이 이름으로 입력하면 결산서에 자동 정리."""
    d=newdoc()
    para(d,f"{CHURCH} 재정 표준 항목표",20,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,"재정 기록 시 아래 표준 항목명을 쓰시면, 주일 결산서에 표준 순서로 자동 정리됩니다.",10,GRAY,AL.CENTER,after=6,font=SANS); hr(d,"1F5C9E",12)
    def _grid(title,arr,hexc):
        para(d,title,13,NAVY,True,before=8,after=3,font=SANS)
        cols=3; rows=(len(arr)+cols-1)//cols
        t=d.add_table(rows=rows,cols=cols)
        for i,name in enumerate(arr):
            c=t.rows[i//cols].cells[i%cols]; ctext(c,"· "+name,11,GRAY,font=SANS)
        for rr in t.rows:
            for c in rr.cells: c.width=Cm(5.3)
        tborders(t,"DDDDDD",4)
    _grid("■ 헌금 종류 (수입)",_heongeum(),"1F5C9E")
    _grid("■ 지출 항목",_jichul(),"8B2E2E")
    para(d,"※ 위 목록에 없는 항목도 자유롭게 입력할 수 있습니다(결산서 맨 아래에 표시됩니다).",9.5,LG,before=8,font=SANS)
    out=os.path.join(CAT("06"),f"[재정항목표] {today()}.docx"); _savedoc(d,out); _ext_copy(out)
    print(f"✅ 재정 표준 항목표: {out}")
    print("  헌금 종류:", " · ".join(_heongeum()[:12]), "...")
    print("  지출 항목:", " · ".join(_jichul()[:12]), "...")
def _budget_all(db): return db.setdefault("예산",{})   # {"2026":{"수입":{항목:액},"지출":{항목:액}}}
def budget_set(a):
    """예산 편성 — 연도·구분(수입/지출)·항목·예산액을 한 줄씩 등록/수정. 제직회·공동의회에서 정한 연간 예산을 입력. 금액 0=삭제."""
    yr=str(getattr(a,"year","") or today()[:4]).strip()
    gubun="지출" if (a.kind or "").strip() in ("지출","expense","out") else "수입"
    item=(a.item or "").strip()
    if not item:
        return budget_status(a)   # 항목 없으면 예산 대비 집행현황(편성 & 집행현황 통합)
    amt=getattr(a,"amount",0) or 0
    db=load(); yb=_budget_all(db).setdefault(yr,{"수입":{},"지출":{}})
    if amt<=0:
        yb[gubun].pop(item,None); save(db); print(f"🗑️ {yr}년 예산 삭제: [{gubun}] {item}"); return
    yb[gubun][item]=int(amt); save(db)
    tin=sum(yb["수입"].values()); tout=sum(yb["지출"].values())
    print(f"✅ {yr}년 예산 편성: [{gubun}] {item} {int(amt):,}원")
    print(f"   현재 편성 합계 — 수입예산 {tin:,}원 · 지출예산 {tout:,}원 · 수지 {tin-tout:,}원")
def budget_status(a):
    """예산 대비 집행 현황 — 편성 예산과 실제 재정기록을 대조해 집행률·잔액·초과를 계산한 제직회 보고서(DOCX). 전문가식 예산 결산."""
    yr=str(getattr(a,"year","") or today()[:4]).strip()
    db=load(); yb=_budget_all(db).get(yr)
    def _amt(r):
        try: return int(r.get("금액") or 0)
        except Exception: return 0
    recs=[r for r in db.get("재정",[]) if (r.get("날짜") or "").startswith(yr)]
    act_in={}; act_out={}
    for r in recs:
        k=(r.get("항목") or "").strip() or "기타"
        tgt=act_out if r.get("구분")=="지출" else act_in
        tgt[k]=tgt.get(k,0)+_amt(r)
    if not yb or (not yb.get("수입") and not yb.get("지출")):
        print(f"⚠ {yr}년 예산이 아직 편성되지 않았습니다. 먼저 '예산 편성' 카드에서 항목·금액을 입력해 주세요.")
        if recs: print(f"   (참고: {yr}년 실제 기록은 수입 {sum(act_in.values()):,}원·지출 {sum(act_out.values()):,}원 있습니다.)")
        return
    d=newdoc()
    para(d,f"{CHURCH}  {yr}년 예산 대비 집행 현황",21,NAVY,bold=True,align=AL.CENTER,after=2,font=SANS)
    para(d,f"작성일: {today()}",12,MAROON,bold=True,align=AL.CENTER,after=8,font=SANS); hr(d,"1F5C9E",12)
    def _bud_table(title,bud,act,hexc,canon,is_exp):
        para(d,title,13,NAVY,bold=True,before=8,after=3,font=SANS)
        order=[]; seen=set()
        for k in list(canon)+list(bud)+list(act):
            if (k in bud or k in act) and k not in seen: seen.add(k); order.append(k)
        if not order: order=["(기록 없음)"]
        t=d.add_table(rows=len(order)+2,cols=5)
        for j,h in enumerate(["항목","예산(원)","집행(원)","잔액(원)","집행률"]):
            ctext(t.rows[0].cells[j],h,10.5,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],hexc)
        tb=ta=0
        for i,k in enumerate(order,1):
            bv=int(bud.get(k,0)); av=int(act.get(k,0)); tb+=bv; ta+=av
            rem=bv-av; rate=(av/bv*100) if bv else 0
            ctext(t.rows[i].cells[0],k,10.5,NAVY,font=SANS)
            ctext(t.rows[i].cells[1],f"{bv:,}",10.5,GRAY,font=SANS,align=AL.RIGHT)
            ctext(t.rows[i].cells[2],f"{av:,}",10.5,GRAY,font=SANS,align=AL.RIGHT)
            ctext(t.rows[i].cells[3],f"{rem:,}",10.5,(MAROON if (is_exp and rem<0) else GRAY),font=SANS,align=AL.RIGHT)
            ctext(t.rows[i].cells[4],(f"{rate:.0f}%" if bv else "—"),10.5,NAVY,font=SANS,align=AL.CENTER)
            if is_exp and bv and av>bv:
                shade(t.rows[i].cells[3],"F6D9D5"); shade(t.rows[i].cells[4],"F6D9D5")
        lr=len(order)+1; trem=tb-ta; trate=(ta/tb*100) if tb else 0
        cells=[("합계",AL.CENTER),(f"{tb:,}",AL.RIGHT),(f"{ta:,}",AL.RIGHT),(f"{trem:,}",AL.RIGHT),(f"{trate:.0f}%" if tb else "—",AL.CENTER)]
        for j,(txt,al) in enumerate(cells):
            ctext(t.rows[lr].cells[j],txt,11,MAROON,True,SANS,al); shade(t.rows[lr].cells[j],"EAF0F7")
        widths=[Cm(5.2),Cm(3.4),Cm(3.4),Cm(3.4),Cm(2.2)]
        for i in range(len(order)+2):
            for j in range(5): t.rows[i].cells[j].width=widths[j]
        tborders(t,"BBBBBB",4)
        return tb,ta
    ib,ia=_bud_table("■ 수입 예산 대비 실적",yb.get("수입",{}),act_in,"1F5C9E",_heongeum(),False)
    xb,xa=_bud_table("■ 지출 예산 대비 집행",yb.get("지출",{}),act_out,"8B2E2E",_jichul(),True)
    para(d,"■ 종합",13,NAVY,bold=True,before=10,after=3,font=SANS)
    srows=[("수입: 예산 / 실적",(f"{ib:,} / {ia:,} 원  (달성 {ia/ib*100:.0f}%)" if ib else "—")),
           ("지출: 예산 / 집행",(f"{xb:,} / {xa:,} 원  (집행 {xa/xb*100:.0f}%)" if xb else "—")),
           ("예산상 수지 (수입예산−지출예산)",f"{ib-xb:,} 원"),
           ("실제 수지 (실적−집행)",f"{ia-xa:,} 원")]
    t=d.add_table(rows=len(srows),cols=2)
    for i,(k,v) in enumerate(srows):
        bd=i>=2
        ctext(t.rows[i].cells[0],k,11,(MAROON if bd else NAVY),bd,SANS)
        ctext(t.rows[i].cells[1],v,11,(MAROON if bd else GRAY),bd,SANS,AL.RIGHT)
        t.rows[i].cells[0].width=Cm(9.5); t.rows[i].cells[1].width=Cm(5.5)
        if bd: shade(t.rows[i].cells[0],"FBF3E4"); shade(t.rows[i].cells[1],"FBF3E4")
    tborders(t,"BBBBBB",4)
    para(d,"재정부장 (인)              회      계 (인)              담임목사 (인)",11,GRAY,align=AL.CENTER,before=20,font=SANS)
    out=os.path.join(CAT("06"),f"[예산집행현황] {yr}년_{today()}.docx"); _savedoc(d,out); _ext_copy(out)
    print(f"✅ {yr}년 예산 대비 집행 현황: {out}")
    print(f"   수입 예산 {ib:,} / 실적 {ia:,}원 · 지출 예산 {xb:,} / 집행 {xa:,}원")
    over=[k for k in yb.get('지출',{}) if act_out.get(k,0)>yb['지출'].get(k,0)]
    if over: print("   ⚠ 예산 초과 지출:", ", ".join(over))
def budget_plan(a):
    """예산 편성 양식 — 작년 실제 수입·지출을 참고로 채운 금년 예산안 심의표(제직회·공동의회용 DOCX). 이 표로 예산을 정한 뒤 '예산 편성' 카드에 입력."""
    yr=str(getattr(a,"year","") or today()[:4]).strip(); prev=str(int(yr)-1)
    db=load()
    def _amt(r):
        try: return int(r.get("금액") or 0)
        except Exception: return 0
    pin={}; pout={}
    for r in db.get("재정",[]):
        if not (r.get("날짜") or "").startswith(prev): continue
        k=(r.get("항목") or "").strip() or "기타"
        (pout if r.get("구분")=="지출" else pin).setdefault(k,0)
        (pout if r.get("구분")=="지출" else pin)[k]+=_amt(r)
    d=newdoc()
    para(d,f"{CHURCH}  {yr}년도 예산(안)",21,NAVY,bold=True,align=AL.CENTER,after=2,font=SANS)
    para(d,f"제직회·공동의회 심의용 · 작성일 {today()}",12,MAROON,bold=True,align=AL.CENTER,after=8,font=SANS); hr(d,"1F5C9E",12)
    def _plan_table(title,canon,prevmap,hexc):
        para(d,title,13,NAVY,bold=True,before=8,after=3,font=SANS)
        rows=[k for k in canon]+[k for k in prevmap if k not in canon]
        t=d.add_table(rows=len(rows)+1,cols=4)
        for j,h in enumerate([f"항목","{}년 실적".format(prev),f"{yr}년 예산(안)","비고"]):
            ctext(t.rows[0].cells[j],h,10.5,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],hexc)
        for i,k in enumerate(rows,1):
            pv=int(prevmap.get(k,0))
            ctext(t.rows[i].cells[0],k,10.5,NAVY,font=SANS)
            ctext(t.rows[i].cells[1],(f"{pv:,}" if pv else "—"),10.5,GRAY,font=SANS,align=AL.RIGHT)
            ctext(t.rows[i].cells[2]," ",10.5,NAVY,font=SANS,align=AL.RIGHT)
            ctext(t.rows[i].cells[3]," ",10.5,GRAY,font=SANS)
        widths=[Cm(4.6),Cm(3.6),Cm(3.8),Cm(3.0)]
        for i in range(len(rows)+1):
            for j in range(4): t.rows[i].cells[j].width=widths[j]
        tborders(t,"BBBBBB",4)
    _plan_table("■ 수입 예산 (헌금)",_heongeum(),pin,"1F5C9E")
    _plan_table("■ 지출 예산",_jichul(),pout,"8B2E2E")
    para(d,f"※ {prev}년 실적을 참고해 각 항목 금년 예산을 정하신 뒤, 확정 금액을 '예산 편성' 카드에 입력하시면 연중 집행률이 자동 관리됩니다.",9.5,LG,before=8,font=SANS)
    para(d,"재정부장 (인)              회      계 (인)              담임목사 (인)",11,GRAY,align=AL.CENTER,before=20,font=SANS)
    out=os.path.join(CAT("06"),f"[예산안] {yr}년도_{today()}.docx"); _savedoc(d,out); _ext_copy(out)
    print(f"✅ {yr}년도 예산(안) 편성 양식: {out}")
    print(f"   {prev}년 실적 참고: 수입 {sum(pin.values()):,}원 · 지출 {sum(pout.values()):,}원")
def finance_ledger(a):
    """재정 출납부(원장) — 모든 수입·지출을 날짜순으로 잔액과 함께 정리한 재정부 장부(DOCX). 월/연/기간별."""
    db=load(); recs=db.get("재정",[])
    def _amt(r):
        try: return int(r.get("금액") or 0)
        except Exception: return 0
    month=getattr(a,"month",None); year=getattr(a,"year",None)
    d1=getattr(a,"dfrom",None); d2=getattr(a,"dto",None)
    if month: d1=month+"-01"; d2=month+"-31"; label=month
    elif year: d1=str(year)+"-01-01"; d2=str(year)+"-12-31"; label=str(year)+"년"
    elif d1 or d2: d1=d1 or "0000-00-00"; d2=d2 or "9999-99-99"; label=f"{d1} ~ {d2}"
    else: label=today()[:7]; d1=label+"-01"; d2=label+"-31"
    inp=sorted([r for r in recs if d1<=(r.get("날짜") or "")<=d2], key=lambda r:(r.get("날짜") or ""))
    prev=[r for r in recs if (r.get("날짜") or "")<d1]
    start=sum(_amt(r) for r in prev if r.get("구분")=="수입")-sum(_amt(r) for r in prev if r.get("구분")=="지출")
    d=newdoc()
    para(d,f"{CHURCH}  재정 출납부(원장)",21,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,f"기간: {label}",12,MAROON,True,AL.CENTER,after=8,font=SANS); hr(d,"1F5C9E",12)
    t=d.add_table(rows=len(inp)+3,cols=6)
    for j,h in enumerate(["날짜","항목","적요(헌금자·메모)","수입","지출","잔액"]):
        ctext(t.rows[0].cells[j],h,10.5,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],"1F5C9E")
    ctext(t.rows[1].cells[0],"전기이월",10,GRAY,True,SANS); ctext(t.rows[1].cells[5],f"{start:,}",10,NAVY,True,SANS,AL.RIGHT)
    run=start; tin=tout=0
    for i,r in enumerate(inp,2):
        amt=_amt(r); isout=r.get("구분")=="지출"; run+=(-amt if isout else amt)
        tin+=(0 if isout else amt); tout+=(amt if isout else 0)
        adj=((r.get("교인") or "")+("  "+r.get("메모") if r.get("메모") else "")).strip()
        ctext(t.rows[i].cells[0],r.get("날짜") or "",9.5,GRAY,font=SANS,align=AL.CENTER)
        ctext(t.rows[i].cells[1],r.get("항목") or "",9.5,NAVY,font=SANS)
        ctext(t.rows[i].cells[2],adj,9,GRAY,font=SANS)
        ctext(t.rows[i].cells[3],("" if isout else f"{amt:,}"),9.5,GRAY,font=SANS,align=AL.RIGHT)
        ctext(t.rows[i].cells[4],(f"{amt:,}" if isout else ""),9.5,MAROON,font=SANS,align=AL.RIGHT)
        ctext(t.rows[i].cells[5],f"{run:,}",9.5,NAVY,font=SANS,align=AL.RIGHT)
    lr=len(inp)+2
    ctext(t.rows[lr].cells[0],"합계",11,MAROON,True,SANS,AL.CENTER)
    ctext(t.rows[lr].cells[3],f"{tin:,}",10.5,NAVY,True,SANS,AL.RIGHT)
    ctext(t.rows[lr].cells[4],f"{tout:,}",10.5,MAROON,True,SANS,AL.RIGHT)
    ctext(t.rows[lr].cells[5],f"{run:,}",10.5,NAVY,True,SANS,AL.RIGHT)
    for c in range(6): shade(t.rows[lr].cells[c],"EAF0F7")
    for rr in t.rows:
        for j,w in enumerate([2.2,2.6,4.0,2.3,2.3,2.8]): rr.cells[j].width=Cm(w)
    tborders(t,"BBBBBB",4)
    para(d,"재정부장 (인)              회      계 (인)              담임목사 (인)",11,GRAY,AL.CENTER,before=16,font=SANS)
    out=os.path.join(CAT("06"),f"[재정출납부] {label.replace(' ','')}.docx"); _savedoc(d,out); _ext_copy(out)
    print(f"✅ 재정 출납부(원장): {out}")
    print(f"   전기이월 {start:,} · 수입 {tin:,} · 지출 {tout:,} · 잔액 {run:,}원 ({len(inp)}건)")
    if not inp: print("   ※ 이 기간 재정 기록이 없습니다 — '재정 기록'으로 먼저 입력하세요.")

# ───────── 명령: 주보 ─────────
def bulletin(a):
    db=load(); d=newdoc()
    para(d,f"{CHURCH} 주보",22,NAVY,True,AL.CENTER,after=1,font=SANS)
    para(d,f"주후 {a.date or today()}  ·  담임 "+PASTOR,10,LG,AL.CENTER,after=4,font=SANS); hr(d,"A67C1E",12)
    para(d,"예배 순서",13,MAROON,True,before=4,after=3,font=SANS)
    order=(a.order or "묵도;찬송;기도;성경봉독;찬양;설교;헌금;광고;축도").split(";")
    for i,o in enumerate(order,1): para(d,f"  {i}. {o.strip()}",11,after=2)
    if a.sermon: para(d,f"※ 설교: {a.sermon}",11,NAVY,True,before=4)
    para(d,"교회 소식",13,MAROON,True,before=8,after=3,font=SANS)
    for n in (a.notice or "").split(";"):
        if n.strip(): para(d,"· "+n.strip(),11,after=2)
    # 주간 심방/행사 자동
    para(d,"주간 일정",13,MAROON,True,before=8,after=3,font=SANS)
    for w in (a.week or "").split(";"):
        if w.strip(): para(d,"· "+w.strip(),11,after=2)
    para(d,"이 주의 기도",13,MAROON,True,before=8,after=3,font=SANS)
    para(d,a.pray or "  ",11,after=2)
    out=os.path.join(CAT("07"),f"[주보] {a.date or today()}.docx"); _savedoc(d,out)
    print(f"✅ 주보 생성: {out}")

# ───────── 명령: 교안 ─────────
def lesson(a):
    d=newdoc(); TEAL=RGBColor(0x1E,0x6E,0x6A)
    para(d,"교안 · Lesson Plan",11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,a.title or "제목",20,NAVY,True,AL.CENTER,before=4,after=2,font=SANS); hr(d,"1E6E6A",12)
    meta=[("부서/대상",a.target or ""),("본문",a.text or ""),("주제",a.theme or ""),("날짜",a.date or today())]
    for k,v in meta: para(d,f"{k} :  {v}",11,GRAY,after=3)
    para(d,"학습 목표",13,TEAL,True,before=8,after=2,font=SANS)
    para(d,a.goal or "  (이 시간 후 학생이 알게/느끼게/행하게 될 것을 한 문장으로)",11,italic=(not a.goal),after=4)
    guide=[("도입 (마음 열기)","관심을 끄는 질문·이야기·활동으로 시작해, 학생의 삶과 본문을 잇는 다리를 놓습니다(약 5~10분)."),
           ("전개 (말씀 이해)","본문을 함께 읽고 핵심 진리를 눈높이에 맞게. 질문으로 스스로 발견하게 하고, 결국 그리스도로 귀결(약 15~20분)."),
           ("적용 (삶으로)","'그래서 이번 주 어떻게 살까?' 구체적이고 실천 가능한 한 가지로. 나눔 질문으로 결단까지(약 5~10분)."),
           ("정리·기도","핵심을 한 문장으로 정리 + 암송구절 + 함께 기도. 다음 시간 예고(약 5분).")]
    for sec,tip in guide:
        para(d,sec,13,TEAL,True,before=8,after=1,font=SANS)
        para(d,f"  ▷ {tip}",9.5,LG,italic=True,after=3)
        para(d,"  ",11,after=8)
    hr(d,"CCCCCC",6)
    para(d,"★ 좋은 교안 작성 가이드",12,MAROON,True,before=6,after=3,font=SANS)
    for t in ["① 목표는 하나로: 이 한 시간에 남길 '한 문장'을 먼저 정하세요. 목표가 흐리면 교안도 흐려집니다.",
              "② 눈높이 언어: 대상(유아·초등·청소년·청년·장년)의 말과 경험으로. 어려운 신학 용어는 쉬운 말로 풀어서.",
              "③ 질문이 힘: 가르치기보다 물어서 스스로 발견하게. 정답형이 아닌 열린 질문(왜·어떻게)을 준비하세요.",
              "④ 활동·예화 하나: 잘 고른 활동이나 예화 하나가 열 마디 설명보다 오래 남습니다(대상에 맞게).",
              "⑤ 그리스도로 귀결: 어떤 본문이든 도덕 교훈에서 멈추지 말고 예수 그리스도와 복음으로 연결하세요.",
              "⑥ 적용은 구체적으로: '착하게 살자'가 아니라 '이번 주 ○○을 해보자'처럼 실천 가능한 한 가지로.",
              "⑦ 시간 배분: 도입10·전개20·적용10·정리5(예시). 전개에 다 쓰지 말고 적용·나눔 시간을 꼭 남기세요.",
              "⑧ 준비물·대안: 필요한 준비물을 미리 적고, 시간이 남거나 모자랄 때의 대안도 한 줄 준비해 두세요."]:
        para(d,t,10,after=2)
    para(d,"※ 더 깊은 내용은 '주일학교 지도자료'·'성경교육 자료' 참고자료 카드를 함께 보세요.",9,LG,before=6,font=SANS)
    out=os.path.join(CAT("03"),f"[교안] {sanit(a.title or '교안')}_{a.date or today()}.docx"); _savedoc(d,out)
    print(f"✅ 교안 생성: {out} (작성 가이드 포함)")

# ───────── 명령: Excel 동기화 ─────────
def export_excel(a):
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    db=load()
    wb=openpyxl.Workbook();
    thin=Side(style='thin',color='BBBBBB'); bd=Border(thin,thin,thin,thin)
    def sheet(name,headers,rows,color):
        ws=wb.create_sheet(name); hf=Font(bold=True,color='FFFFFF',name='맑은 고딕',size=10)
        fill=PatternFill('solid',fgColor=color)
        for j,h in enumerate(headers,1):
            c=ws.cell(1,j,h); c.font=hf; c.fill=fill; c.alignment=Alignment('center','center'); c.border=bd
        for i,row in enumerate(rows,2):
            for j,v in enumerate(row,1):
                c=ws.cell(i,j,v); c.font=Font(name='맑은 고딕',size=10); c.border=bd
        for j,h in enumerate(headers,1):
            ws.column_dimensions[chr(64+j)].width=max(10,len(str(h))*2+4)
        return ws
    wb.remove(wb.active)
    sheet("교인명부",["ID","이름","성별","생년월일","연락처","주소","직분","세례","소속셀","인도자","등록일","상태"],
        [[m["id"],m["이름"],m["성별"],m["생년월일"],m["연락처"],m["주소"],m["직분"],m["세례"],m["소속셀"],m["인도자"],m["등록일"],m["상태"]] for m in db["교인"]],"3C6E3C")
    vrows=[]
    for m in db["교인"]:
        for v in m["심방이력"]:
            vrows.append([m["이름"],v["날짜"],v["구분"],v["심방자"],v["말씀"],v["내용"],"; ".join(v["기도제목"]),v["후속"]])
    sheet("심방이력",["이름","날짜","구분","심방자","말씀","내용","기도제목","후속"],vrows,"3C6E3C")
    sheet("출석",["날짜","예배","남","여","새신자","계"],[[r["날짜"],r["예배"],r["남"],r["여"],r["새신자"],r["계"]] for r in db["출석"]],"1B2A4A")
    sheet("재정",["날짜","구분","항목","금액","교인","메모"],[[r["날짜"],r["구분"],r["항목"],r["금액"],r["교인"],r["메모"]] for r in db["재정"]],"A67C1E")
    sheet("설교이력",["날짜","예배","제목","본문","시리즈"],[[s.get("날짜",""),s.get("예배",""),s.get("제목",""),s.get("본문",""),s.get("시리즈","")] for s in db.get("설교",[])],"5A4FCF")
    sheet("경조사",["날짜","종류","대상","내용","경조금"],[[e.get("날짜",""),e.get("종류",""),e.get("대상",""),e.get("내용",""),e.get("경조금","")] for e in db.get("경조사",[])],"8B2E2E")
    sheet("일정",["날짜","유형","장소","주제","담당"],[[e.get("날짜",""),e.get("유형",""),e.get("장소",""),e.get("주제",""),e.get("담당자","")] for e in db.get("일정",[])],"C25A2E")
    sheet("태신자",["이름","담당","연락처","접촉수","상태"],[[v.get("이름",""),v.get("담당",""),v.get("연락처",""),len(v.get("접촉",[])),v.get("상태","")] for v in db.get("태신자",[])],"1F5C9E")
    grows=[[g.get("이름",""),g.get("리더",""),g.get("요일",""),g.get("장소",""),len([m for m in db["교인"] if m.get("소속셀")==g["이름"]])] for g in db.get("소그룹",[])]
    sheet("소그룹",["소그룹","리더","요일","장소","조원수"],grows,"3C6E3C")
    trows=[[m["이름"],t.get("과정",""),t.get("단계",""),t.get("상태",""),t.get("날짜","")] for m in db["교인"] for t in m.get("양육이력",[])]
    sheet("양육이력",["이름","과정","단계","상태","날짜"],trows,"1E6E6A")
    sheet("노회",["날짜","구분","내용","담당직무","기한"],[[e.get("날짜",""),e.get("구분",""),e.get("내용",""),e.get("담당직무",""),e.get("기한","")] for e in db.get("노회",[])],"6A3D7A")
    xd=os.path.join(ROOT,"_관리대장(Excel)"); os.makedirs(xd,exist_ok=True)
    out=os.path.join(xd,f"{CHURCH} 관리대장.xlsx"); wb.save(out); _ext_copy(out); _openfile(out)
    print(f"✅ Excel 동기화: {out}")
    print(f"   시트: 교인·심방·출석·재정·설교이력·경조사·일정·태신자·소그룹·양육·노회 (모든 기록) · 이중저장 됨")

def setlist(a):
    """찬양 콘티 — 곡 목록 → 악보(A4/A3)·key·BPM·가사. 찬양팀 연습·인도용. 자작곡 라이브러리 연동."""
    import json as _j
    def norm(s): return re.sub(r'\s','',s)
    # 통합 곡사전: song_index(이미지 색인) + song_library(등록·자작곡)
    idx={}
    idxp=os.path.join(BASE,"song_index.json")
    if os.path.exists(idxp):
        for k,v in _j.load(open(idxp,encoding='utf-8')).items():
            idx[k]={"title":v["title"],"paths":v.get("paths",[]),"meta":{}}
    if os.path.exists(SONGLIB):
        for s in _j.load(open(SONGLIB,encoding='utf-8')):
            pth=[s["악보"]] if s.get("악보") and os.path.exists(s["악보"]) and s["악보"].lower().endswith(('.jpg','.jpeg','.png')) else []
            idx[norm(s["제목"])]={"title":s["제목"],"paths":pth,
                "meta":{"key":s.get("key",""),"bpm":s.get("bpm",""),"작곡":s.get("작곡가",""),"가사":s.get("가사",""),"유튜브":s.get("유튜브","")}}
    if not idx: print("✗ 곡이 없습니다 — '찬양곡·자작곡 등록'에서 곡을 추가해 주세요."); return
    songs=[s.strip() for s in (a.songs or "").split(";") if s.strip()]
    size=(a.size or "A4").upper()
    W,H=(Cm(29.7),Cm(42.0)) if size=="A3" else (Cm(21.0),Cm(29.7))
    d=Document()
    for s in d.sections:
        s.page_width=W; s.page_height=H
        s.top_margin=Cm(1.2); s.bottom_margin=Cm(1.2); s.left_margin=Cm(1.2); s.right_margin=Cm(1.2)
    para(d,f"{CHURCH} 찬양팀",12,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,a.title or "찬양 콘티",24,NAVY,True,AL.CENTER,before=4,after=2,font=SANS)
    para(d,f"{a.date or today()}  ·  {size} 인쇄용",10,LG,AL.CENTER,after=6,font=SANS); hr(d,"A67C1E",12)
    para(d,"오늘의 찬양 순서",14,MAROON,True,before=4,after=4,font=SANS)
    def mtag(m):
        return "  ".join(x for x in [f"[{m.get('key')}]" if m.get('key') else "",
                f"{m.get('bpm')}BPM" if m.get('bpm') else "", f"· {m.get('작곡')} 작곡" if m.get('작곡') else ""] if x)
    found=[]; miss=[]
    for i,q in enumerate(songs,1):
        nk=norm(q); hit=None
        for k,v in idx.items():
            if nk in k or k in nk: hit=v; break
        if hit: found.append((q,hit)); para(d,f"  {i}. {hit['title']}   {mtag(hit['meta'])}",13,NAVY,False,after=3)
        else: miss.append(q); para(d,f"  {i}. {q}   (미등록 — '찬양곡·자작곡 등록'에서 추가)",12,MAROON,False,after=3)
    imgw=W-Cm(2.6)
    for q,hit in found:
        d.add_page_break(); m=hit["meta"]
        para(d,hit["title"]+(f"   {mtag(m)}" if mtag(m) else ""),15,NAVY,True,AL.CENTER,after=4,font=SANS)
        if m.get("유튜브"): para(d,f"🎬 유튜브: {m['유튜브']}",10,RGBColor(0x1F,0x5C,0x9E),align=AL.CENTER,after=4,font=SANS)
        if hit["paths"]:
            try: d.add_picture(hit["paths"][0],width=imgw); d.paragraphs[-1].alignment=AL.CENTER
            except Exception as e: para(d,f"[악보 삽입 실패: {e}]",10,MAROON)
        elif m.get("가사"):
            for ln in m["가사"].split("\n"): para(d,ln,13,GRAY,align=AL.CENTER,after=2)
        else: para(d,"(악보·가사 미등록)",11,MAROON,align=AL.CENTER)
    out=os.path.join(CAT("07"),f"[찬양콘티] {sanit(a.title or '콘티')}_{a.date or today()}_{size}.docx"); _savedoc(d,out)
    print(f"✅ 찬양 콘티({size}) 생성: {out}")
    print(f"   수록 {len(found)}곡"+(f" · 미등록 {len(miss)}곡: {', '.join(miss)}" if miss else ""))
SONGDIR=os.path.join(BASE,"_내자료","찬양곡"); SONGLIB=os.path.join(SONGDIR,"song_library.json")
def song_add(a):
    """찬양곡 추가(자작곡 포함) — 제목·작곡가·key·BPM·주제·가사·악보. 콘티에서 자동 사용."""
    if not a.title: print("⚠ 곡 제목을 입력해 주세요."); return
    import shutil, json as _j
    os.makedirs(SONGDIR,exist_ok=True)
    lib=_j.load(open(SONGLIB,encoding='utf-8')) if os.path.exists(SONGLIB) else []
    sheet=""
    if a.sheet and os.path.exists(a.sheet):
        dst=os.path.join(SONGDIR,os.path.basename(a.sheet))
        if os.path.abspath(a.sheet)!=os.path.abspath(dst): shutil.copy2(a.sheet,dst)
        sheet=dst
    lib.append({"제목":a.title or "","작곡가":a.composer or "","작사":a.lyricist or "","key":a.key or "","bpm":a.bpm or "",
        "주제":a.theme or "","성경":a.bible or "","가사":a.lyrics or "","악보":sheet,"유튜브":a.youtube or "",
        "발표일":a.release or "","앨범":a.album or "","저작권":a.copyright or "",
        "유형":a.type or ("자작곡" if a.composer else "찬양")})
    _j.dump(lib,open(SONGLIB,'w',encoding='utf-8'),ensure_ascii=False,indent=2)
    _ext_copy(SONGLIB)            # 곡 목록 USB·D 이중저장
    if sheet: _ext_copy(sheet)    # 악보 파일도 이중저장
    print(f"✅ 찬양곡 추가: {a.title} [{a.key or '-'}/{a.bpm or '-'}BPM]{' · 자작곡('+a.composer+')' if a.composer else ''}{' 🎬' if a.youtube else ''} → 총 {len(lib)}곡")
    if _C.get("백업폴더"): print("  💾 USB·D에도 자동 저장됨(작곡·악보 안전 이중보관)")
def song_list(a):
    import json as _j
    if not os.path.exists(SONGLIB): print("등록 곡 없음 — '찬양곡·자작곡 등록'에서 추가해 주세요."); return
    lib=_j.load(open(SONGLIB,encoding='utf-8'))
    if a.theme: lib=[s for s in lib if a.theme in (s.get("주제","")+s.get("성경","")+s.get("제목",""))]
    print(f"■ 내 찬양곡 {len(lib)}곡"+(f" (주제:{a.theme})" if a.theme else ""))
    for i,s in enumerate(lib):
        print(f"  {i}. {s['제목']} [{s.get('key','')}/{s.get('bpm','')}BPM] {s.get('유형','')}"+(f" — {s['작곡가']} 작곡" if s.get('작곡가') else ""))
        if s.get("유튜브"): print(f"       🎬 {s['유튜브']}")
def song_open(a):
    """등록된 찬양곡의 유튜브를 브라우저로 열기 — 작곡 목사님이 자기 올린 영상 바로 보기."""
    import json as _j, webbrowser
    if not os.path.exists(SONGLIB): print("등록 곡 없음."); return
    lib=_j.load(open(SONGLIB,encoding='utf-8'))
    yt=[s for s in lib if s.get("유튜브")]
    if a.title:
        hit=[s for s in yt if a.title in s["제목"]]
        if hit and hit[0].get("유튜브"): webbrowser.open(hit[0]["유튜브"]); print(f"▶ 여는 중: {hit[0]['제목']} — {hit[0]['유튜브']}")
        else: print(f"✗ '{a.title}' 유튜브 없음")
    else:
        print(f"■ 유튜브 등록 찬양곡 {len(yt)}곡 (song-open --title 제목 으로 열기)")
        for s in yt: print(f"  🎬 {s['제목']} — {s['유튜브']}")
def compose(a):
    """찬양 작곡 도구 — 찬양 가사 틀 + 작곡도구(Suno/Udio) 프롬프트 생성."""
    d=newdoc()
    para(d,"찬양 작곡 · 작업지",11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,a.title or "새 찬양",20,NAVY,True,AL.CENTER,before=4,after=2,font=SANS); hr(d,"A67C1E",12)
    for k,v in [("주제",a.theme or ""),("성경 본문",a.bible or ""),("분위기",a.mood or "은혜롭고 벅찬"),
                ("조성 / 템포",f"{a.key or 'G'} / {a.bpm or '72'}BPM")]:
        para(d,f"{k} :  {v}",11,GRAY,after=3)
    para(d,"가사 (직접 창작 — 저작권 무관)",13,MAROON,True,before=8,after=3,font=SANS)
    if a.lyrics:
        for ln in a.lyrics.split("\n"): para(d,ln,11,after=2)
    else:
        for sec in ["[1절]","[후렴]","[2절]","[브릿지]"]:
            para(d,sec,12,NAVY,True,before=4,after=2,font=SANS); para(d,"  (가사)",11,LG,after=4)
    para(d,"작곡 프롬프트 (Suno · Udio 등에 붙여넣기)",13,MAROON,True,before=8,after=3,font=SANS)
    prompt=f"Korean worship CCM, {a.mood or 'uplifting and reverent'}, key {a.key or 'G'} major, {a.bpm or '72'} BPM, theme of {a.theme or 'praise and grace'}, warm lead vocal, piano and strings, singable congregational chorus, studio quality"
    para(d,prompt,11,GRAY,after=4)
    para(d,"※ 실제 멜로디·음원 생성은 외부 작곡도구(Suno 등, 무료/유료)에서 위 프롬프트+가사로 만드세요. 완성곡을 유튜브에 올린 뒤 '찬양곡·자작곡 등록'에서 유튜브 링크와 함께 등록하면 콘티에 자동 반영됩니다.",9,LG,before=4,font=SANS)
    folder=CAT("07")
    out=os.path.join(folder,f"[작곡작업지] {sanit(a.title or '새찬양')}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 작곡 작업지 생성: {out}"); print(f"   작곡 프롬프트: {prompt}")
def video_plan(a):
    """행사 홍보영상 제작 작업지 — 스토리보드·나레이션·자막·영상AI 프롬프트·무료 편집 안내를 한 장에.
       수노(음악)처럼 영상 AI(Runway·Pika·Kling·Canva)와 무료 편집기(Clipchamp·Canva·CapCut)를 연동해 초보도 홍보영상을 만들게 한다."""
    title=(getattr(a,'title','') or "우리 교회 행사").strip()
    purpose=(getattr(a,'purpose','') or "성도·지역 초청").strip()
    tone=(getattr(a,'tone','') or "따뜻하고 설레는").strip()
    message=(getattr(a,'message','') or "함께 예배하며 은혜받는 시간에 초대합니다").strip()
    when=(getattr(a,'when','') or "(일시·장소를 적어 주세요)").strip()
    cta=(getattr(a,'cta','') or "신청·문의: 교회 사무실").strip()
    try: L=int(re.sub(r'[^0-9]','',str(getattr(a,'length','') or "30")) or "30")
    except Exception: L=30
    d=newdoc()
    para(d,"홍보영상 제작 작업지 · 영상 AI 연동",11,GOLD,bold=True,align=AL.CENTER,after=0,font=SANS)
    para(d,title,20,NAVY,bold=True,align=AL.CENTER,before=4,after=2,font=SANS); hr(d,"A67C1E",12)
    para(d,"■ 영상 개요",13,NAVY,bold=True,before=6,after=3,font=SANS)
    for k,v in [("행사·제목",title),("목적·대상",purpose),("길이",f"{L}초 (카톡·릴스 15~30초 / 유튜브 30~60초 권장)"),
                ("톤·분위기",tone),("핵심 메시지",message),("일시·장소",when),
                ("배포처","카카오톡 · 교회 유튜브 · 인스타 릴스 · 밴드  (세로 9:16 / 가로 16:9)")]:
        para(d,f"· {k} :  {v}",11,GRAY,after=2,font=SANS)
    para(d,"■ 장면 구성(스토리보드) — 이 순서대로 만들면 됩니다",13,NAVY,bold=True,before=8,after=3,font=SANS)
    scenes=[
      ("씬1 · 오프닝",f"0~{max(2,round(L*0.1))}초","시선을 끄는 한 컷 + 행사명 크게",title,"임팩트 있는 음악 시작"),
      ("씬2 · 초대",f"~{round(L*0.33)}초","따뜻한 교회·성도 모습, 웃는 얼굴","당신을 초대합니다",message),
      ("씬3 · 핵심 정보",f"~{round(L*0.66)}초","일시·장소·대상을 큰 자막으로",when,"언제, 어디서, 누구나"),
      ("씬4 · 감성·비전",f"~{round(L*0.9)}초","은혜로운 장면(예배·찬양·공동체)",message,"기대와 은혜를 담아"),
      ("씬5 · 클로징",f"~{L}초","행사명 + 신청/문의 + 교회 이름",cta,"잔잔한 마무리 음악"),
    ]
    t=d.add_table(rows=len(scenes)+1,cols=5)
    for j,h in enumerate(["장면","시간","화면(무엇을)","자막","나레이션·음악"]):
        ctext(t.rows[0].cells[j],h,10,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],"7A3EA6")
    for i,(nm,tm,scr,sub,nar) in enumerate(scenes,1):
        ctext(t.rows[i].cells[0],nm,9.5,NAVY,True,SANS)
        ctext(t.rows[i].cells[1],tm,9.5,GRAY,font=SANS,align=AL.CENTER)
        ctext(t.rows[i].cells[2],scr,9.5,GRAY,font=SANS)
        ctext(t.rows[i].cells[3],sub,9.5,MAROON,font=SANS)
        ctext(t.rows[i].cells[4],nar,9.5,GRAY,font=SANS)
    widths=[Cm(2.4),Cm(1.9),Cm(4.4),Cm(3.4),Cm(4.0)]
    for i in range(len(scenes)+1):
        for j in range(5): t.rows[i].cells[j].width=widths[j]
    tborders(t,"BBBBBB",4)
    para(d,"■ 나레이션 대본 (초안 — 자유롭게 다듬으세요)",13,NAVY,bold=True,before=10,after=3,font=SANS)
    for s in [f"“{title}에 여러분을 초대합니다.”",f"“{message}”",f"“{when}.”",f"“{cta}. 함께해요.”"]:
        para(d,s,11,GRAY,after=2,font=SANS)
    para(d,"■ 화면 자막 문구 (크게 넣을 것)",13,NAVY,bold=True,before=8,after=3,font=SANS)
    for s in [title,when,message,cta]: para(d,f"· {s}",11,GRAY,after=2,font=SANS)
    para(d,"■ 영상 AI 프롬프트 (Runway · Pika · Kling · Canva Magic Media 등에 붙여넣기)",13,MAROON,bold=True,before=10,after=3,font=SANS)
    prompts=[
      ("오프닝 분위기 컷","warm cinematic establishing shot, welcoming church community atmosphere, soft golden morning light, gentle camera push-in, hopeful inviting mood, 16:9, high quality"),
      ("감성·공동체 컷","heartwarming scene of people in joyful fellowship, soft bokeh background, warm tones, uplifting worship atmosphere, cinematic, 9:16 vertical"),
      ("클로징 배경","peaceful serene background with soft light rays, calm reverent mood, subtle motion, space for text overlay, 16:9"),
    ]
    for nm,pr in prompts:
        para(d,f"[{nm}]",11,NAVY,bold=True,after=1,font=SANS); para(d,pr,10.5,GRAY,after=3,font=SANS)
    para(d,"※ 실제 영상 생성은 외부 영상 AI(무료 한도 있음)에서 위 프롬프트로 만드세요. 실제 사람 얼굴·우리 교회 건물은 직접 촬영/사진이 더 자연스럽습니다 — AI는 오프닝·배경·분위기 컷에 쓰는 것을 추천합니다.",9,LG,before=2,font=SANS)
    para(d,"■ 무료로 조립하는 법 (가장 쉬운 순서)",13,NAVY,bold=True,before=10,after=3,font=SANS)
    for s in [
      "1) Canva 또는 Clipchamp에서 ‘홍보영상 템플릿’을 엽니다 (배경·자막 틀이 이미 준비됨).",
      "2) 위 스토리보드대로 사진·영상 클립을 끼워 넣습니다 (직접 찍은 것 + AI 생성 컷).",
      "3) 자막을 붙입니다 — CapCut·Clipchamp의 ‘자동 자막’이 나레이션을 글자로 자동 변환해 줍니다.",
      "4) 나레이션은 Clipchamp·CapCut의 ‘AI 음성(TTS)’으로 대본을 읽게 하거나 직접 녹음합니다. (한국어는 네이버 클로바더빙도 자연스럽습니다.)",
      "5) 배경음악은 저작권 무료 음원(유튜브 오디오 보관함 등)이나, ‘AI 찬양 작곡 (수노·작업지)’ 카드로 수노에서 직접 만들어 넣습니다.",
      "6) 내보내기 — 카톡·릴스용은 세로(9:16), 유튜브용은 가로(16:9)로 저장합니다.",
    ]: para(d,s,11,GRAY,after=2,font=SANS)
    para(d,"■ 무료 영상 편집 프로그램 (초보 목사님 추천 순)",13,NAVY,bold=True,before=8,after=3,font=SANS)
    tools=[("Clipchamp","윈도우11에 이미 설치됨(설치 불필요) · 자동 자막·AI 음성·템플릿이 쉬움","★★★ 가장 쉬움"),
           ("Canva","홍보영상 템플릿 풍부 · 클릭만으로 완성 · 포스터·썸네일도 한 곳에","★★★"),
           ("CapCut","자동 자막·AI 기능 강력 · 트렌디한 템플릿","★★☆"),
           ("DaVinci Resolve","무료인데 전문가급 · 익숙해지면 고급 편집","★☆☆")]
    tt=d.add_table(rows=len(tools)+1,cols=3)
    for j,h in enumerate(["프로그램","특징","초보 적합"]):
        ctext(tt.rows[0].cells[j],h,10,WHITE,True,SANS,AL.CENTER); shade(tt.rows[0].cells[j],"1F5C9E")
    for i,(nm,ft,st) in enumerate(tools,1):
        ctext(tt.rows[i].cells[0],nm,10.5,NAVY,True,SANS)
        ctext(tt.rows[i].cells[1],ft,10,GRAY,font=SANS)
        ctext(tt.rows[i].cells[2],st,10,MAROON,font=SANS,align=AL.CENTER)
    for i in range(len(tools)+1):
        tt.rows[i].cells[0].width=Cm(3.4); tt.rows[i].cells[1].width=Cm(9.0); tt.rows[i].cells[2].width=Cm(3.2)
    tborders(tt,"BBBBBB",4)
    para(d,"※ 더 자세한 도구·AI 사용법·목적별 템플릿은 ‘영상 제작·AI 완전 가이드’ 카드를 보세요.",9.5,LG,before=8,font=SANS)
    out=os.path.join(CAT("07"),f"[홍보영상작업지] {sanit(title)}_{today()}.docx"); _savedoc(d,out); _ext_copy(out)
    print(f"✅ 홍보영상 제작 작업지: {out}")
    print(f"   길이 {L}초 · 스토리보드 5장면 · 영상AI 프롬프트 3종 · 무료 편집기 안내 포함")
    print("   → 이 작업지대로 무료 편집기(Clipchamp·Canva·CapCut)로 만드시면 됩니다. ('무료 영상편집·영상 AI 완전 가이드' 카드 참고)")
    print("   ※ 사진만 넣으면 자동으로 영상을 만들어 주는 기능은 다음 업데이트에 정식 지원 예정입니다.")
def _vwrap(s,n=16):
    """자막용 한글 줄바꿈 — n자 넘으면 공백·구분점에서 줄바꿈."""
    s=str(s).strip(); out=[]; line=""
    for ch in s:
        line+=ch
        if len(line)>=n and ch in " ·,.!?":
            out.append(line.strip()); line=""
        elif len(line)>=n+5:
            out.append(line.strip()); line=""
    if line.strip(): out.append(line.strip())
    return "\n".join([x for x in out if x]) or s
def video_render(a):
    """홍보영상 자동 생성 — 대본/스토리보드 + 목사님 사진·짧은 영상·음악을 합쳐 실제 영상 파일(.mp4)을 만든다.
       소스는 _내자료/영상소스 폴더에 넣는다. ffmpeg 사용(없으면 안내). 영상 장비 못 다뤄도 완성 홍보영상이 나온다."""
    import subprocess, tempfile, shutil
    ff=None   # 번들(_시스템/ffmpeg.exe) 우선, 없으면 시스템 PATH
    for c in (os.path.join(BASE,"ffmpeg.exe"), os.path.join(BASE,"_bin","ffmpeg.exe"), os.path.join(BASE,"ffmpeg","ffmpeg.exe")):
        if os.path.exists(c): ff=c; break
    if not ff: ff=shutil.which("ffmpeg")
    title=(getattr(a,'title','') or "우리 교회 행사").strip()
    message=(getattr(a,'message','') or "함께 예배하며 은혜받는 시간에 초대합니다").strip()
    when=(getattr(a,'when','') or "").strip()
    cta=(getattr(a,'cta','') or "신청·문의: 교회 사무실").strip()
    try: L=int(re.sub(r'[^0-9]','',str(getattr(a,'length','') or "30")) or "30")
    except Exception: L=30
    asp=(getattr(a,'aspect','') or "")
    vertical=("세로" in asp) or ("9:16" in asp)
    W,H=(720,1280) if vertical else (1280,720)
    SRC=os.path.join(BASE,"_내자료","영상소스"); os.makedirs(SRC,exist_ok=True)
    guide=os.path.join(SRC,"여기에 사진·짧은영상·음악을 넣으세요.txt")
    if not os.path.exists(guide):
        open(guide,'w',encoding='utf-8').write(
            "이 폴더에 홍보영상 재료를 넣고 '홍보영상 자동 생성' 카드를 실행하면, 대본 자막과 합쳐 영상(mp4)을 만들어 드립니다.\n\n"
            "· 사진(jpg·png): 교회 전경, 행사 사진, 성도 모습 등 (파일 이름 순서대로 영상에 나옵니다 — 01, 02 … 로 이름 붙이면 순서 조절 쉬움)\n"
            "· 짧은 영상(mp4·mov): 핸드폰으로 찍은 10초 내외 클립 (각 클립은 최대 6초만 사용)\n"
            "· 음악(음악.mp3): 배경음악 파일 이름을 '음악.mp3'로 (저작권 무료 음원 권장)\n"
            "· 자막(자막.txt): 원하시면, 사진/영상 순서대로 넣을 자막을 한 줄에 하나씩. (없으면 대본의 자막을 자동 사용)\n\n"
            "팁: 사진 5~8장 + 짧은 영상 1~2개 + 음악이면 30초짜리 멋진 홍보영상이 됩니다.")
    if not ff:
        print("ℹ 영상 자동 생성 엔진(ffmpeg)이 이 컴퓨터에 없어, 지금은 영상 파일을 바로 만들지 못합니다.")
        print("  1) '행사 홍보영상 기획' 카드의 작업지 + 무료 편집기(Clipchamp·Canva)로 만드실 수 있습니다.")
        print("  2) ffmpeg(무료)를 설치하면, 이 카드가 사진·영상·음악을 합쳐 자동으로 mp4를 만들어 드립니다.")
        return
    exts_img=('jpg','jpeg','png','webp','bmp','gif'); exts_vid=('mp4','mov','m4v','avi','mkv')
    allmedia=sorted([os.path.join(SRC,f) for f in os.listdir(SRC)
                     if '.' in f and f.lower().rsplit('.',1)[-1] in (exts_img+exts_vid)])
    capfile=os.path.join(SRC,"자막.txt")
    if os.path.exists(capfile):
        caps=[l.strip() for l in open(capfile,encoding='utf-8') if l.strip()]
    else:
        caps=[c for c in [message, when, message, "함께해요"] if c]
    music=None
    for m in ("음악.mp3","음악.m4a","음악.wav","music.mp3","music.m4a"):
        if os.path.exists(os.path.join(SRC,m)): music=os.path.join(SRC,m); break
    # 사용자 음악이 없으면 → 분위기(밝은/잔잔한/경건한)에 맞춰 무료(CC0) 배경음악 자동 선곡
    if not music:
        import datetime as _dt
        mood=(getattr(a,'mood','') or "").strip()
        if mood not in ("밝은","잔잔한","경건한"):
            blob=(title+" "+message+" "+(getattr(a,'purpose','') or "")+" "+(getattr(a,'tone','') or "")+" "+cta).replace(" ","")
            if any(k in blob for k in ("여름성경학교","전도","축제","새생명","어린이","청년","수련회","캠프","VBS","활기","경쾌","신나")): mood="밝은"
            elif any(k in blob for k in ("부활","성탄","절기","예배","기도","경건","추수감사","고난","십자가","묵상")): mood="경건한"
            else: mood="잔잔한"
        for md in (os.path.join(BASE,"_음악",mood), os.path.join(BASE,"_음악")):
            if os.path.isdir(md):
                trk=sorted(os.path.join(md,f) for f in os.listdir(md)
                           if '.' in f and f.lower().rsplit('.',1)[-1] in ("mp3","m4a","wav","ogg"))
                if trk: music=trk[_dt.date.today().toordinal()%len(trk)]; break
    FONT=r"C:/Windows/Fonts/malgun.ttf".replace(':',r'\:')
    fs_t=int(W/20); fs_c=int(W/28)
    tmp=tempfile.mkdtemp(); jobs=[]; segs=[]; durs=[]; sc=[0]
    def _cap(txt,n=16):
        k=len(sc); sc.append(1); p=os.path.join(tmp,f"cap{k}.txt")
        open(p,'w',encoding='utf-8').write(_vwrap(txt,n)); return f"cap{k}.txt"
    OUTLINE="borderw=4:bordercolor=black@0.85:shadowcolor=black@0.5:shadowx=2:shadowy=2"
    FADEIN="alpha='if(lt(t\\,0.45)\\,t/0.45\\,1)'"   # 자막이 부드럽게 나타남
    def _dt(cf,i,center=False,big=False):
        st=2 if center else (i%3)
        fs=int(W/15) if (big or st==2) else int(W/21)
        if st==1:   # 하단 좌측 강조
            pos=f"x=80:y=h-th-{int(H/7)}"
        elif st==2: # 중앙 크게(제목·강조)
            pos=f"x=(w-tw)/2:y=(h-th)/2"
        else:       # 하단 중앙
            pos=f"x=(w-tw)/2:y=h-th-{int(H/7)}"
        return (f"drawtext=fontfile='{FONT}':textfile='{cf}':fontcolor=white:fontsize={fs}:"
                f"{pos}:line_spacing=14:{OUTLINE}:{FADEIN}")
    def _run(inputs,vf,dur,outp):   # 실제 렌더는 나중에 병렬로 — 여기선 작업만 모음
        cmd=[ff,'-y']+inputs+['-vf',vf,'-t',str(dur),'-r','30',
             '-pix_fmt','yuv420p','-an','-c:v','libx264','-preset','ultrafast','-crf','23',outp]
        jobs.append((cmd,os.path.basename(outp),float(dur))); return True
    fill=f"scale={W}:{H}:force_original_aspect_ratio=increase,crop={W}:{H},setsar=1"
    PW,PH=int(W*1.16),int(H*1.16)   # 켄 번스(줌+패닝) 여유 프레임 — 가벼워 빠름
    def _ken(i,fr):
        b=f"scale={PW}:{PH}:force_original_aspect_ratio=increase,crop={PW}:{PH},"
        cx="iw/2-(iw/zoom/2)"; cy="ih/2-(ih/zoom/2)"
        P=[  # 6종 순환 — 줌인·좌우패닝·줌아웃·상하패닝으로 파노라마처럼 살아 움직임
          f"zoompan=z='min(zoom+0.0012,1.14)':d={fr}:x='{cx}':y='{cy}'",                       # 줌인
          f"zoompan=z='1.14':d={fr}:x='(iw-iw/zoom)*(on/{fr})':y='{cy}'",                        # 좌→우 패닝
          f"zoompan=z='if(eq(on,0),1.14,max(1.001,zoom-0.0012))':d={fr}:x='{cx}':y='{cy}'",      # 줌아웃
          f"zoompan=z='1.14':d={fr}:x='(iw-iw/zoom)*(1-on/{fr})':y='{cy}'",                       # 우→좌 패닝
          f"zoompan=z='min(zoom+0.0011,1.13)':d={fr}:x='{cx}':y='(ih-ih/zoom)*(on/{fr})'",        # 줌인+아래로
          f"zoompan=z='1.13':d={fr}:x='{cx}':y='(ih-ih/zoom)*(1-on/{fr})'",                       # 위로 패닝
        ]
        return b+P[i%len(P)]+f":s={W}x{H}:fps=30,setsar=1"
    # 인트로(제목) — 중앙 크게
    _run(['-f','lavfi','-i',f'color=c=0x14213D:s={W}x{H}'], _dt(_cap(title,14),0,center=True,big=True), 3, os.path.join(tmp,"seg000.mp4"))
    base=max(3.0, min(7.0, (L-6.5)/max(len(allmedia),1))) if allmedia else 3.5
    vdur=min(8.0, max(4.0, base))
    n=0
    for i,mp in enumerate(allmedia):
        ext=mp.lower().rsplit('.',1)[-1]; cap=caps[i] if i<len(caps) else ""
        if ext=='gif':   # 애니메이션 GIF → 첫 프레임 PNG로 변환(정지 사진처럼 사용)
            gp=os.path.join(tmp,f"gif{i}.png")
            subprocess.run([ff,'-y','-i',mp,'-frames:v','1',gp],cwd=tmp,capture_output=True,text=True)
            if os.path.exists(gp): mp=gp; ext='png'
        cf=_cap(cap,18) if cap else None
        outp=os.path.join(tmp,f"seg{i+1:03d}.mp4")
        capf=(","+_dt(cf,i)) if cf else ""
        if ext in exts_img:
            fr=max(2,int(round(base*30)))
            ok=_run(['-loop','1','-i',mp], _ken(i,fr)+capf, round(base,1), outp)
        else:
            ok=_run(['-i',mp], fill+capf, round(vdur,1), outp)
        if ok: n+=1
    # 아웃트로(제목+신청)
    _run(['-f','lavfi','-i',f'color=c=0x0F2A1F:s={W}x{H}'], _dt(_cap(title+"\n"+cta,16),0,center=True), 3.5, os.path.join(tmp,"seg999.mp4"))
    # ── 세그먼트 병렬 렌더(속도 대폭 개선) — 파일명(seg000·seg001…seg999)이 순서를 담아 정렬로 순서 보존 ──
    import concurrent.futures as _cf
    def _ex(job):
        cmd,bn,dur=job
        subprocess.run(cmd,cwd=tmp,capture_output=True,text=True)
        return bn,dur,os.path.exists(os.path.join(tmp,bn))
    done={}
    try:
        with _cf.ThreadPoolExecutor(max_workers=max(2,min(6,(os.cpu_count() or 4)-1))) as _pool:
            for bn,dur,ok in _pool.map(_ex,jobs):
                if ok: done[bn]=dur
    except Exception:
        for job in jobs:                     # 폴백: 순차 렌더
            bn,dur,ok=_ex(job)
            if ok: done[bn]=dur
    segs=[bn for bn in sorted(done)]; durs=[done[bn] for bn in segs]
    if not segs:
        print("✗ 영상을 만들지 못했습니다. _내자료/영상소스 폴더에 사진을 몇 장 넣고 다시 실행해 주세요."); return
    outdir=CAT("07"); os.makedirs(outdir,exist_ok=True)
    final=os.path.join(outdir,f"[홍보영상] {sanit(title)}_{today()}.mp4")
    # ── 다양한 전환(크로스페이드·슬라이드·디졸브)으로 이어붙임 — PPT 컷 느낌 제거 ──
    T=0.6; TR=['fade','dissolve','slideleft','wiperight','smoothup','circleopen','slideright','fadegrays']
    ins=[]
    for s in segs: ins+=['-i',s]
    music_i=len(segs)
    if music: ins+=['-stream_loop','-1','-i',music]
    fc=[]; prev="[0:v]"; cum=durs[0]
    for i in range(1,len(segs)):
        off=max(0.1,cum-T); lbl=f"[x{i}]"
        fc.append(f"{prev}[{i}:v]xfade=transition={TR[(i-1)%len(TR)]}:duration={T}:offset={off:.2f}{lbl}")
        cum=cum+durs[i]-T; prev=lbl
    fc.append(f"{prev}fade=t=in:st=0:d=0.5,fade=t=out:st={max(0.3,cum-0.8):.2f}:d=0.8[vout]")
    cmd=[ff,'-y']+ins+['-filter_complex']
    if music:
        fc.append(f"[{music_i}:a]afade=t=in:d=1.2,afade=t=out:st={max(0.5,cum-1.5):.2f}:d=1.5[aout]")
        cmd+=[';'.join(fc),'-map','[vout]','-map','[aout]','-c:a','aac','-b:a','192k']
    else:
        cmd+=[';'.join(fc),'-map','[vout]']
    cmd+=['-t',f'{cum:.2f}','-c:v','libx264','-preset','superfast','-crf','23','-pix_fmt','yuv420p',final]
    r=subprocess.run(cmd,cwd=tmp,capture_output=True,text=True)
    if r.returncode!=0 or not os.path.exists(final):
        # 폴백: 전환 실패 시 단순 이어붙이기
        open(os.path.join(tmp,"list.txt"),'w',encoding='utf-8').write("".join(f"file '{s}'\n" for s in segs))
        subprocess.run([ff,'-y','-f','concat','-safe','0','-i','list.txt','-c','copy',final],cwd=tmp,capture_output=True,text=True)
    try: _ext_copy(final)
    except Exception: pass
    try: _openfile(final)
    except Exception: pass
    approx=int(round(3 + n*base + 3.5))
    print(f"✅ 홍보영상 생성 완료: {final}")
    print(f"   {W}x{H} · 약 {approx}초 · 장면 {len(segs)}개(사진·영상 {n}개) · 음악 {'있음' if music else '없음(음악.mp3 넣으면 추가)'}")
    if n==0: print("   ※ _내자료/영상소스 폴더에 사진·영상을 넣으면 훨씬 멋진 영상이 됩니다(지금은 제목·자막만).")
    else: print("   ※ 더 길게(3분+) 만들려면 사진·짧은 영상을 더 많이 넣고 다시 실행하세요.")
def sermon_slides(a):
    """설교 슬라이드(PPT) — 설교 제목·본문 성구·대지를 넣으면 예배 프로젝터·화면공유용 슬라이드를 자동 생성."""
    title=(getattr(a,'title','') or "주일 설교").strip()
    text=(getattr(a,'text','') or "").strip()
    theme=(getattr(a,'theme','') or "").strip()
    points=[re.sub(r'^\s*(\d{1,2}\s*[.)]\s*|[①-⑳]\s*)','',p.strip()) for p in re.split(r'[;\n]', getattr(a,'points','') or "") if p.strip()]
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt as PPt
        from pptx.dml.color import RGBColor as PColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("설교 슬라이드는 python-pptx가 설치된 환경에서 만들어집니다(배포판엔 포함)."); return
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); blank=prs.slide_layouts[6]
    def slide(lines,size=44,sub=None):
        s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=PColor(0x0f,0x12,0x26)
        tb=s.shapes.add_textbox(Inches(0.7),Inches(2.2),Inches(11.9),Inches(3.1)); tf=tb.text_frame; tf.word_wrap=True
        for i,ln in enumerate(lines):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER
            r=p.add_run(); r.text=ln; r.font.size=PPt(size); r.font.bold=True
            r.font.color.rgb=PColor(0xff,0xff,0xff); r.font.name="맑은 고딕"
        if sub:
            p=tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=sub
            r.font.size=PPt(24); r.font.color.rgb=PColor(0xC9,0xA2,0x27); r.font.name="맑은 고딕"
    slide([title], 54, sub=(text or theme or None))
    if text: slide([f"본문 · {text}"], 44)
    for i,pt in enumerate(points,1): slide([f"{i}. {pt}"], 40)
    slide([CHURCH], 44, sub="함께 은혜받는 예배")
    pout=os.path.join(CAT("07"),f"[설교슬라이드] {sanit(title)}_{today()}.pptx"); prs.save(pout); _ext_copy(pout); _openfile(pout)
    print(f"✅ 설교 슬라이드 PPT: {pout} ({2+len(points)}장) — 예배 프로젝터·화면공유용")
    if not points: print("   ※ 설교 대지(;로 구분)를 넣으면 대지별 슬라이드도 함께 만들어집니다.")
def song_catalog(a):
    """찬양 작품집(카탈로그) — 작곡 목사님의 전체 자작곡을 발표순으로 정리(유튜브·앨범·저작권). 사역 유산 보존."""
    import json as _j
    lib=_j.load(open(SONGLIB,encoding='utf-8')) if os.path.exists(SONGLIB) else []
    comp=[s for s in lib if s.get("작곡가")]
    comp.sort(key=lambda s:(s.get("발표일") or "9999"))
    d=newdoc()
    para(d,f"{CHURCH} 찬양 작품집",24,NAVY,True,AL.CENTER,after=2,font=SANS)
    composers=sorted({s["작곡가"] for s in comp if s.get("작곡가")})
    if composers: para(d,"작곡 "+", ".join(composers),12,MAROON,True,AL.CENTER,after=2,font=SANS)
    para(d,f"총 {len(comp)}곡  ·  {today()}",10,LG,AL.CENTER,after=6,font=SANS); hr(d,"A67C1E",12)
    for i,s in enumerate(comp,1):
        para(d,f"{i}. {s['제목']}   [{s.get('key','')}/{s.get('bpm','')}BPM]",13,NAVY,True,before=6,after=2,font=SANS)
        info=[]
        if s.get("작곡가"): info.append("작곡 "+s["작곡가"])
        if s.get("작사"): info.append("작사 "+s["작사"])
        if s.get("발표일"): info.append("발표 "+s["발표일"])
        if s.get("앨범"): info.append("앨범 "+s["앨범"])
        if info: para(d," · ".join(info),10,GRAY,after=1,font=SANS)
        if s.get("주제") or s.get("성경"): para(d,f"주제 {s.get('주제','')}  성경 {s.get('성경','')}",10,GRAY,after=1)
        if s.get("유튜브"): para(d,f"🎬 {s['유튜브']}",10,RGBColor(0x1F,0x5C,0x9E),after=1,font=SANS)
        if s.get("저작권"): para(d,f"저작권 등록: {s['저작권']}",9.5,LG,after=1,font=SANS)
    out=os.path.join(CAT("07"),f"[찬양작품집] {sanit(a.title or CHURCH)}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 찬양 작품집 생성: {out} (자작곡 {len(comp)}곡)")
def songbook(a):
    """찬양집(악보집) 출판 — 선택/전체 곡을 표지·목차·곡별(제목·작곡·key·가사·악보)로 편집·인쇄."""
    import json as _j
    lib=_j.load(open(SONGLIB,encoding='utf-8')) if os.path.exists(SONGLIB) else []
    if a.songs:
        want=[re.sub(r'\s','',x) for x in a.songs.split(";")]
        sel=[s for s in lib if any(w in re.sub(r'\s','',s['제목']) for w in want)]
    else: sel=lib
    if not sel: print("✗ 곡이 없습니다 — '찬양곡·자작곡 등록'에서 등록해 주세요."); return
    d=newdoc()
    para(d,"",after=60); hr(d,"A67C1E",12)
    para(d,a.title or f"{CHURCH} 찬양집",26,NAVY,True,AL.CENTER,after=4,font=SANS)
    para(d,"WORSHIP SONGBOOK",10,LG,AL.CENTER,after=6,font=SANS); hr(d,"A67C1E",12)
    para(d,f"총 {len(sel)}곡  ·  {today()}",10,LG,AL.CENTER,before=8,font=SANS)
    d.add_page_break()
    para(d,"목  차",18,NAVY,True,AL.CENTER,after=6,font=SANS); hr(d,"1B2A4A",6)
    for i,s in enumerate(sel,1): para(d,f"{i}. {s['제목']}   [{s.get('key','')}]",12,TEAL if False else GRAY,after=3,font=SANS)
    imgw=Cm(16)
    for i,s in enumerate(sel,1):
        d.add_page_break()
        para(d,f"{i}. {s['제목']}",17,NAVY,True,AL.CENTER,after=2,font=SANS)
        sub=[x for x in [f"작곡 {s['작곡가']}" if s.get('작곡가') else "", f"[{s.get('key','')}/{s.get('bpm','')}BPM]" if s.get('key') else ""] if x]
        if sub: para(d," · ".join(sub),10,GRAY,AL.CENTER,after=4,font=SANS)
        if s.get("악보") and os.path.exists(s["악보"]) and s["악보"].lower().endswith(('.jpg','.jpeg','.png')):
            try: d.add_picture(s["악보"],width=imgw); d.paragraphs[-1].alignment=AL.CENTER
            except Exception: pass
        elif s.get("가사"):
            for ln in s["가사"].split("\n"): para(d,ln,13,GRAY,align=AL.CENTER,after=2)
    out=os.path.join(CAT("07"),f"[찬양집] {sanit(a.title or '찬양집')}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 찬양집(악보집) 출판: {out} ({len(sel)}곡)")
def lyrics_screen(a):
    """찬양 가사 스크린 — 예배 프로젝터용 큰글씨 가사 슬라이드(등록곡 가사 사용)."""
    import json as _j
    from docx.enum.section import WD_ORIENT
    lib=_j.load(open(SONGLIB,encoding='utf-8')) if os.path.exists(SONGLIB) else []
    byn={re.sub(r'\s','',s['제목']):s for s in lib}
    songs=[s.strip() for s in (a.songs or "").split(";") if s.strip()]
    d=Document()
    sec=d.sections[0]; sec.orientation=WD_ORIENT.LANDSCAPE
    sec.page_width,sec.page_height=Cm(33.87),Cm(19.05)
    sec.top_margin=Cm(2.0); sec.bottom_margin=Cm(2.0); sec.left_margin=Cm(2.5); sec.right_margin=Cm(2.5)
    miss=[]
    for gi,q in enumerate(songs):
        nk=re.sub(r'\s','',q); s=None
        for k,v in byn.items():
            if nk in k or k in nk: s=v; break
        if gi>0: d.add_page_break()
        title=s['제목'] if s else q
        para(d,title,32,NAVY,True,AL.CENTER,before=60,after=10,font=SANS)
        lyr=(s.get('가사') if s else '') or ''
        if not lyr: miss.append(q); para(d,"(가사 미등록 — '찬양곡·자작곡 등록'에서 가사 추가)",16,MAROON,align=AL.CENTER); continue
        for block in re.split(r'\n\s*\n',lyr):
            d.add_page_break()
            for ln in block.split("\n"):
                para(d,ln.strip(),40,NAVY,True,AL.CENTER,before=6,after=6,font=SANS)
    out=os.path.join(CAT("07"),f"[가사스크린] {sanit(a.title or '찬양')}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 가사 스크린 생성: {out}"+(f" · 가사없음 {len(miss)}곡: {', '.join(miss)}" if miss else ""))
def worship_roster(a):
    """찬양팀 배정표 — 파트별 배정(인도·건반·기타·베이스·드럼·싱어) + 카톡 리마인드."""
    d=newdoc()
    para(d,f"{CHURCH} 찬양팀 배정표",20,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,f"{a.date or today()}",11,GRAY,AL.CENTER,after=6,font=SANS); hr(d,"A67C1E",12)
    parts=[("인도",a.leader),("건반",a.keys),("기타",a.guitar),("베이스",a.bass),("드럼",a.drums),("싱어",a.singers)]
    t=d.add_table(rows=len(parts),cols=2)
    for i,(role,who) in enumerate(parts):
        ctext(t.rows[i].cells[0],role,12,WHITE,True,SANS,AL.CENTER); shade(t.rows[i].cells[0],"3C6E3C")
        ctext(t.rows[i].cells[1],"  "+(who or ""),12,GRAY,font=SERIF)
        t.rows[i].cells[0].width=Cm(3.5); t.rows[i].cells[1].width=Cm(12)
    tborders(t,"BBBBBB",4)
    if a.songs: para(d,"콘티: "+a.songs.replace(";"," · "),11,MAROON,before=8,font=SANS)
    out=os.path.join(CAT("07"),f"[찬양팀배정] {a.date or today()}.docx"); _savedoc(d,out)
    print(f"✅ 찬양팀 배정표: {out}")
    who=[f"{r} {w}" for r,w in parts if w]
    print(f"  카톡 리마인드 초안: [찬양팀] {a.date or today()} 배정 — "+", ".join(who)+". 예배 30분 전 리허설 부탁드립니다 🙏")

ILBANK=_C.get("예화은행경로") or os.path.join(BASE,"_내자료","예화은행.json")
COMDIR=os.path.join(BASE,"_내자료","성경주석")   # 목사님이 주석 PDF를 넣는 폴더
def illus_add(a):
    """예화 추가 — 목사님 자신의 예화은행에 저장(설교·묵상에서 자동 추천)."""
    import json as _j
    os.makedirs(os.path.dirname(ILBANK),exist_ok=True)
    bank=_j.load(open(ILBANK,encoding='utf-8')) if os.path.exists(ILBANK) else []
    if isinstance(bank,dict): bank=bank.get("예화",[])
    bank.append({"예화명":a.name or "","출처":a.source or "","주제키워드":[k.strip() for k in (a.topic or "").split(";") if k.strip()],
                 "성경연결":a.bible or "","전문":a.text or ""})
    _j.dump(bank,open(ILBANK,'w',encoding='utf-8'),ensure_ascii=False,indent=2)
    print(f"✅ 예화 추가: {a.name} → {ILBANK} (총 {len(bank)}건)")
def illus_list(a):
    import json as _j
    if not os.path.exists(ILBANK): print("예화 없음. 먼저 예화 자료를 입력해 주세요."); return
    bank=_j.load(open(ILBANK,encoding='utf-8')); bank=bank if isinstance(bank,list) else bank.get("예화",[])
    print(f"■ 내 예화은행 {len(bank)}건")
    for it in bank: print(f"  · {it.get('예화명','')} ({it.get('출처','')}) — {it.get('성경연결','')}")
_BIBLE_ABBR={"창":"창세기","출":"출애굽기","레":"레위기","민":"민수기","신":"신명기","수":"여호수아","삿":"사사기","룻":"룻기",
 "삼상":"사무엘상","삼하":"사무엘하","왕상":"열왕기상","왕하":"열왕기하","대상":"역대상","대하":"역대하","스":"에스라","느":"느헤미야","에":"에스더",
 "욥":"욥기","시":"시편","잠":"잠언","전":"전도서","아":"아가","사":"이사야","렘":"예레미야","애":"예레미야애가","겔":"에스겔","단":"다니엘",
 "호":"호세아","욜":"요엘","암":"아모스","옵":"오바댜","욘":"요나","미":"미가","나":"나훔","합":"하박국","습":"스바냐",
 "학":"학개","슥":"스가랴","말":"말라기","마":"마태복음","막":"마가복음","눅":"누가복음","요":"요한복음","행":"사도행전",
 "롬":"로마서","고전":"고린도전서","고후":"고린도후서","갈":"갈라디아서","엡":"에베소서","빌":"빌립보서","골":"골로새서",
 "살전":"데살로니가전서","살후":"데살로니가후서","딤전":"디모데전서","딤후":"디모데후서","딛":"디도서","몬":"빌레몬서",
 "히":"히브리서","약":"야고보서","벧전":"베드로전서","벧후":"베드로후서","유":"유다서","계":"요한계시록"}
def _nlm_book(book):
    """책 이름 → NotebookLM 노트북ID (설정 'nlm책노트북'). 축약형·부분일치 허용."""
    m=_C.get("nlm책노트북") or {}
    if not book or not m: return None
    b=str(book).strip()
    if b in m: return m[b]
    full=_BIBLE_ABBR.get(b,b)
    if full in m: return m[full]
    for k in m:
        if full and (full in k or k.startswith(full)): return m[k]
    return None
def commentary_add(a):
    """성경주석 PDF 추가 — 내 주석 라이브러리에 등록 + (설정 시) NotebookLM 업로드. 계속 쌓아 설교가 풍성해짐."""
    import shutil, subprocess
    os.makedirs(COMDIR,exist_ok=True); src=a.pdf
    if not src or not os.path.exists(src): print(f"✗ PDF 경로 확인: {src}"); return
    dst=os.path.join(COMDIR,os.path.basename(src))
    if os.path.abspath(src)!=os.path.abspath(dst): shutil.copy2(src,dst)
    _ext_copy(dst)   # 주석 PDF USB·D 이중저장
    print(f"✅ 주석 등록: {os.path.basename(dst)}")
    nb=a.notebook or _nlm_book(getattr(a,'book',None)) or _C.get("nlm노트북")
    if nb:
        try:
            r=subprocess.run(["nlm","source","add",nb,"--file",dst,"--wait"],capture_output=True,text=True,timeout=300)
            if r.returncode==0: print(f"  ☁ NotebookLM 업로드 완료 · 노트북 {nb}")
            else: print(f"  ⚠ NotebookLM 업로드 실패: {(r.stderr or r.stdout or '').strip()[:90]}\n     (로그인이 안 돼 있으면 → nlm login 먼저 실행)")
        except FileNotFoundError: print("  ※ NotebookLM 연동 도구(nlm)가 없어 로컬 등록만 했습니다. 웹(notebooklm.google.com)에서 올려도 됩니다.")
        except Exception as e: print(f"  ⚠ NotebookLM 미연동: {str(e)[:60]}")
    else:
        print("  ※ NotebookLM 자동업로드: church_config.json에 \"nlm노트북\":\"<ID>\" 설정 또는 --notebook 지정")
    cnt=len([f for f in os.listdir(COMDIR) if f.lower().endswith('.pdf')])
    print(f"  주석 라이브러리 현재 {cnt}개 — 설교 시 자동 참고")
def commentary_list(a):
    pdfs=[f for f in os.listdir(COMDIR) if f.lower().endswith('.pdf')] if os.path.isdir(COMDIR) else []
    print(f"■ 성경주석 라이브러리 {len(pdfs)}개 ({COMDIR})")
    for p in pdfs: print("  ·",p)
def nlm_add(a):
    """NotebookLM에 자료 올리기 — 성경주석·찬양곡·가사·자료 등 어떤 파일이든 내 NotebookLM 노트북에 올려 검색에 활용."""
    import subprocess
    f=(a.file or "").strip().strip('"').strip("'")
    if f and not os.path.exists(f):
        # 전체 경로가 안 맞으면 파일 '이름'으로 흔한 폴더에서 찾아준다(경로 타이핑 실수 방지)
        want=os.path.basename(f).lower()
        for d in [os.path.join(BASE,"_내자료"), os.path.join(os.path.expanduser("~"),"Desktop"), os.path.join(os.path.expanduser("~"),"Downloads"), ROOT]:
            if not os.path.isdir(d): continue
            hit=None
            for root,_,files in os.walk(d):
                for fn in files:
                    if fn.lower()==want or (len(want)>=2 and want in fn.lower()): hit=os.path.join(root,fn); break
                if hit: break
            if hit: f=hit; print(f"  (파일 찾음: {f})"); break
    if not f or not os.path.exists(f):
        print("✗ 파일을 찾지 못했습니다. 올리실 파일을  _시스템\\_내자료  폴더나 바탕화면에 넣고, '올릴 파일 경로' 칸에 파일 이름을 넣어 주세요.")
        nd=os.path.join(BASE,"_내자료"); avail=[]
        if os.path.isdir(nd):
            for root,_,files in os.walk(nd):
                for fn in files:
                    if not fn.startswith("여기에") and fn.lower().endswith((".pdf",".docx",".txt",".hwp",".md",".pptx")): avail.append(fn)
        print("  올릴 수 있는 파일: "+(", ".join(avail[:12]) if avail else "(없음 — 먼저 파일을 _내자료 폴더에 복사해 넣으세요)")); return
    nb=a.notebook or _nlm_book(getattr(a,'book',None)) or _C.get("nlm노트북")
    if not nb:
        print("✗ 노트북ID가 필요합니다 — --notebook 노트북ID 로 지정하거나 church_config.json에 \"nlm노트북\":\"<ID>\" 설정.")
        print("  NotebookLM 무료 가입: notebooklm.google.com → 새 노트북 → 주소창 끝의 노트북ID 사용."); return
    try:
        r=subprocess.run(["nlm","source","add",nb,"--file",f,"--wait"],capture_output=True,text=True,timeout=550)
        if r.returncode==0:
            print(f"✅ NotebookLM 업로드 완료: {os.path.basename(f)} · 노트북 {nb}")
            print("   이제 NotebookLM에서 이 자료를 검색·질문해 참고하실 수 있습니다.")
        else:
            print(f"⚠ 업로드 실패: {(r.stderr or r.stdout or '').strip()[:90]}\n   (로그인이 안 돼 있으면 → nlm login 먼저 실행하세요)")
    except FileNotFoundError:
        print("⚠ NotebookLM 연동 도구(nlm)가 설치돼 있지 않습니다.")
        print("   → 더 간단하게: notebooklm.google.com 웹사이트에서 이 파일을 직접 '소스 추가'로 올리셔도 됩니다.")
    except Exception as e:
        print(f"⚠ 업로드 오류: {str(e)[:80]}")
def _match_illus(keywords):
    import json as _j
    if not os.path.exists(ILBANK): return []
    try:
        bank=_j.load(open(ILBANK,encoding='utf-8'))
        items=bank if isinstance(bank,list) else bank.get("예화",[])
    except Exception: return []
    out=[]
    for it in items:
        tv=it.get("주제키워드",[])
        tags=" ".join(tv) if isinstance(tv,list) else str(tv)
        blob=tags+" "+str(it.get("성경연결",""))+" "+str(it.get("주제",""))
        if any(k and k in blob for k in keywords.split()): out.append(it)
    return out
def _illus_of_day(datestr=None):
    """날짜별 순환 예화 — 그날 날짜로 예화은행에서 한 편을 뽑는다.
    ordinal(연·월·일 절대일수) % 총편수 → 매일 다르고, 해가 바뀌어도 같은 날짜가 다른 예화가 됨
    (365·366이 편수의 배수가 아니라 매년 위치가 밀림). 전용판=채굴 예화, 배포판=목사님이 넣으신 예화. 비면 None."""
    import json as _j, datetime as _dt
    if not os.path.exists(ILBANK): return None
    try:
        bank=_j.load(open(ILBANK,encoding='utf-8'))
        items=bank if isinstance(bank,list) else bank.get("예화",[])
    except Exception: return None
    if not items: return None
    try: od=_dt.date.fromisoformat(str(datestr or today())[:10]).toordinal()
    except Exception: od=_dt.date.today().toordinal()
    return items[od % len(items)]
def sermon(a):
    """설교문 작성 — 설교체(장면→진리→그리스도) 초안. 예화은행·성경주석(PDF/NotebookLM) 연동."""
    import json as _j
    points=[p.strip() for p in (a.points or "").split(";") if p.strip()] or ["대지 1","대지 2","대지 3"]
    # 예화은행에서 주제 매칭
    illus=_match_illus((a.theme or "")+" "+(a.title or "")+" "+(a.text or ""))
    d=newdoc()
    svc=a.service or "주일오전예배"
    para(d,f"{CHURCH} · {svc}",11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,a.title or "설교 제목",22,NAVY,True,AL.CENTER,before=4,after=2,font=SANS); hr(d,"8B2E2E",12)
    para(d,f"본문 {a.text or ''}   ·   {a.date or today()}   ·   담임 "+PASTOR,10,GRAY,AL.CENTER,after=8,font=SANS)
    para(d,"서론 — 마음 열기",14,MAROON,True,before=4,after=3,font=SANS)
    para(d,"  (도입 예화/질문으로 본문의 문을 연다)",10.5,LG,italic=True,after=8)
    for i,pt in enumerate(points,1):
        para(d,f"대지 {i}.  {pt}",14,NAVY,True,before=8,after=4,font=SANS)
        para(d,"① 장면 — ",11,MAROON,True,after=2); para(d,"   (본문/예화의 구체적 장면을 그린다)",10.5,LG,italic=True,after=4)
        para(d,"② 진리 — ",11,MAROON,True,after=2); para(d,"   (그 장면이 드러내는 성경의 진리)",10.5,LG,italic=True,after=4)
        para(d,"③ 그리스도 — ",11,MAROON,True,after=2); para(d,"   (진리가 그리스도와 십자가로 귀결)",10.5,LG,italic=True,after=6)
    para(d,"결론 — 적용과 결단",14,MAROON,True,before=8,after=3,font=SANS)
    para(d,"  (한 문장 핵심 + 삶의 적용 + 결단 촉구)",10.5,LG,italic=True,after=8)
    if illus:
        hr(d,"CCCCCC",4); para(d,"◈ 예화은행 추천 (주제 매칭)",12,GOLD,True,before=4,after=4,font=SANS)
        for it in illus[:5]:
            para(d,f"· {it.get('예화명','')}  ({it.get('출처','')})",11,NAVY,True,after=1)
            para(d,f"   성경연결: {it.get('성경연결','')}",10,GRAY,after=4)
    # 목사님 자신의 성경주석 PDF 목록(참고자료)
    pdfs=[f for f in os.listdir(COMDIR) if f.lower().endswith('.pdf')] if os.path.isdir(COMDIR) else []
    if pdfs:
        para(d,"◈ 참고 성경주석 (내 자료 PDF)",11.5,GOLD,True,before=8,after=3,font=SANS)
        for p in pdfs[:10]: para(d,f"· {p}",10,GRAY,after=1,font=SANS)
    if pdfs or _C.get("nlm노트북") or _C.get("주석연동"):   # 주석 PDF·NotebookLM·연동설정 있을 때만(배포판엔 없으니 미표시)
        pass
    folder=os.path.join(CAT("09"),sanit(svc)); os.makedirs(folder,exist_ok=True)  # 예배유형별 폴더
    _t=sanit(a.title or '제목없음'); _b=sanit((a.text or '').strip())
    out=os.path.join(folder,f"[{sanit(svc)}] {_t}"+(f" ({_b})" if _b else "")+f"_{a.date or today()}.docx"); _savedoc(d,out)
    _ext_copy(out)   # USB·D 외부 이중보관(설정 시). HWP는 별도 '한글 변환' 카드로(웹 지연 방지)
    # 설교 이력 자동 기록 + 같은 본문 과거 설교 경고(반복 방지)
    db=load(); db.setdefault("설교",[])
    prev=[s for s in db["설교"] if a.text and s.get("본문") and (a.text in s["본문"] or s["본문"] in a.text)]
    db["설교"].append({"날짜":a.date or today(),"제목":a.title or "","본문":a.text or "","예배":svc,"시리즈":a.series or "","대지":a.points or ""})
    db["설교"].sort(key=lambda s:s.get("날짜","")); save(db)
    print(f"✅ 설교 초안 생성(DOCX): {out}")
    print(f"   대지 {len(points)}개 · 예화은행 추천 {len(illus)}건 · 설교이력 자동기록")
    print("   ※ 한글(.hwp)로 원하시면 '한글 변환' 카드를 눌러주세요. (DOCX도 한글에서 그대로 열립니다)")
    if prev: print("   ⚠ 이 본문으로 전에 설교함: "+", ".join(f"{s['날짜']} '{s['제목']}'" for s in prev[:3]))

BIBLEDIR=os.path.join(BASE,"bible")
BOOKALIAS={"창":"창세기","출":"출애굽기","시":"시편","잠":"잠언","사":"이사야","마":"마태복음","막":"마가복음",
 "눅":"누가복음","요":"요한복음","행":"사도행전","롬":"로마서","고전":"고린도전서","고후":"고린도후서",
 "갈":"갈라디아서","엡":"에베소서","빌":"빌립보서","골":"골로새서","히":"히브리서","약":"야고보서","계":"요한계시록"}
def bible(a):
    """성경본문 불러오기 — bible/<역본>.json 에서. 역본: 개역개정 niv kjv"""
    import json as _j
    ver=a.version or "개역개정"; path=os.path.join(BIBLEDIR,ver+".json")
    if not os.path.exists(path):
        avail=[f[:-5] for f in os.listdir(BIBLEDIR)] if os.path.isdir(BIBLEDIR) else []
        print(f"✗ '{ver}' 성경 데이터 없음. 보유 역본: {avail or '없음'}")
        print(f"  → 성경 텍스트 파일을 {BIBLEDIR}\\{ver}.json 형식으로 넣으면 즉시 작동합니다.")
        print(f"  형식: {{\"요한복음\": {{\"3\": {{\"16\": \"하나님이 세상을...\"}}}}}}"); return
    bib=_j.load(open(path,encoding='utf-8'))
    # 관대한 인식 — 장/절/전각콜론/물결/띄어쓰기 다 허용. 예: 요한복음3:16 · 요 3:16-18 · 요한복음 3장 16절 · 요한복음 3(장 전체)
    ref=(a.ref or "").strip().replace("：",":").replace("．",":").replace("장",":").replace("절","").replace("~","-").replace(" -","-").replace("- ","-")
    m=re.match(r'\s*([1-9]?[가-힣]+|[A-Za-z]+)\s*[:：]?\s*(\d+)\s*(?::\s*(\d+)(?:\s*-\s*(\d+))?)?',ref)
    if not m: print("✗ 참조 형식: '요한복음 3:16' · '요 3:16-17' · '요한복음 3'(장 전체)"); return
    book=BOOKALIAS.get(m.group(1),m.group(1)); ch=m.group(2)
    chap=bib.get(book,{}).get(ch,{})
    if not chap: print(f"✗ {book} {ch}장 본문 없음(역본·책이름 확인)"); return
    if m.group(3):   # 절 지정
        v1=int(m.group(3)); v2=int(m.group(4) or m.group(3)); verses=list(range(v1,v2+1)); label=f"{ch}:{v1}{'-'+str(v2) if v2!=v1 else ''}"
    else:            # 장 전체
        vs=sorted(int(k) for k in chap.keys() if str(k).isdigit()); verses=vs; label=f"{ch}장 전체({len(vs)}절)"
    out=[f"{v} {chap[str(v)]}" for v in verses if str(v) in chap]
    if not out: print(f"✗ {book} {ch} 본문 없음(데이터 확인)"); return
    print(f"【{book} {label} · {ver}】")
    for line in out: print("  "+line)

# ───────── 성경자료 (v2 batch B) ─────────
BOOKS66=[("창세기",50),("출애굽기",40),("레위기",27),("민수기",36),("신명기",34),
 ("여호수아",24),("사사기",21),("룻기",4),("사무엘상",31),("사무엘하",24),
 ("열왕기상",22),("열왕기하",25),("역대상",29),("역대하",36),("에스라",10),
 ("느헤미야",13),("에스더",10),("욥기",42),("시편",150),("잠언",31),
 ("전도서",12),("아가",8),("이사야",66),("예레미야",52),("예레미야애가",5),
 ("에스겔",48),("다니엘",12),("호세아",14),("요엘",3),("아모스",9),
 ("오바댜",1),("요나",4),("미가",7),("나훔",3),("하박국",3),
 ("스바냐",3),("학개",2),("스가랴",14),("말라기",4),
 ("마태복음",28),("마가복음",16),("누가복음",24),("요한복음",21),("사도행전",28),
 ("로마서",16),("고린도전서",16),("고린도후서",13),("갈라디아서",6),("에베소서",6),
 ("빌립보서",4),("골로새서",4),("데살로니가전서",5),("데살로니가후서",3),("디모데전서",6),
 ("디모데후서",4),("디도서",3),("빌레몬서",1),("히브리서",13),("야고보서",5),
 ("베드로전서",5),("베드로후서",3),("요한일서",5),("요한이서",1),("요한삼서",1),
 ("유다서",1),("요한계시록",22)]  # 개신교 66권 표준 장수(합계 1189 — 코드에서 검증)
def BIBLECAT():
    """성경자료(통독표·퀴즈·암송) 저장 폴더."""
    ar=_C.get("아카이브루트")
    if ar and os.path.isdir(ar):
        real=os.path.join(ar,"성경자료")
        try: os.makedirs(real,exist_ok=True); return real
        except Exception: pass
    d=os.path.join(ROOT,"성경자료"); os.makedirs(d,exist_ok=True); return d
def _fmt_range(seg):
    """연속 장 목록 → '창세기 1-3' / '창세기 48 ~ 출애굽기 2' 형식."""
    (b1,c1),(b2,c2)=seg[0],seg[-1]
    if b1==b2: return f"{b1} {c1}" if c1==c2 else f"{b1} {c1}-{c2}"
    return f"{b1} {c1} ~ {b2} {c2}"
def bible_plan(a):
    """📖 성경 통독표 — 유형(1년/90일/맥체인)·시작일·부서. 날짜별 읽을 범위+진도 체크칸(docx). 본문 없이 범위만 자체 계산."""
    typ=(a.type or "1년").strip()
    try:
        start=datetime.date.fromisoformat(a.date) if a.date else datetime.date.today()
    except ValueError:
        print("⚠ 시작일 형식이 올바르지 않습니다(예: 2026-01-01). 오늘 날짜로 진행합니다."); start=datetime.date.today()
    seq=[(b,c) for b,n in BOOKS66 for c in range(1,n+1)]
    if len(seq)!=1189: print(f"✗ 내부 오류: 장수 합계가 {len(seq)}(1189이어야 함) — 데이터 확인 필요"); return
    total=len(seq); dept=a.dept or ""
    d=newdoc()
    para(d,CHURCH,11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,f"성경 통독표 · {typ}",22,NAVY,True,AL.CENTER,before=4,after=2,font=SANS)
    para(d,f"시작 {start.isoformat()}"+(f" · {dept}" if dept else ""),11,MAROON,True,AL.CENTER,after=8,font=SANS); hr(d,"1F5C9E",12)
    if typ.startswith("맥체인"):
        streams=[seq[i::4] for i in range(4)]; days=max(len(s) for s in streams)
        para(d,f"맥체인식 4갈래 병행 통독(자체 구성) — 하루 4장 · 총 {days}일 · 66권 전체",10,GRAY,after=6,font=SANS)
        t=d.add_table(rows=days+1,cols=6)
        for j,h in enumerate(["일차","날짜","①","②","③","④"]): ctext(t.rows[0].cells[j],h,10,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],"1F5C9E")
        for i in range(days):
            ctext(t.rows[i+1].cells[0],str(i+1),9,GRAY,font=SANS,align=AL.CENTER)
            ctext(t.rows[i+1].cells[1],(start+datetime.timedelta(days=i)).isoformat(),9,GRAY,font=SANS,align=AL.CENTER)
            for k in range(4):
                v=f"{streams[k][i][0]} {streams[k][i][1]}" if i<len(streams[k]) else ""
                ctext(t.rows[i+1].cells[2+k],v,9,NAVY,font=SERIF)
        tborders(t)
    else:
        ndays={"1년":365,"90일":90,"180일":180,"30일":30,"1개월":30}.get(typ)
        if ndays is None:
            digits=re.sub(r'\D','',typ); ndays=int(digits) if digits else 365
        ndays=max(1,ndays); base=total//ndays; rem=total%ndays
        para(d,f"하루 약 {base}~{base+1}장 · 총 {ndays}일 · {total}장 전체",10,GRAY,after=6,font=SANS)
        t=d.add_table(rows=ndays+1,cols=4)
        for j,h in enumerate(["일차","날짜","읽을 범위","✓"]): ctext(t.rows[0].cells[j],h,10,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],"1F5C9E")
        idx=0
        for i in range(ndays):
            cnt=base+(1 if i<rem else 0); seg=seq[idx:idx+cnt]; idx+=cnt
            ctext(t.rows[i+1].cells[0],str(i+1),9,GRAY,font=SANS,align=AL.CENTER)
            ctext(t.rows[i+1].cells[1],(start+datetime.timedelta(days=i)).isoformat(),9,GRAY,font=SANS,align=AL.CENTER)
            ctext(t.rows[i+1].cells[2],_fmt_range(seg) if seg else "-",10,NAVY,font=SERIF)
            ctext(t.rows[i+1].cells[3],"☐",11,LG,font=SANS,align=AL.CENTER)
        tborders(t)
    para(d,f"· {CHURCH} · 발급일 {today()}",9,LG,align=AL.RIGHT,before=8,font=SANS)
    out=os.path.join(BIBLECAT(),f"[성경통독표] {typ}{('_'+sanit(dept)) if dept else ''}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 성경 통독표({typ}) 생성: {out}  (총 {total}장)")

def _quiz_words(t):
    """본문에서 빈칸/보기용 실단어 추출 — 한글 어절은 2자↑, 영문은 4자↑."""
    out=[]
    for w in t.split():
        c=w.strip('.,;:!?"\'()[]')
        if not c.isalpha(): continue
        han=any('가'<=ch<='힣' for ch in c)
        if len(c)>=(2 if han else 4): out.append(c)
    return out
def _blank_word(text, word):
    """어절 단위 빈칸 치환(한글·영문 공통, \\b 미의존)."""
    toks=text.split()
    for i,tk in enumerate(toks):
        if tk.strip('.,;:!?"\'()[]')==word:
            toks[i]=tk.replace(word,"( ________ )",1); return " ".join(toks)
    return text.replace(word,"( ________ )",1)
def _bible_ver(requested):
    """L 성경본문 2갈래 — 개역개정(한국어) 우선. requested 있으면 그대로. 없으면 개역개정.json 있으면 개역개정, 없으면 (None,None)=주소중심."""
    if requested:
        p=os.path.join(BIBLEDIR,requested+".json")
        return (requested, p if os.path.exists(p) else None)
    kp=os.path.join(BIBLEDIR,"개역개정.json")
    if os.path.exists(kp): return ("개역개정", kp)
    return (None, None)
def bible_quiz(a):
    """📝 성경 퀴즈 — L분기: 개역개정.json 있으면 한국어 본문, 없으면 주소중심(성구주소만·본문은 본인 성경으로). --version web/kjv로 공개영문 명시 가능. 원문복사X·답은 본문에서 검증."""
    import random
    if not a.ref: print("⚠ 성경 범위를 입력해 주세요. (예: bible-quiz --ref '요한복음 3' --count 10 --type 빈칸)"); return
    ver, path = _bible_ver(a.version); addr_mode=(path is None)
    struct_path = path or os.path.join(BIBLEDIR,"web.json")   # 절 구조는 web에서라도(주소중심)
    if not os.path.exists(struct_path): print("✗ 성경 구조 데이터가 없습니다(bible/web.json 확인)."); return
    bib=json.load(open(struct_path,encoding='utf-8'))
    m=re.match(r'\s*([1-9]?[가-힣]+)\s*(?:(\d+)(?::(\d+)(?:-(\d+))?)?)?\s*$', a.ref)
    if not m: print("✗ 범위 형식: '요한복음 3' 또는 '요한복음 3:1-16' 또는 '요한복음'"); return
    book=BOOKALIAS.get(m.group(1),m.group(1))
    if book not in bib: print(f"✗ '{book}' 없음(성경 데이터 확인)"); return
    verses=[]
    if m.group(2) is None:
        for ch,vs in bib[book].items():
            for v,t in vs.items(): verses.append((f"{book} {ch}:{v}", t))
    elif m.group(3) is None:
        ch=m.group(2)
        for v,t in bib[book].get(ch,{}).items(): verses.append((f"{book} {ch}:{v}", t))
    else:
        ch=m.group(2); v1=int(m.group(3)); v2=int(m.group(4) or m.group(3))
        for v in range(v1,v2+1):
            t=bib[book].get(ch,{}).get(str(v))
            if t: verses.append((f"{book} {ch}:{v}", t))
    if not verses: print(f"✗ 범위에 본문이 없습니다: {a.ref}"); return
    n=min(_intarg(a.count,10,"문항수"), len(verses)); picked=random.sample(verses,n)
    typ=a.type or "빈칸"; level=a.level or ""; dept=a.dept or ""
    questions=[]
    if addr_mode:   # 주소중심 — 본문 미포함(저작권), 성구를 본인 성경에서 찾아 쓰기
        typ="성구찾아쓰기"
        for ref,_ in picked: questions.append((f"{ref} 말씀을 성경에서 찾아 아래에 적으세요.", ref))
        srclabel="주소중심(본문 미포함 · 본인 성경으로 채움)"
    else:
        pool=list({w for _,t in picked for w in _quiz_words(t)})
        for ref,text in picked:
            cws=_quiz_words(text)
            if typ.startswith("성구") or not cws:
                questions.append((text, ref)); continue
            word=max(cws,key=len); blanked=_blank_word(text,word)
            if typ.startswith("사지") or typ.startswith("4"):
                others=[w for w in pool if w.lower()!=word.lower()]; random.shuffle(others)
                opts=[word]+others[:3]; random.shuffle(opts); lab="①②③④"
                qtext=blanked+"\n    "+"   ".join(f"{lab[i]} {o}" for i,o in enumerate(opts))
                questions.append((qtext, f"{lab[opts.index(word)]} {word}  ({ref})"))
            else:
                questions.append((blanked, f"{word}  ({ref})"))
        srclabel=f"{ver}" + (" (개역개정 한국어)" if ver=="개역개정" else " (공개역본·저작권 자유)")
    d=newdoc()
    para(d,CHURCH,11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,"성경 퀴즈",22,NAVY,True,AL.CENTER,before=4,after=2,font=SANS)
    meta=" · ".join(x for x in [dept,(f"난이도 {level}" if level else ""),f"범위 {a.ref}",(ver or "주소중심"),f"{n}문항",typ] if x)
    para(d,meta,10,MAROON,True,AL.CENTER,after=8,font=SANS); hr(d,"1F5C9E",12)
    para(d,"■ 문제",13,NAVY,True,after=4,font=SANS)
    for i,(q,_) in enumerate(questions,1):
        para(d,f"{i}. {q}",11,GRAY,after=2,font=SERIF)
        para(d,"   답: ________________________",10,LG,after=7,font=SANS)
    d.add_page_break()
    para(d,"■ 정답지",16,MAROON,True,AL.CENTER,after=6,font=SANS); hr(d,"8B2E2E",10)
    for i,(_,ans) in enumerate(questions,1): para(d,f"{i}. {ans}",11,NAVY,after=3,font=SERIF)
    para(d,f"· {CHURCH} · {today()} · 본문 출처: {srclabel}",9,LG,align=AL.RIGHT,before=8,font=SANS)
    out=os.path.join(BIBLECAT(),f"[성경퀴즈] {sanit(a.ref)}_{typ}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 성경 퀴즈 생성({n}문항·{typ}·{ver or '주소중심'}): {out}")

def memory_verse(a):
    """💳 성경 암송표·지갑카드 — 성구주소 목록(--refs '요한복음 3:16; 로마서 8:28')으로 암송표+지갑용 카드 생성. 본문은 공개역본(원문 저작권 회피)·세트명은 라벨(저작권 큐레이션 목록 내장 안 함)."""
    if not a.refs: print("⚠ 암송할 성구 주소를 입력해 주세요. (예: memory-verse --refs '요한복음 3:16; 로마서 8:28; 시편 23:1' --set 구원)"); return
    ver, path = _bible_ver(a.version); addr_mode=(path is None)   # L: 개역개정 없으면 주소중심
    bib=json.load(open(path,encoding='utf-8')) if path else {}
    items=[]
    for raw in re.split(r'[;\n]', a.refs):
        raw=raw.strip()
        if not raw: continue
        mm=re.match(r'\s*([1-9]?[가-힣]+)\s*(\d+):(\d+)(?:-(\d+))?', raw)
        if not mm: items.append((raw,"(주소 형식 확인: 요한복음 3:16)")); continue
        book=BOOKALIAS.get(mm.group(1),mm.group(1)); ch=mm.group(2); v1=int(mm.group(3)); v2=int(mm.group(4) or mm.group(3))
        disp=f"{book} {ch}:{v1}"+(f"-{v2}" if v2!=v1 else "")
        if addr_mode:
            items.append((disp,"(성경에서 찾아 적으세요)"))
        else:
            parts=[bib.get(book,{}).get(ch,{}).get(str(v)) for v in range(v1,v2+1)]; parts=[p for p in parts if p]
            items.append((disp," ".join(parts) if parts else "(본문 없음)"))
    setname=a.set or "암송"; dept=a.dept or ""; verlabel=(ver or "주소중심")
    d=newdoc()
    para(d,CHURCH,11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,f"성경 암송표 · {setname}",22,NAVY,True,AL.CENTER,before=4,after=2,font=SANS)
    para(d," · ".join(x for x in [dept,verlabel,f"{len(items)}구절"] if x),10,MAROON,True,AL.CENTER,after=8,font=SANS); hr(d,"1F5C9E",12)
    t=d.add_table(rows=len(items)+1,cols=4)
    _bodyhdr = "본문(직접 기입)" if addr_mode else ("본문(개역개정)" if ver=="개역개정" else "본문(공개역본)")
    for j,h in enumerate(["번호","성구",_bodyhdr,"암송 ✓"]): ctext(t.rows[0].cells[j],h,10,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],"1F5C9E")
    for i,(ref,text) in enumerate(items,1):
        ctext(t.rows[i].cells[0],str(i),9,GRAY,font=SANS,align=AL.CENTER)
        ctext(t.rows[i].cells[1],ref,10,NAVY,True,SANS)
        ctext(t.rows[i].cells[2],text,10,GRAY,font=SERIF)
        ctext(t.rows[i].cells[3],"☐",11,LG,font=SANS,align=AL.CENTER)
    tborders(t)
    d.add_page_break()
    para(d,f"지갑용 암송카드 · {setname} (잘라서 사용)",14,MAROON,True,AL.CENTER,after=6,font=SANS)
    rows=(len(items)+1)//2; ct=d.add_table(rows=max(1,rows),cols=2)
    for i,(ref,text) in enumerate(items):
        cell=ct.rows[i//2].cells[i%2]; cell.text=""
        p=cell.paragraphs[0]; r=p.add_run(f"[{setname}] {ref}"); r.bold=True; r.font.size=Pt(10); r.font.color.rgb=NAVY; setf(r,SANS)
        p2=cell.add_paragraph(); r2=p2.add_run(text); r2.font.size=Pt(9); r2.font.color.rgb=GRAY; setf(r2,SERIF)
        p3=cell.add_paragraph(); r3=p3.add_run(f"— {CHURCH}"); r3.font.size=Pt(8); r3.font.color.rgb=LG; setf(r3,SANS)
    tborders(ct,"888888",6)
    _src = "주소중심(본문은 본인 성경으로)" if addr_mode else ("개역개정(한국어)" if ver=="개역개정" else f"{ver}(공개역본·저작권 자유)")
    para(d,f"· {CHURCH} · {today()} · 본문 출처: {_src}",9,LG,align=AL.RIGHT,before=8,font=SANS)
    out=os.path.join(BIBLECAT(),f"[성경암송표] {sanit(setname)}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 성경 암송표·지갑카드 생성({len(items)}구절·{setname}·{verlabel}): {out}")

def kakao(a):
    """카톡에 붙여넣기 목록(무료·붙여넣기형) — 대상별 완성 메시지를 txt로. 목사님이 카톡방에 붙여넣기.
    개인 카톡 자동발송은 카카오 정책상 불가 → 생성→붙여넣기가 비용0 실용안."""
    db=load(); mode=a.type or "custom"; blocks=[]; title=""
    if mode=="birthday":
        ahead=int(a.ahead or 1); tgt=(datetime.date.today()+datetime.timedelta(days=ahead)); md=tgt.isoformat()[5:]
        title=f"생일 축하 ({tgt.isoformat()})"
        for m in db["교인"]:
            if (m.get("생년월일") or "")[5:]==md and md:
                msg=f"{m['이름']} {m['직분']}님, 생일을 진심으로 축하드립니다! 🎂\n{PASTOR}와 {CHURCH} 온 성도가 늘 축복하며 기도합니다. 샬롬 🙏"
                blocks.append((m['이름'],m['연락처'],m['소속셀'],msg))
    elif mode=="devotion":
        title="오늘의 묵상 발송"; msg=a.text or f"[{CHURCH} 오늘의 묵상] (묵상 내용)"
        for m in db["교인"]: blocks.append((m['이름'],m['연락처'],m['소속셀'],msg))
    else:
        title=a.title or "공지"; msg=a.text or ""
        who=[m for m in db["교인"] if (not a.to or a.to in ("all","전체") or m['소속셀']==a.to)]
        for m in who: blocks.append((m['이름'],m['연락처'],m['소속셀'],msg))
    lines=[f"════ 카톡에 붙여넣기 목록 · {title} ════",
           f"※ 아래 각 블록을 복사해 해당 성도 카톡방에 붙여넣으세요(무료). 총 {len(blocks)}명.",""]
    from collections import defaultdict
    bycell=defaultdict(list)
    for b in blocks: bycell[b[2] or '미배정'].append(b)
    for cell,bs in bycell.items():
        lines.append(f"───── [{cell}] {len(bs)}명 ─────")
        for name,tel,_,msg in bs:
            lines.append(f"▼ {name} ({tel})"); lines.append(msg); lines.append("")
    folder=CAT("07")
    out=os.path.join(folder,f"[카톡발송] {sanit(title)}_{today()}.txt")
    open(out,'w',encoding='utf-8').write("\n".join(lines))
    print(f"✅ 카톡에 붙여넣기 목록: {out}  ({len(blocks)}명)")
    print("  → 파일 열어 블록별로 복사→카톡 붙여넣기(비용0). 완전자동 원하면 유료 알림톡 채널 필요.")

def _days_since(dstr):
    try: return (datetime.date.today()-datetime.date.fromisoformat(dstr)).days
    except Exception: return None
_MUKSANG_THEMES={
 "사랑":{"해설":"하나님의 사랑은 조건이 없습니다. 우리가 사랑스러워서가 아니라, 하나님이 사랑이시기에 먼저 사랑하셨습니다. 그 사랑을 받은 사람은 이제 그 사랑으로 이웃을 품게 됩니다.",
        "질문":["이 말씀에서 하나님의 사랑은 어떻게 드러납니까?","내가 조건 없이 사랑해야 할 사람은 누구입니까?","오늘 그 사랑을 어떤 말과 행동으로 표현할 수 있을까요?"]},
 "믿음":{"해설":"믿음은 보이지 않는 하나님을 신뢰하며 그분의 약속을 붙드는 것입니다. 상황이 아니라 하나님을 바라볼 때, 흔들리던 마음에 평안이 찾아옵니다.",
        "질문":["지금 내 믿음을 흔드는 상황은 무엇입니까?","그 상황에서 하나님은 나에게 무엇을 신뢰하라 하십니까?","오늘 믿음으로 내딛을 한 걸음은 무엇일까요?"]},
 "은혜":{"해설":"은혜는 받을 자격 없는 우리에게 하나님이 값없이 베푸시는 선물입니다. 우리가 무언가를 이루어서가 아니라, 그리스도 안에서 이미 받은 은혜로 오늘을 살아갑니다.",
        "질문":["최근 하나님의 은혜를 느낀 순간은 언제였습니까?","내가 당연하게 여겼던 은혜는 무엇이 있을까요?","그 은혜에 어떻게 반응하며 살고 싶습니까?"]},
 "순종":{"해설":"순종은 내 생각을 내려놓고 하나님의 뜻을 따르는 사랑의 표현입니다. 온전히 이해되지 않아도 먼저 순종할 때, 하나님이 예비하신 길이 열립니다.",
        "질문":["하나님이 이미 알려주셨지만 미루고 있는 순종이 있습니까?","순종을 망설이게 하는 두려움은 무엇입니까?","오늘 순종으로 옮길 구체적인 한 가지는 무엇일까요?"]},
 "기도":{"해설":"기도는 하나님과의 대화이며, 내 힘이 아니라 하나님의 능력을 구하는 자리입니다. 응답보다 먼저 하나님과 가까워지는 것이 기도의 가장 큰 은혜입니다.",
        "질문":["요즘 내 기도의 자리는 어떠합니까?","하나님께 정직하게 아뢰고 싶은 마음은 무엇입니까?","오늘 누구를 위해 기도하기 원하십니까?"]},
 "감사":{"해설":"감사는 이미 받은 은혜를 발견하는 눈입니다. 환경이 완벽해서가 아니라, 그 가운데 함께하시는 하나님을 볼 때 감사가 흘러나옵니다.",
        "질문":["오늘 감사할 작은 일 세 가지는 무엇입니까?","힘든 상황 속에서도 감사할 수 있는 이유는 무엇일까요?","그 감사를 어떻게 표현하고 싶습니까?"]},
 "시련":{"해설":"시련은 우리를 무너뜨리려는 것이 아니라, 더 깊이 하나님을 신뢰하도록 빚어가는 과정입니다. 우리가 약할 때, 하나님의 능력이 온전히 드러납니다.",
        "질문":["지금 내가 지나고 있는 광야는 무엇입니까?","이 시련을 통해 하나님이 빚으시는 것은 무엇일까요?","오늘 이 자리에서 붙들 하나님의 약속은 무엇입니까?"]},
 "소망":{"해설":"그리스도인의 소망은 막연한 낙관이 아니라, 신실하신 하나님께 뿌리내린 확신입니다. 오늘이 흔들려도, 하나님이 이루실 내일을 바라보며 다시 일어섭니다.",
        "질문":["요즘 내 마음의 소망은 어디에 있습니까?","하나님이 주신 약속 중 붙들고 싶은 것은 무엇입니까?","그 소망이 오늘 나의 하루를 어떻게 바꿀 수 있을까요?"]},
 "공동체":{"해설":"교회는 건물이 아니라 그리스도 안에서 한 가족 된 사람들입니다. 서로의 다름을 품고 함께 세워질 때, 우리 가운데 하나님이 거하십니다.",
        "질문":["나는 공동체의 지체들을 한 가족으로 여기고 있습니까?","연합을 위해 내가 더 품고 인내해야 할 부분은 무엇입니까?","오늘 한 지체를 위해 할 수 있는 섬김은 무엇일까요?"]},
 "회개":{"해설":"회개는 정죄가 아니라, 하나님께로 방향을 돌이키는 은혜의 문입니다. 정직하게 내어놓을 때, 하나님은 용서하시고 새롭게 하십니다.",
        "질문":["하나님 앞에 정직하게 내려놓아야 할 것은 무엇입니까?","돌이켜야 할 마음의 방향은 어디입니까?","용서받은 자로서 오늘 어떻게 새롭게 살고 싶습니까?"]},
 "사명":{"해설":"하나님은 우리를 구원하시고 세상으로 보내십니다. 거창한 일이 아니라, 내가 선 자리에서 빛과 소금으로 사는 것이 사명의 시작입니다.",
        "질문":["하나님이 지금 나를 두신 자리는 어디입니까?","그 자리에서 나를 통해 전하고 싶은 복음은 무엇일까요?","오늘 한 사람에게 사랑으로 다가갈 방법은 무엇입니까?"]},
 "평안":{"해설":"세상이 주는 평안과 하나님이 주시는 평안은 다릅니다. 상황이 잔잔해서가 아니라, 폭풍 속에서도 함께하시는 주님 때문에 우리는 안식합니다.",
        "질문":["지금 내 마음을 어지럽히는 염려는 무엇입니까?","그 염려를 하나님께 맡긴다는 것은 어떤 의미일까요?","오늘 하나님이 주시는 평안을 어디서 누리고 싶습니까?"]},
}
_MUKSANG_GENERAL=["이 말씀에서 하나님은 어떤 분으로 나타나십니까?",
 "이 말씀이 오늘 나에게 주시는 도전이나 위로는 무엇입니까?",
 "이 말씀에 순종한다면, 오늘 내 삶에서 무엇이 달라질까요?"]
_MUKSANG_APPLY={
 "사랑":"오늘 한 사람을 정하여, 그에게 하나님의 사랑을 따뜻한 말이나 작은 섬김으로 구체적으로 표현해 보십시오.",
 "믿음":"지금 마음을 짓누르는 염려 하나를 하나님께 맡기고, 믿음으로 오늘 한 걸음을 내디뎌 보십시오.",
 "은혜":"당연하게 여겼던 은혜 하나를 적어 감사하고, 그 은혜를 오늘 누군가에게 흘려보내 보십시오.",
 "순종":"미뤄 왔던 순종 한 가지를 오늘 곧바로 실천에 옮겨 보십시오.",
 "기도":"오늘 정한 시간에 잠시 멈추어, 한 사람의 이름을 부르며 짧게라도 기도해 보십시오.",
 "감사":"오늘 감사할 세 가지를 적어 보고, 그중 하나를 그 대상에게 직접 표현해 보십시오.",
 "시련":"지금의 어려움 속에서 붙들 하나님의 약속 한 구절을 적어, 하루 종일 되새겨 보십시오.",
 "소망":"낙심되는 상황 하나를 하나님이 이루실 내일의 눈으로 다시 바라보며, 소망의 한 문장을 적어 보십시오.",
 "공동체":"오늘 한 지체에게 안부를 전하거나, 작은 섬김으로 함께함을 표현해 보십시오.",
 "회개":"하나님 앞에 정직히 내려놓을 것 하나를 고백하고, 돌이켜 새 걸음을 시작해 보십시오.",
 "사명":"지금 선 자리에서 한 사람에게 빛과 소금이 될 작은 실천 하나를 정해 오늘 행해 보십시오.",
 "평안":"마음을 어지럽히는 염려 하나를 하나님께 맡기는 기도를 드리고, 그분이 주시는 평안 안에 머물러 보십시오.",
}
def _muksang_pick(theme, blob, datestr=None):
    """주제(입력) 또는 본문·제목 키워드로 해설·질문을 고른다. 없으면 일반 질문."""
    key=None; t=(theme or "").strip()
    ALIAS={"고난":"시련","고통":"시련","환난":"시련","교회":"공동체","연합":"공동체","가족":"공동체",
           "전도":"사명","선교":"사명","제자":"사명","십자가":"은혜","용서":"회개","돌이킴":"회개","걱정":"평안","염려":"평안"}
    hay=(t+" "+(blob or "")).strip()
    for k in _MUKSANG_THEMES:
        if k in hay: key=k; break
    if not key:
        for al,k in ALIAS.items():
            if al in hay: key=k; break
    if key:
        th=_MUKSANG_THEMES[key]; return key, th["해설"], list(th["질문"])
    return None, "", list(_MUKSANG_GENERAL)
def _muksang_image(ref, phrase, church, outpath, mon=None):
    """말씀 카드 이미지(png) — 따뜻한 계절 그라데이션 배경 + 성구 주소 + 짧은 문구. 저작권 안전(본문 전문 없음)."""
    import subprocess, shutil, datetime
    ff=None
    for c in (os.path.join(BASE,"ffmpeg.exe"), os.path.join(BASE,"_bin","ffmpeg.exe")):
        if os.path.exists(c): ff=c; break
    if not ff: ff=shutil.which("ffmpeg")
    if not ff: return None
    W=H=1080; FONT=r"C:/Windows/Fonts/malgun.ttf".replace(':',r'\:')
    # 배경 이미지 폴더 준비(목사님이 교회 사진·풍경·무료 이미지를 넣으면 그 위에 성구를 얹음)
    bgdir=os.path.join(BASE,"_내자료","묵상배경"); os.makedirs(bgdir,exist_ok=True)
    _gd=os.path.join(bgdir,"여기에 배경 그림을 넣으세요.txt")
    if not os.path.exists(_gd):
        open(_gd,'w',encoding='utf-8').write(
            "이 폴더에 사진·그림(jpg·png)을 넣으면, 오늘의 묵상 '말씀 카드' 배경으로 쓰입니다(날짜별로 돌아가며).\n"
            "· 교회 사진, 풍경, 좋아하는 이미지 등 무엇이든 좋습니다.\n"
            "· 무료 이미지는 픽사베이(pixabay.com)·언스플래시(unsplash.com)·펙셀스(pexels.com)에서 받으실 수 있습니다(라이선스 확인).\n"
            "· 비워 두시면 따뜻한 계절 그라데이션 배경이 자동으로 쓰입니다.")
    bgs=[os.path.join(bgdir,f) for f in sorted(os.listdir(bgdir))
         if '.' in f and f.lower().rsplit('.',1)[-1] in ('jpg','jpeg','png','webp')] if os.path.isdir(bgdir) else []
    import tempfile as _tf
    tmp=_tf.mkdtemp()
    open(os.path.join(tmp,"ref.txt"),'w',encoding='utf-8').write(str(ref or "오늘의 말씀"))
    open(os.path.join(tmp,"ph.txt"),'w',encoding='utf-8').write(_vwrap(str(phrase or "오늘도 은혜 안에서"),14))
    open(os.path.join(tmp,"ch.txt"),'w',encoding='utf-8').write(str(church or ""))
    AL=f"drawbox=x={W//2-70}:y=352:w=140:h=3:color=%s:t=fill"   # 성구 아래 강조선
    if bgs:   # 배경 사진 위 — 어두운 막 + 비네트(깊이감) + 흰 글자 + 금색 프레임
        pick=bgs[datetime.date.today().toordinal() % len(bgs)]
        inputs=['-i',pick]
        vf=(f"scale={W}:{H}:force_original_aspect_ratio=increase,crop={W}:{H},"
            f"drawbox=x=0:y=0:w={W}:h={H}:color=black@0.42:t=fill,vignette=a=PI/4.6,"
            f"drawbox=x=66:y=66:w={W-132}:h={H-132}:color=white@0.55:t=2,"
            f"drawtext=fontfile='{FONT}':textfile='ref.txt':fontcolor=0xF3E1B0:fontsize=48:x=(w-tw)/2:y=270:shadowcolor=black@0.8:shadowx=2:shadowy=2,"
            +(AL%"0xF3E1B0")+","
            f"drawtext=fontfile='{FONT}':textfile='ph.txt':fontcolor=white:fontsize=70:x=(w-tw)/2:y=(h-th)/2+30:line_spacing=24:shadowcolor=black@0.85:shadowx=3:shadowy=3,"
            f"drawtext=fontfile='{FONT}':textfile='ch.txt':fontcolor=0xF3D9A0:fontsize=36:x=(w-tw)/2:y=h-175:shadowcolor=black@0.8:shadowx=2:shadowy=2")
    else:     # 계절 그라데이션 + 은은한 패널 + 금색 프레임 + 어두운 글자
        m=mon or datetime.date.today().month
        pal={3:("FCE3EC","F6C8DE"),4:("FCE3EC","F6C8DE"),5:("E9F6DE","C7E9B0"),6:("DCEEFB","BBD8F5"),7:("DCEEFB","BBD8F5"),8:("DCEEFB","BBD8F5"),
             9:("FBEAD0","F4D4A6"),10:("FBEAD0","F4D4A6"),11:("FBEAD0","F4D4A6"),12:("EAF0FA","D2E0F2"),1:("EAF0FA","D2E0F2"),2:("EAF0FA","D2E0F2")}.get(m,("FBEAD0","F4D4A6"))
        inputs=['-f','lavfi','-i',f'gradients=s={W}x{H}:c0=0x{pal[0]}:c1=0x{pal[1]}:d=1']
        vf=(f"drawbox=x=90:y=90:w={W-180}:h={H-180}:color=white@0.48:t=fill,"
            f"drawbox=x=66:y=66:w={W-132}:h={H-132}:color=0xB98A3C@0.55:t=2,"
            f"drawtext=fontfile='{FONT}':textfile='ref.txt':fontcolor=0x9A6A2E:fontsize=48:x=(w-tw)/2:y=270:box=0,"
            +(AL%"0xB98A3C")+","
            f"drawtext=fontfile='{FONT}':textfile='ph.txt':fontcolor=0x243244:fontsize=70:x=(w-tw)/2:y=(h-th)/2+30:line_spacing=24,"
            f"drawtext=fontfile='{FONT}':textfile='ch.txt':fontcolor=0x9A6A2E:fontsize=36:x=(w-tw)/2:y=h-175")
    r=subprocess.run([ff,'-y']+inputs+['-vf',vf,'-frames:v','1',outpath],cwd=tmp,capture_output=True,text=True)
    return outpath if (r.returncode==0 and os.path.exists(outpath)) else None
def devotion(a):
    """오늘의 묵상 말씀 — 성도 발송용 초안(제목·본문·묵상·기도·문자). 내용을 채워 발송."""
    d=newdoc()
    # ★저작권 안전: 성경 본문은 목사님이 직접 입력(프로그램에 성경 미내장). 주제로 해설·질문을 자동 제안.
    _blob=" ".join([getattr(a,'title','') or "", a.text or "", getattr(a,'body','') or "", getattr(a,'verse','') or ""])
    _tkey,_seed,_questions=_muksang_pick(getattr(a,'theme',''), _blob, a.date)
    para(d,f"{CHURCH} · 오늘의 묵상",11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,a.title or "제목을 입력",20,NAVY,True,AL.CENTER,before=4,after=2,font=SANS); hr(d,"1E6E6A",12)
    para(d,f"{a.date or today()}  ·  본문 {a.text or ''}"+(f"  ·  주제 {_tkey}" if _tkey else ""),10,GRAY,AL.CENTER,after=8,font=SANS)
    para(d,"말씀 (본문은 목사님이 직접 적어 넣으세요 — 저작권 안전)",13,RGBColor(0x1E,0x6E,0x6A),True,before=2,after=3,font=SANS); para(d,a.verse or "  (본문 말씀을 여기에 입력)",11.5,NAVY,italic=True,after=6)
    para(d,"묵상 / 해설"+("  (아래는 주제 참고 문장 — 목사님 말로 다듬으세요)" if (not (getattr(a,'body','') or '').strip() and _seed) else ""),13,RGBColor(0x1E,0x6E,0x6A),True,before=4,after=3,font=SANS)
    para(d,(getattr(a,'body','') or '').strip() or _seed or "  (짧은 묵상 3~5문장)",11,after=6)
    _apply_seed=_MUKSANG_APPLY.get(_tkey,"") if _tkey else ""
    para(d,"오늘의 적용"+("  (아래는 주제 참고 문장 — 목사님 말로 다듬으세요)" if (not (a.apply or '').strip() and _apply_seed) else ""),13,RGBColor(0x1E,0x6E,0x6A),True,before=4,after=3,font=SANS); para(d,(a.apply or '').strip() or _apply_seed or "  ",11,after=6)
    para(d,"묵상 질문 (성도와 나눌 질문 — 그대로 보내도 좋습니다)",13,RGBColor(0x1E,0x6E,0x6A),True,before=4,after=3,font=SANS)
    for _i,_q in enumerate(_questions,1): para(d,f"{_i}. {_q}",11,after=2)
    para(d,"기도",13,RGBColor(0x1E,0x6E,0x6A),True,before=6,after=3,font=SANS); para(d,a.pray or "  ",11,after=6)
    # ── 오늘의 예화 — 날짜별 자동 순환(해가 바뀌어도 매일 다른 예화가 나옴) ──
    #    전용판=채굴 예화 / 배포판=목사님이 넣으신 예화. 저작권 안전: 창작 재서술.
    day_ill=_illus_of_day(a.date)
    theme_ill=[it for it in _match_illus((a.theme or "")+(a.title or "")+(a.text or "")) if it is not day_ill]
    if day_ill:
        hr(d,"CCCCCC",4)
        para(d,"◈ 오늘의 예화 (날짜별 자동 순환 · 묵상에 녹이는 은혜의 소재)",11.5,GOLD,True,before=4,after=3,font=SANS)
        para(d,f"· {day_ill.get('예화명','')}  ({day_ill.get('출처','')})",11,NAVY,True,after=1)
        if day_ill.get('성경연결'): para(d,f"   성경연결: {day_ill.get('성경연결','')}",10,GRAY,after=3)
        _full=str(day_ill.get('전문','')).strip()
        if _full: para(d,_full,10.5,after=4)
        if theme_ill:
            para(d,"◈ 주제로도 어울리는 예화",11,GOLD,True,before=3,after=2,font=SANS)
            for it in theme_ill[:2]:
                para(d,f"· {it.get('예화명','')} ({it.get('출처','')}) — {it.get('성경연결','')}",10.5,NAVY,after=1)
    else:
        hr(d,"CCCCCC",4)
        para(d,"◈ 오늘의 예화 — 목사님만의 예화은행을 채워보세요",11.5,GOLD,True,before=4,after=3,font=SANS)
        para(d,"예화를 한 편씩 넣어두시면(명령 illus-add 또는 웹의 '예화 추가'), 넣으신 예화가 날짜별로 하나씩 자동으로 '오늘의 묵상' 소재로 순환합니다. 해가 바뀌어도 날짜마다 다른 예화가 나와, 매일의 묵상 준비가 한결 가벼워집니다.",10,GRAY,after=4)
    hr(d,"CCCCCC",4)
    # ★저작권 안전: 본문(개역개정/NIV) 전문 복제 배포 금지 → 성구 주소 + 목사님 창작 묵상만 발송
    # 성도 발송문 — 목사님이 쓰신 제목·묵상·적용·기도 그대로(창작 글이라 저작권 안전), 본문은 성구 주소만.
    _mbody=(getattr(a,'body','') or '').strip() or _seed
    _parts=[f"[{CHURCH} 오늘의 묵상] {(a.title or '').strip()}".rstrip(),
            (f"📖 본문 {a.text}  (각자 성경으로 펴서 읽어보세요)" if (a.text or '').strip() else ""),
            (f"\n{_mbody}" if _mbody else ""),
            (f"\n[오늘의 적용] {(a.apply or '').strip()}" if (a.apply or '').strip() else ""),
            (("\n[묵상 질문]\n"+"\n".join(f"{i}. {q}" for i,q in enumerate(_questions,1))) if _questions else ""),
            (f"\n[기도] {(a.pray or '').strip()}" if (a.pray or '').strip() else ""),
            "\n오늘도 말씀 안에서 은혜 충만하세요! 🙏"]
    send_msg="\n".join(p for p in _parts if p)
    para(d,"■ 성도 발송문 (저작권 안전형: 본문 전문 대신 성구주소 + 창작 묵상)",10.5,MAROON,True,before=6,after=3,font=SANS)
    para(d,send_msg,10.5,GRAY)
    para(d,"※ 성도들에게 보내는 법: 함께 만들어지는 '카톡에 붙여넣기 목록'(메모장) 파일을 열어 → 성도별 내용을 복사 → 카톡에 붙여넣어 보내시면 됩니다(무료·비용0).",8.5,LG,before=3,font=SANS)
    para(d,"※ 개역개정·NIV 본문은 저작권 보호 대상 → 전문 복제 배포하지 않고 성구 주소만 안내. 묵상·적용·기도문은 목사님이 새로 쓰신 창작 글이라 저작권 걱정 없이 성도들께 보내실 수 있습니다.",8.5,LG,before=4,font=SANS)
    folder=CAT("09")
    _t=sanit(a.title or '묵상'); _b=sanit((a.text or '').strip())
    out=os.path.join(folder,f"[오늘의묵상] {_t}"+(f" ({_b})" if _b else "")+f"_{a.date or today()}.docx"); _savedoc(d,out)
    print(f"✅ 오늘의 묵상 생성: {out}")
    # ── 말씀 카드 이미지(png) — 성구 주소 + 짧은 문구 + 계절 배경(저작권 안전: 본문 전문 없음) ──
    try:
        _ref=(a.text or '').strip() or (a.title or '오늘의 말씀').strip()
        _phrase=(getattr(a,'theme','') or '').strip() or (a.title or '').strip() or "오늘도 은혜 안에서"
        _imgp=os.path.join(folder,f"[묵상카드] {_t}_{a.date or today()}.png")
        _img=_muksang_image(_ref,_phrase,CHURCH,_imgp)
        if _img:
            print(f"✅ 묵상 카드 이미지: {_img}")
            print(f"__OPENFILE__{_img}")
        else:
            print("   (묵상 카드 이미지는 영상 엔진(ffmpeg)이 있는 환경에서 함께 만들어집니다.)")
    except Exception as _e: print(f"   (묵상 카드 이미지 건너뜀: {str(_e)[:60]})")
    print("─"*34)
    print("※ 배경 그림을 바꾸려면: _내자료 › 묵상배경 폴더에 교회·풍경 사진을 넣으세요(없으면 계절 배경 자동).")
    print("※ 성경 본문은 목사님이 직접 넣는 구조라, 저작권 걱정 없이 성도들께 카톡으로 보내실 수 있습니다.")
    print("※ 주제(사랑·믿음·공동체·시련 등)를 넣으면 그에 맞는 해설·질문이 나옵니다(비우면 자동 판별).")
    # ── 교인 카톡에 붙여넣기 목록 자동 생성 — 묵상 내용이 있고 등록 교인이 있으면 카톡 붙여넣기 txt까지 한 번에 ──
    if (a.body or '').strip():
        try:
            db=load(); mem=db.get("교인",[])
            if mem:
                from collections import defaultdict
                bycell=defaultdict(list)
                for m in mem: bycell[m.get('소속셀') or '미배정'].append(m)
                L=[f"════ 오늘의 묵상 · 카톡에 붙여넣기 목록 ({a.date or today()}) ════",
                   f"※ 이 메모장을 열어, 아래 각 사람 블록을 복사해서 카톡에 붙여넣어 보내세요(무료·비용0). 총 {len(mem)}명.",""]
                for cell,ms in bycell.items():
                    L.append(f"───── [{cell}] {len(ms)}명 ─────")
                    for m in ms:
                        L.append(f"▼ {m['이름']} ({m.get('연락처','') or ''})"); L.append(send_msg); L.append("")
                kit=os.path.join(CAT("07"),f"[오늘의묵상_카톡발송] {_t}_{a.date or today()}.txt")
                open(kit,'w',encoding='utf-8').write("\n".join(L)); _openfile(kit)
                print(f"✅ 교인 카톡에 붙여넣기 목록: {kit}  ({len(mem)}명) — 열어서 성도별 카톡 붙여넣기(무료)")
            else:
                print("   (등록 교인이 없어 붙여넣기 목록은 생략 — 교인 등록 후 다시 만드시면 카톡에 붙여넣기 목록이 함께 나옵니다)")
        except Exception as e:
            print(f"   붙여넣기 목록 생략: {str(e)[:60]}")
    else:
        print("   (묵상 내용을 채우시면, 교인에게 보낼 카톡에 붙여넣기 목록이 함께 만들어집니다)")
def care(a):
    """돌봄 필요 성도 감지 — 심방한 지 오래된 순 (벤치마크 #2)"""
    db=load(); days=_intarg(a.days, 90, "일수"); rows=[]
    for m in db["교인"]:
        last=m["심방이력"][-1]["날짜"] if m["심방이력"] else None
        ds=_days_since(last) if last else 9999
        if ds>=days: rows.append((ds,m,last))
    rows.sort(key=lambda t:t[0],reverse=True)
    print(f"■ 돌봄 필요 성도 (심방 {days}일 이상 경과): {len(rows)}명")
    for ds,m,last in rows:
        print(f"  ⚠ {m['이름']} {m['직분']} ({m['소속셀'] or '-'}) — 마지막 심방 {'없음' if last is None else last+f' ({ds}일 전)'}")
    if not rows: print("  (전원 최근 심방 완료)")
def newcomer(a):
    """새가족 정착 파이프라인 (벤치마크 #1)"""
    db=load(); days=_intarg(a.days, 90, "일수"); rows=[]
    for m in db["교인"]:
        reg=_days_since(m.get("등록일",""))
        if reg is not None and reg<=days: rows.append((reg,m))
    rows.sort(key=lambda t:t[0])
    print(f"■ 새가족 정착 관리 (최근 {days}일 등록): {len(rows)}명")
    for reg,m in rows:
        vc=len(m["심방이력"]); wk=reg//7
        step="첫 심방(1주 내)" if vc==0 else ("소그룹 배정·4주 양육" if vc<3 else "정착 확인")
        flag=" ❗미이행" if (reg>=7 and vc==0) else ""
        print(f"  · {m['이름']} — 등록 {reg}일({wk}주) 전 · 심방 {vc}회 · 다음단계: {step}{flag}")
    if not rows: print("  (최근 등록 새가족 없음)")

def weekly_brief(a):
    """주간 목회 브리핑 — 교인·출석·재정·이번주 생일·돌봄대상을 한 장 요약(카톡 발송 가능)."""
    import datetime as _dt
    db=load(); act=[m for m in db["교인"] if m.get("상태","재적")=="재적"]
    L=[f"■ {CHURCH} 주간 목회 브리핑  ({today()})",f"· 재적 교인 {len(act)}명"]
    if db["출석"]: last=db["출석"][-1]; L.append(f"· 최근 출석: {last['날짜']} {last['예배']} {last['계']}명")
    ym=today()[:7]; rec=[r for r in db["재정"] if r['날짜'].startswith(ym)]
    inc=sum(r['금액'] for r in rec if r['구분']=='수입'); exp=sum(r['금액'] for r in rec if r['구분']=='지출')
    L.append(f"· 이번달 재정: 수입 {inc:,} · 지출 {exp:,} · 잔액 {inc-exp:,}원")
    td=_dt.date.fromisoformat(today()); bd=[]
    for m in act:
        b=(m.get('생년월일') or '')[5:]
        if b and any(b==(td+_dt.timedelta(days=o)).isoformat()[5:] for o in range(7)): bd.append(f"{m['이름']}")
    L.append(f"· 이번주 생일: {', '.join(bd) if bd else '없음'}")
    caren=[m['이름'] for m in act if (not m['심방이력']) or (_days_since(m['심방이력'][-1]['날짜']) or 0)>=90]
    L.append(f"· 돌봄 필요(심방 오래): {', '.join(caren[:10]) if caren else '없음'}")
    up=sorted([e for e in db.get("일정",[]) if e.get("날짜","")>=today()],key=lambda x:x["날짜"])[:5]
    if up: L.append("· 다가오는 일정: "+", ".join(f"{e['날짜']} {e.get('유형','')}{('·'+e['장소']) if e.get('장소') else ''}" for e in up))
    nf=[m for m in db["교인"] if m.get("유형")=="새가족" and m.get("정착여부")=="진행"]
    if nf: L.append(f"· 정착 중 새가족: {len(nf)}명 (newfamily-board 로 확인)")
    txt="\n".join(L); print(txt)
    out=os.path.join(CAT("01"),f"[주간브리핑] {today()}.txt"); open(out,'w',encoding='utf-8').write(txt)
    print(f"  → 저장: {out} (카톡 발송: kakao --type custom --to all --text ...)")
def prayer_add(a):
    """기도제목 접수 — 주제 태깅해 누적."""
    if not (a.text or a.name): print("⚠ 기도제목 내용이나 대상을 입력해 주세요."); return
    db=load(); db.setdefault("기도제목",[])
    db["기도제목"].append({"날짜":a.date or today(),"교인":a.name or "","내용":a.text or "","주제":a.theme or "일반","상태":"기도중"})
    save(db); print(f"✅ 기도제목 접수: {a.name or ''} — {a.text or ''} [{a.theme or '일반'}]")
def prayer_digest(a):
    """기도제목 정리 — 주제별로 묶어 주간 기도목록."""
    db=load(); ps=[p for p in db.get("기도제목",[]) if p["상태"]!="응답"]
    from collections import defaultdict
    by=defaultdict(list)
    for p in ps: by[p["주제"]].append(p)
    print(f"■ 기도제목 정리 ({len(ps)}건 기도중)")
    for t,items in by.items():
        print(f" [{t}]")
        for p in items: print(f"  · {p['교인']}: {p['내용']} ({p['날짜']})")
    if not ps: print("  (기도중 항목 없음)")
def giving_insight(a):
    """헌금 케어 신호 — 정기헌금 끊긴 교인·새가족 첫 헌금 감지(심방 우선순위)."""
    db=load(); days=_intarg(a.days, 60, "일수"); last={}
    for r in db["재정"]:
        if r["구분"]=="수입" and r.get("교인"): last[r["교인"]]=max(last.get(r["교인"],""),r["날짜"])
    lapsed=[(n,d) for n,d in last.items() if (_days_since(d) or 0)>=days]
    print(f"■ 헌금 케어 신호 (정기헌금 {days}일+ 끊김): {len(lapsed)}명")
    for n,d in sorted(lapsed,key=lambda x:x[1]): print(f"  ⚠ {n} — 마지막 헌금 {d}")
    newg=[m['이름'] for m in db['교인'] if m['이름'] in last and (_days_since(m.get('등록일','')) or 999)<=90]
    if newg: print(f"  🌱 최근 새가족 중 헌금 시작: {', '.join(newg)}")
    if not lapsed and not newg: print("  (신호 없음 — 재정 기록에 '교인'을 넣어야 정확)")
def schedule_add(a):
    """집회·외부 일정 (등록 & 목록) 통합 — 장소·주제 등을 넣으면 등록, 비우면 일정 목록을 본다."""
    if not (a.place or a.host or a.theme):
        return schedule_list(a)   # 내용 비우면 일정 목록 보기(통합)
    db=load(); db.setdefault("일정",[])
    db["일정"].append({"날짜":a.date or today(),"유형":a.type or "집회","장소":a.place or "","담당자":a.host or "",
        "연락처":a.tel or "","주제":a.theme or "","사례":a.fee or "","메모":a.memo or "","상태":"예정"})
    db["일정"].sort(key=lambda e:e["날짜"]); save(db)
    print(f"✅ 일정 등록: {a.date or today()} [{a.type or '집회'}] {a.place or ''} — {a.theme or ''}")
def schedule_list(a):
    import datetime as _dt
    db=load(); evs=sorted(db.get("일정",[]),key=lambda e:e["날짜"])
    if a.upcoming: evs=[e for e in evs if e["날짜"]>=today()]
    print(f"■ 집회·초청 일정 {len(evs)}건"+(" (다가오는)" if a.upcoming else ""))
    for e in evs:
        try:
            dd=(_dt.date.fromisoformat(e["날짜"])-_dt.date.today()).days
            tag=f"D-{dd}" if dd>0 else ("D-DAY" if dd==0 else f"({-dd}일전)")
        except Exception: tag=""
        line=f"  · {e['날짜']} {tag} [{e['유형']}] {e['장소']} — {e.get('주제','')}"
        if e.get("담당자"): line+=f" · {e['담당자']}({e.get('연락처','')})"
        if e.get("사례"): line+=f" · 사례 {e['사례']}"
        print(line)
    if not evs: print("  아직 등록된 집회 일정이 없습니다.\n  → 먼저 '집회·외부설교 일정' 카드로 집회를 등록하시면 여기에 나타납니다.")
def student_add(a):
    """학생 등록(청소년부) — 학교·학년·부서·보호자 연락처까지. 청소년 사역용."""
    if not a.name: print("⚠ 학생 이름을 입력해 주세요."); return
    db=load(); db["_seq"]+=1
    m={"id":db["_seq"],"이름":a.name,"성별":a.sex or "","생년월일":a.birth or "","연락처":a.tel or "",
       "주소":a.addr or "","직분":"학생","세례":"","소속셀":a.dept or "","인도자":a.leader or "","등록일":a.date or today(),
       "가족":[],"상태":"재적","심방이력":[],"기도제목누적":[],"메모":a.memo or "",
       "유형":"학생","학교":a.school or "","학년":a.grade or "","보호자":a.guardian or "","보호자연락처":a.gtel or "","부서":a.dept or "","전공":a.major or ""}
    db["교인"].append(m); save(db)
    print(f"✅ 학생 등록: {a.name} · {a.school or ''} {a.grade or ''} · [{a.dept or ''}]{' · 전공 '+a.major if a.major else ''} · 보호자 {a.guardian or ''}({a.gtel or ''})")
def student_list(a):
    db=load(); ss=[m for m in db["교인"] if m.get("유형")=="학생"]
    if a.dept: ss=[m for m in ss if m.get("부서")==a.dept]
    if a.grade: ss=[m for m in ss if m.get("학년")==a.grade]
    print(f"■ 학생 {len(ss)}명"+(f" ({a.dept or ''}{a.grade or ''})" if (a.dept or a.grade) else ""))
    for m in ss:
        print(f"  · {m['이름']} {m.get('학교','')} {m.get('학년','')} [{m.get('부서','')}] {m['연락처']} · 보호자 {m.get('보호자','')}({m.get('보호자연락처','')})")
    if not ss: print("  (등록 학생 없음 — '학생 등록(청소년)'에서 등록해 주세요)")
def exam_cheer(a):
    """시험기간 응원 — 학생·보호자에게 보낼 격려 카톡 꾸러미(무료 붙여넣기형). 청소년 사역용."""
    db=load(); ss=[m for m in db["교인"] if m.get("유형")=="학생"]
    if a.dept: ss=[m for m in ss if m.get("부서")==a.dept]
    ex=a.exam or "시험"; L=[f"════ 시험기간 응원 · {ex} ════","※ 블록별로 복사해 카톡에 붙여넣으세요(무료).",""]
    for m in ss:
        L.append(f"▼ {m['이름']} ({m['연락처']})")
        L.append(f"{m['이름']}아! {ex} 준비하느라 애쓰지? 결과보다 네가 하나님의 자녀인 게 더 소중해. '내가 네게 명한 것이 아니냐 마음을 강하게 하고 담대히 하라'(수1:9). {CHURCH}가 응원하고 기도한다! 🙏"); L.append("")
        if m.get("보호자연락처"):
            L.append(f"▼ (보호자) {m.get('보호자','')} ({m['보호자연락처']})")
            L.append(f"{m.get('보호자','') or '학부모'}님, {m['이름']} 학생이 {ex}을 잘 치르도록 {CHURCH}가 함께 기도합니다. 늘 감사드립니다 🙏"); L.append("")
    out=os.path.join(CAT("03"),f"[시험응원] {sanit(ex)}_{today()}.txt"); open(out,'w',encoding='utf-8').write("\n".join(L))
    print(f"✅ 시험 응원 카톡에 붙여넣기 목록: {out} (학생 {len(ss)}명 + 보호자)")
def production_add(a):
    """뮤지컬·공연 프로덕션 등록 — 작품·공연일·장소·연출·음악감독."""
    db=load(); db.setdefault("공연",[])
    db["공연"].append({"작품":a.title or "","공연일":a.date or "","장소":a.place or "","연출":a.director or "",
        "음악감독":a.music or "","상태":a.status or "준비","메모":a.memo or ""})
    db["공연"].sort(key=lambda e:e.get("공연일") or "9999"); save(db)
    print(f"✅ 공연 등록: 〈{a.title or ''}〉 · {a.date or ''} · {a.place or ''}")
def production_list(a):
    import datetime as _dt
    db=load(); ps=sorted(db.get("공연",[]),key=lambda e:e.get("공연일") or "9999")
    print(f"■ 공연·프로덕션 {len(ps)}건")
    for p in ps:
        dd=""
        try:
            n=(_dt.date.fromisoformat(p["공연일"])-_dt.date.today()).days; dd=f"D-{n}" if n>0 else ("D-DAY" if n==0 else f"({-n}일전)")
        except Exception: pass
        print(f"  · {p.get('공연일','')} {dd} 〈{p['작품']}〉 {p.get('장소','')} · 연출 {p.get('연출','')} · 음악 {p.get('음악감독','')} [{p.get('상태','')}]")
    if not ps: print("  (등록 공연 없음)")
def casting(a):
    """캐스팅표 — 작품의 배역-배우 배정 + 파트·연락처(docx). 뮤지컬·전문찬양팀용."""
    d=newdoc()
    para(d,f"〈{a.title or '작품'}〉 캐스팅표",20,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,f"{CHURCH}  ·  공연 {a.date or ''}",10,GRAY,AL.CENTER,after=6,font=SANS); hr(d,"6A3D7A",12)
    roles=[r.strip() for r in (a.roles or "").split(";") if r.strip()]
    t=d.add_table(rows=len(roles)+1,cols=3)
    for j,h in enumerate(["배역","배우","파트/연락처"]):
        ctext(t.rows[0].cells[j],h,11,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],"6A3D7A")
    for i,r in enumerate(roles,1):
        pp=r.split(":"); role=pp[0].strip(); actor=pp[1].strip() if len(pp)>1 else ""; extra=pp[2].strip() if len(pp)>2 else ""
        ctext(t.rows[i].cells[0],role,11,NAVY,True,SANS,AL.CENTER)
        ctext(t.rows[i].cells[1],actor,11,GRAY,font=SERIF); ctext(t.rows[i].cells[2],extra,10,GRAY,font=SERIF)
        for j,w in enumerate([4.0,5.0,7.0]): t.rows[i].cells[j].width=Cm(w)
    tborders(t,"BBBBBB",4)
    out=os.path.join(CAT("07"),f"[캐스팅] {sanit(a.title or '작품')}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 캐스팅표 생성: {out} (배역 {len(roles)}) · 입력형식: '배역:배우:파트' 를 ; 로 구분")
def mission_add(a):
    """단기선교 프로젝트 등록 — 팀명·국가·기간·목적."""
    if not a.title: print("⚠ 팀명을 입력해 주세요."); return
    db=load(); db.setdefault("선교",[])
    db["선교"].append({"팀명":a.title or "","국가":a.country or "","기간":a.period or "","목적":a.purpose or "","팀원":[],"상태":"준비"})
    save(db); print(f"✅ 선교 프로젝트: {a.title or ''} · {a.country or ''} · {a.period or ''}")
def mission_member(a):
    """선교 팀원 등록 — 여권·역할·비상연락·건강특이."""
    db=load(); ms=db.get("선교",[])
    if not ms: print("✗ 먼저 '선교 준비 타임라인'에서 프로젝트를 등록해 주세요"); return
    proj=next((p for p in ms if a.title and a.title in p["팀명"]), ms[-1])
    proj["팀원"].append({"이름":a.name or "","역할":a.role or "","여권만료":a.expire or "","비상연락":a.emergency or "","건강":a.health or ""})
    save(db); print(f"✅ [{proj['팀명']}] 팀원: {a.name} ({a.role or ''}) → 총 {len(proj['팀원'])}명")
def mission_roster(a):
    """선교 팀 명단·역할·비상연락망 + 여권만료 점검(docx)."""
    db=load(); ms=db.get("선교",[])
    proj=next((p for p in ms if a.title and a.title in p["팀명"]), ms[-1]) if ms else None
    if not proj: print("✗ 선교 프로젝트 없음 — 먼저 '선교 준비 타임라인'에서 등록해 주세요"); return
    d=newdoc()
    para(d,f"{proj['팀명']} · 선교팀 명단",20,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,f"{proj.get('국가','')} · {proj.get('기간','')} · {CHURCH}",10,GRAY,AL.CENTER,after=6,font=SANS); hr(d,"1F5C9E",12)
    t=d.add_table(rows=len(proj["팀원"])+1,cols=5)
    for j,h in enumerate(["이름","역할","여권만료","비상연락","건강특이"]):
        ctext(t.rows[0].cells[j],h,10,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],"1F5C9E")
    for i,mm in enumerate(proj["팀원"],1):
        for j,k in enumerate(["이름","역할","여권만료","비상연락","건강"]):
            ctext(t.rows[i].cells[j],mm.get(k,""),9.5,GRAY,font=SERIF)
    tborders(t,"BBBBBB",4); para(d,f"총 {len(proj['팀원'])}명",10,GRAY,before=6,font=SANS)
    out=os.path.join(CAT("05"),f"[선교팀명단] {sanit(proj['팀명'])}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 선교 팀 명단: {out} ({len(proj['팀원'])}명)")
def mission_checklist(a):
    """단기선교 준비 체크리스트 — 서류·건강·이동·사역·준비물(docx)."""
    d=newdoc()
    para(d,f"{a.title or '단기선교'} 준비 체크리스트",19,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,f"{CHURCH} · {today()}",10,GRAY,AL.CENTER,after=6,font=SANS); hr(d,"1F5C9E",12)
    groups=[("서류·증명",["여권(잔여 6개월+ 유효)","비자(단체/개별)","여행자보험(전원)","미성년자 부모동의서·위임장","여권 사본(인솔자 전원 보관)","비상연락망(팀원·보호자·현지)"]),
      ("건강·안전",["현지 필수 예방접종","말라리아 예방약(해당시)","개인 상비약(복용약)","팀 구급함","팀원 알레르기·지병 파악","출국 전 안전교육"]),
      ("청소년 인솔(학생 동반)",["학생 여권 인솔자 일괄 보관","조·인솔교사 배정표","학생 보호자 비상연락 카드","현지 코디네이터 연락처","분실·미아·응급 대응 매뉴얼","귀국 후 학생 인계 확인"]),
      ("이동·숙박·통신",["항공권 발권(전원)","현지 교통","숙소 예약","환전·현지화폐","로밍/현지 유심","단체 카톡방"]),
      ("사역 준비",["찬양·율동 콘티","드라마·공연·뮤지컬 준비","전도지·선물·간증문","현지 인사말·통역","팀 티셔츠·배너"]),
      ("개인 준비물",["여권·사본","세면도구","멀티 어댑터","현지 기후 복장","성경·필기구","개인 간식"]),
      ("귀국 후",["선교 보고서","사진·영상 정리","후원자 감사장","주일 간증·보고"])]
    for g,items in groups:
        para(d,g,13,RGBColor(0x1F,0x5C,0x9E),True,before=8,after=3,font=SANS)
        for it in items: para(d,f"  ☐  {it}",11,after=2)
    out=os.path.join(CAT("05"),f"[선교체크리스트] {sanit(a.title or '단기선교')}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 선교 준비 체크리스트: {out}")
def mission_plan(a):
    """단기선교 준비 타임라인 — 출국일 기준 D-day 단계별 준비사항(날짜 자동계산). 미리 보고 준비."""
    import datetime as _dt
    d=newdoc()
    para(d,f"{a.title or '단기선교'} 준비 타임라인",19,NAVY,True,AL.CENTER,after=2,font=SANS)
    depd=None
    try: depd=_dt.date.fromisoformat(a.dday) if a.dday else None
    except Exception: depd=None
    para(d,(f"출국 {a.dday}  ·  " if a.dday else "")+f"{CHURCH} · {today()}",10,GRAY,AL.CENTER,after=6,font=SANS); hr(d,"1F5C9E",12)
    stages=[(-75,"팀 구성·국가/일정 확정",["팀원 모집·확정","국가·목적지·기간 확정","여권 유효기간 확인(6개월+)","항공권 예약","예산·후원 계획 수립"]),
      (-45,"서류·건강 착수",["비자 신청","예방접종 시작(황열·A형간염 등)","여행자보험 가입","숙소 예약","미성년자 부모동의서 배부"]),
      (-30,"사역 준비 시작",["찬양·율동·드라마 콘티","전도지·선물·간증문 준비","팀 훈련 시작","현지 교회·코디네이터 연락"]),
      (-21,"팀 훈련·리허설",["공연·사역 리허설","역할·조 배정","현지 언어 인사·문화 교육","팀 티셔츠·배너 제작"]),
      (-14,"물품·개인 준비",["환전","유심/로밍","팀 구급함·상비약","짐 목록 배부","단체 카톡방 개설"]),
      (-7,"최종 점검·안전",["부모동의서·비상연락망 취합","출국 전 안전교육","여권 사본 일괄 보관","현지 일정 최종 확인","팀원 건강·알레르기 파악"]),
      (-1,"출국 전날",["개인 짐 싸기","여권·항공권·서류 최종 확인","집결 시간·장소 공지","비상약·귀중품 점검"]),
      (0,"출국 당일",["공항 집결·인원 점검","여권·수하물 확인","파송·출국 기도"])]
    for off,title,items in stages:
        ds=""
        if depd: ds="  ("+(depd+_dt.timedelta(days=off)).isoformat()+")"
        head=(f"D{off}" if off<0 else "D-DAY")+ds
        para(d,f"{head} · {title}",13,RGBColor(0x1F,0x5C,0x9E),True,before=8,after=3,font=SANS)
        for it in items: para(d,f"  ☐  {it}",11,after=2)
    para(d,"현지 매일 점검 (일자별)",13,MAROON,True,before=10,after=3,font=SANS)
    for it in ["인원 점검(아침·저녁)","건강·안전 확인","오늘 사역·준비물","현지 팀 소통","감사·나눔·기도","경비 기록"]:
        para(d,f"  ☐  {it}",11,after=2)
    para(d,"귀국 후",13,RGBColor(0x1F,0x5C,0x9E),True,before=10,after=3,font=SANS)
    for it in ["학생 보호자 인계 확인","선교 보고서","사진·영상 정리","후원자 감사장","주일 간증·보고"]:
        para(d,f"  ☐  {it}",11,after=2)
    out=os.path.join(CAT("05"),f"[선교 준비타임라인] {sanit(a.title or '단기선교')}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 선교 준비 타임라인: {out}"+(f" (출국 {a.dday} 기준 날짜 계산)" if depd else " (--dday 출국일 넣으면 날짜 자동계산)"))
    try: mission_checklist(a)   # 체크리스트도 함께 생성(타임라인 & 체크리스트 통합)
    except Exception as _e: print(f"   (체크리스트 생성 건너뜀: {str(_e)[:50]})")
def mission_place(a):
    """현지 정보 등록 — 선교사·숙소·식당·사역지 이름·연락처·주소."""
    db=load(); ms=db.get("선교",[])
    if not ms: print("✗ 먼저 '선교 준비 타임라인'에서 프로젝트를 등록해 주세요"); return
    proj=next((p for p in ms if a.title and a.title in p["팀명"]), ms[-1])
    proj.setdefault("현지",[]).append({"종류":a.kind or "장소","이름":a.name or "","연락처":a.contact or "","주소":a.address or "","메모":a.memo or ""})
    save(db); print(f"✅ [{proj['팀명']}] 현지 {a.kind or '장소'}: {a.name or ''} ({a.contact or ''})")
def mission_field(a):
    """현지 정보 시트 — 선교사·숙소·식당·사역지 연락처 + 일자별 장소 플랜(docx)."""
    from collections import defaultdict
    db=load(); ms=db.get("선교",[])
    proj=next((p for p in ms if a.title and a.title in p["팀명"]), ms[-1]) if ms else None
    if not proj: print("✗ 선교 프로젝트 없음 — 먼저 '선교 준비 타임라인'에서 등록해 주세요"); return
    d=newdoc()
    para(d,f"{proj['팀명']} · 현지 정보/연락",20,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,f"{proj.get('국가','')} · {proj.get('기간','')} · {CHURCH}",10,GRAY,AL.CENTER,after=6,font=SANS); hr(d,"1F5C9E",12)
    by=defaultdict(list)
    for it in proj.get("현지",[]): by[it["종류"]].append(it)
    labels=[("선교사","현지 선교사 연락처"),("숙소","숙소"),("식당","식당"),("사역지","사역지·선교 장소"),("장소","기타 장소")]
    for k,lab in labels:
        if by.get(k):
            para(d,lab,13,RGBColor(0x1F,0x5C,0x9E),True,before=8,after=3,font=SANS)
            t=d.add_table(rows=len(by[k])+1,cols=4)
            for j,h in enumerate(["이름","연락처","주소","메모"]): ctext(t.rows[0].cells[j],h,10,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],"1F5C9E")
            for i,it in enumerate(by[k],1):
                for j,kk in enumerate(["이름","연락처","주소","메모"]): ctext(t.rows[i].cells[j],it.get(kk,""),9.5,GRAY,font=SERIF)
            tborders(t,"BBBBBB",4)
    if not proj.get("현지"): para(d,"(현지 정보 미등록 — mission-place로 추가하세요)",11,MAROON,after=4)
    para(d,"일자별 장소 플랜",13,MAROON,True,before=10,after=3,font=SANS)
    t=d.add_table(rows=8,cols=3)
    for j,h in enumerate(["날짜","일정·장소(오전/오후/저녁)","숙소/식사"]): ctext(t.rows[0].cells[j],h,10,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],"3C6E3C")
    for r in t.rows:
        for j,w in enumerate([3.0,9.0,4.0]): r.cells[j].width=Cm(w)
    tborders(t,"BBBBBB",4)
    out=os.path.join(CAT("05"),f"[선교 현지정보] {sanit(proj['팀명'])}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 현지 정보 시트: {out}")
def mission_english(a):
    """선교지 영어 회화집 — 인사·예배·전도·기도·일상·응급 필수 표현(한/영)."""
    d=newdoc()
    para(d,"선교지 영어 회화집",20,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,f"{CHURCH} 단기선교 · {today()}",10,GRAY,AL.CENTER,after=6,font=SANS); hr(d,"1F5C9E",12)
    sections=[("인사·소개",[("안녕하세요, 만나서 반갑습니다.","Hello, nice to meet you."),("저는 한국에서 온 ○○입니다.","I'm ○○ from Korea."),("우리는 교회에서 왔습니다.","We came from a church in Korea."),("하나님이 당신을 사랑하십니다.","God loves you.")]),
      ("예배·찬양 인도",[("함께 일어나 찬양합시다.","Let's stand and worship together."),("손을 들고 찬양해요.","Let's raise our hands and praise Him."),("다 같이 박수로 찬양해요.","Let's praise Him with clapping."),("한 번 더 부르겠습니다.","Let's sing it one more time."),("주님을 찬양합니다.","We praise the Lord.")]),
      ("복음 제시",[("예수님이 당신을 위해 죽으셨습니다.","Jesus died for you."),("예수님을 믿으면 구원을 받습니다.","If you believe in Jesus, you will be saved."),("하나님은 당신을 위한 계획이 있으십니다.","God has a plan for you."),("함께 기도해도 될까요?","May I pray with you?"),("예수님을 영접하시겠어요?","Would you like to receive Jesus?")]),
      ("기도·축복",[("함께 기도하겠습니다.","Let's pray together."),("하나님의 은혜가 함께하시길 빕니다.","May God's grace be with you."),("주님이 당신을 지키시고 축복하시길.","May the Lord bless you and keep you."),("예수님의 이름으로 기도합니다. 아멘.","In Jesus' name we pray. Amen.")]),
      ("일상·이동",[("감사합니다.","Thank you."),("죄송합니다.","I'm sorry."),("화장실이 어디예요?","Where is the restroom?"),("얼마예요?","How much is it?"),("도와주시겠어요?","Could you help me?"),("물 좀 주시겠어요?","Could I have some water?")]),
      ("응급·안전",[("도와주세요! 응급상황입니다.","Help! It's an emergency."),("병원이 어디예요?","Where is the hospital?"),("이 사람이 아파요.","This person is sick."),("경찰을 불러 주세요.","Please call the police."),("일행을 잃어버렸어요.","I lost my group.")])]
    for title,rows in sections:
        para(d,title,13,RGBColor(0x1F,0x5C,0x9E),True,before=8,after=3,font=SANS)
        for ko,en in rows:
            para(d,f"· {ko}",11,NAVY,after=1); para(d,f"    {en}",11,GRAY,after=3,font=SANS)
    out=os.path.join(CAT("05"),f"[선교 영어회화]_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 선교지 영어 회화집: {out}")
def mission_sermon(a):
    """선교지 영어 설교 준비 — 한/영 대역 구조 + 영어 성경본문(WEB) 자동삽입(통역용)."""
    import json as _j
    eng=""
    m=re.match(r'\s*([1-9]?[가-힣]+|[A-Za-z ]+?)\s*(\d+):(\d+)(?:-(\d+))?',a.text or "")
    webp=os.path.join(BIBLEDIR,"web.json")
    if m and os.path.exists(webp):
        bib=_j.load(open(webp,encoding='utf-8')); book=BOOKALIAS.get(m.group(1).strip(),m.group(1).strip())
        ch=m.group(2); v1=int(m.group(3)); v2=int(m.group(4) or m.group(3))
        eng=" ".join(f"{v} {bib.get(book,{}).get(ch,{}).get(str(v),'')}".strip() for v in range(v1,v2+1) if bib.get(book,{}).get(ch,{}).get(str(v))).strip()
    d=newdoc()
    para(d,"Mission Sermon · 선교지 영어 설교",11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,a.title or "Sermon Title",20,NAVY,True,AL.CENTER,before=4,after=2,font=SANS); hr(d,"8B2E2E",12)
    para(d,f"Text {a.text or ''}  ·  {CHURCH}",10,GRAY,AL.CENTER,after=6,font=SANS)
    if eng:
        para(d,"Scripture (WEB · public domain)",12,MAROON,True,before=4,after=2,font=SANS)
        para(d,eng,11,NAVY,italic=True,after=6)
    para(d,"※ 통역 동반 시 한 문장씩 끊어 말하세요 (interpreter-friendly, short sentences).",9.5,LG,after=6,font=SANS)
    points=[p.strip() for p in (a.points or "").split(";") if p.strip()] or ["Point 1","Point 2","Point 3"]
    para(d,"Introduction / 서론",13,MAROON,True,before=6,after=3,font=SANS); para(d,"  (opening — short, warm)",10.5,LG,italic=True,after=6)
    for i,pt in enumerate(points,1):
        para(d,f"Point {i}. {pt}",13,NAVY,True,before=6,after=3,font=SANS)
        para(d,"  EN: ",11,MAROON,after=1); para(d,"  통역(KR): ",11,GRAY,after=5)
    para(d,"Conclusion / 결론·결단",13,MAROON,True,before=6,after=3,font=SANS)
    para(d,"  (invitation to receive Jesus)",10.5,LG,italic=True,after=4)
    folder=CAT("09")
    out=os.path.join(folder,f"[선교영어설교] {sanit(a.title or 'sermon')}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 선교지 영어 설교 준비: {out}"+(" · 영어 성경본문 자동삽입" if eng else " (본문 지정 시 영어성경 자동삽입)"))
def sermon_log(a):
    """설교 기록 — 언제 어떤 본문·제목으로 설교했는지 이력 저장(반복 방지·시리즈)."""
    db=load(); db.setdefault("설교",[])
    db["설교"].append({"날짜":a.date or today(),"제목":a.title or "","본문":a.text or "","예배":a.service or "주일낮","시리즈":a.series or "","대지":a.points or ""})
    db["설교"].sort(key=lambda s:s["날짜"]); save(db)
    print(f"✅ 설교 기록: {a.date or today()} [{a.service or '주일낮'}] '{a.title or ''}' — {a.text or ''}"+(f" 〔{a.series}〕" if a.series else ""))
def sermon_list(a):
    """설교 이력 — 언제 무슨 설교 했는지(본문·제목·예배유형 검색). '이 본문 설교했나?' 확인."""
    db=load(); ss=db.get("설교",[]); q=a.query or a.text
    if q: ss=[s for s in ss if q in s.get("본문","") or q in s.get("제목","") or q in s.get("시리즈","")]
    if a.service: ss=[s for s in ss if a.service in s.get("예배","")]
    if a.year: ss=[s for s in ss if s.get("날짜","").startswith(a.year)]
    ss=sorted(ss,key=lambda s:s.get("날짜",""),reverse=True)
    print(f"■ 설교 이력 {len(ss)}편"+(f" · 검색 '{q}'" if q else "")+(f" · {a.service}" if a.service else ""))
    for s in ss[:80]:
        print(f"  · {s['날짜']} [{s.get('예배','')}] {s.get('제목','')} — {s.get('본문','')}"+(f" 〔{s['시리즈']}〕" if s.get('시리즈') else ""))
    if not ss: print("  (기록 없음 — sermon-log 로 기록하거나 sermon 작성 시 자동 기록)")
def _sermon_folders():
    """지난 설교 폴더·내 PC 열기 버튼(주보 카드처럼) — C·D·USB의 예전 설교로 직접 이동."""
    _sf=CAT("09")
    if not os.path.isdir(_sf): _sf=ROOT
    print(f"▶열기|{_sf}|설교 폴더 열기 (탐색기 — 여기서 C·D·USB의 예전 설교로 이동해 여세요)")
    print(f"▶열기|shell:MyComputerFolder|내 PC 열기 (C·D·USB 모든 드라이브에서 예전 설교 찾기)")
    _bd=_C.get("백업폴더")
    if _bd and os.path.isdir(_bd):
        print(f"▶열기|{_bd}|백업(USB·D) 폴더 열기")
def sermon_reuse(a):
    """지난 설교 재생성 — 과거 설교(본문·제목 검색)를 바탕으로 새 설교초안 작성(재묵상·재작성)."""
    import types
    db=load(); ss=db.get("설교",[]); q=a.query or a.text
    hits=[s for s in ss if q and (q in s.get("본문","") or q in s.get("제목","") or q in s.get("시리즈",""))]
    if not hits:
        print(f"✗ '{q}' 관련 기록된 지난 설교가 없습니다 — 아래에서 예전 설교(C·D·USB)를 열어 참고하거나 '설교 이력·검색'으로 확인하세요")
        _sermon_folders(); return
    src=sorted(hits,key=lambda s:s.get("날짜",""))[-1]
    print(f"↻ 지난 설교 재생성: {src.get('날짜','')} '{src.get('제목','')}' ({src.get('본문','')}) 기반으로 새 초안 작성")
    na=types.SimpleNamespace(title=a.title or src.get("제목",""), text=src.get("본문",""), theme=a.theme or "",
        points=src.get("대지","") or "", series=src.get("시리즈",""), service=a.service or src.get("예배","") or "", date=a.date or "")
    sermon(na)
def newfamily_add(a):
    """새가족 등록 — 정착 관리 시작(인도자·첫방문·소속셀). 부흥 교회 뒷문 방지."""
    db=load(); db["_seq"]+=1
    m={"id":db["_seq"],"이름":a.name,"성별":a.sex or "","생년월일":a.birth or "","연락처":a.tel or "",
       "주소":a.addr or "","직분":"새가족","세례":"","소속셀":a.cell or "","인도자":a.leader or "","등록일":a.date or today(),
       "가족":[],"상태":"새가족","심방이력":[],"기도제목누적":[],"메모":a.memo or "",
       "유형":"새가족","첫방문":a.date or today(),"정착단계":0,"양육기록":[],"정착여부":"진행"}
    db["교인"].append(m); save(db)
    print(f"✅ 새가족 등록: {a.name} · 인도자 {a.leader or '-'} · 셀 {a.cell or '-'} — 정착관리 시작")
    print("  다음: 1주내 첫 양육 접촉 → newfamily-step 으로 기록")
def newfamily_step(a):
    """새가족 양육 기록 — 주차별 심방/전화/문자 접촉 + 단계 진행."""
    db=load(); hit=[m for m in db["교인"] if m["이름"]==a.name and m.get("유형")=="새가족"] or find(db,a.name)
    if not hit: print(f"✗ 새가족 '{a.name}' 없음 — 먼저 '새가족 등록'에서 등록해 주세요"); return
    m=hit[0]; m.setdefault("양육기록",[]).append({"날짜":a.date or today(),"방법":a.how or "심방","담당":a.by or PASTOR,"내용":a.note or ""})
    m["정착단계"]=m.get("정착단계",0)+1; save(db)
    print(f"✅ [{m['이름']}] {m['정착단계']}차 양육({a.how or '심방'}) 기록 · {a.note or ''}")
def newfamily_board(a):
    """새가족 정착 대시보드 — 진행 새가족·단계·마지막 접촉·미이행 경고(뒷문 방지)."""
    db=load(); nf=[m for m in db["교인"] if m.get("유형")=="새가족" and m.get("정착여부")=="진행"]
    print(f"■ 새가족 정착 현황 {len(nf)}명 (❗=관리 필요)")
    for m in sorted(nf,key=lambda x:x.get("등록일","")):
        yuk=m.get("양육기록",[]); last=yuk[-1]["날짜"] if yuk else m.get("등록일","")
        gap=_days_since(last) or 0; reg=_days_since(m.get("등록일","")) or 0; step=m.get("정착단계",0)
        warn=" ❗연락 필요(2주+ 공백)" if gap>=14 else (" ❗첫 양육 미시작" if step==0 and reg>=7 else "")
        print(f"  · {m['이름']} ({m['소속셀'] or '-'}) 등록 {reg}일차 · {step}차 양육 · 마지막 {last}({gap}일전){warn}")
        if step>=4 and not warn: print(f"      → 4주 양육 완료. 정착 확인 후 newfamily-graduate 로 교인 전환 가능")
    if not nf: print("  (진행 중 새가족 없음)")
def newfamily_graduate(a):
    """새가족 정착 완료 — 정식 교인 전환."""
    db=load(); hit=[m for m in db["교인"] if m["이름"]==a.name and m.get("유형")=="새가족"]
    if not hit: print(f"✗ 새가족 '{a.name}' 없음"); return
    m=hit[0]; m["정착여부"]="정착"; m["상태"]="재적"; m["직분"]=a.role or "성도"; m["유형"]="교인"
    save(db); print(f"✅ [{m['이름']}] 정착 완료 → 정식 교인 전환 (양육 {len(m.get('양육기록',[]))}회)")
CAREMSG={"장례":"삼가 조의를 표합니다","입원":"쾌유를 빕니다","출생":"출산을 축하합니다","결혼":"결혼을 축하합니다","회갑":"회갑을 축하합니다"}
def careevent_add(a):
    """🎉 축하·위로 문자 (상황별·경조사 통합) — 상황/사건과 이름을 넣으면 골라 쓰는 카톡 문구를 만들고,
       경조사(장례·입원·출생·결혼·회갑)는 기록·경조금·심방·결혼기념일까지 자동 연동한다. 상황별(취업·이사 등)은 3문구 kit."""
    name=(getattr(a,'name','') or "").strip()
    if not name: print("⚠ 대상자 이름을 입력해 주세요."); return
    kind=(getattr(a,'kind','') or getattr(a,'occasion','') or "").strip()
    jik=(getattr(a,'jik','') or getattr(a,'role','') or "").strip()
    date=getattr(a,'date','') or today()
    amount=getattr(a,'amount','') or ""; action=getattr(a,'action','') or ""; note=getattr(a,'note','') or ""
    pool=_congrats_pool()
    if not kind:
        allk=list(dict.fromkeys(list(CAREMSG.keys())+list(pool.keys())))
        print("■ 쓸 수 있는 상황(경조사 + 상황별): "+" · ".join(allk))
        print("  예) careevent-add --kind 취업·입사 --name 홍길동 --jik 형제   /   --kind 장례 --name 김성도 --amount 100000"); return
    # ── 경조사 유형이거나 경조금/조치가 입력되면 대장에 기록(상황별 취업·이사 등은 기록 안 함) ──
    CARE=("장례","입원","출생","결혼","회갑","소천","별세","위독","조의","부고")
    is_care=any(c in kind for c in CARE)
    if is_care or amount or action:
        db=load(); db.setdefault("경조사",[])
        db["경조사"].append({"날짜":date,"종류":kind,"대상":name,"내용":note,"경조금":amount,"조치":action,"상태":"진행"})
        db["경조사"].sort(key=lambda e:e["날짜"],reverse=True)
        linked=""
        if "결혼" in kind:
            hit=find(db,name)
            if hit: hit[0]["결혼기념일"]=date; linked=f" · 교인 '{name}' 결혼기념일({date}) 자동 연동"
        save(db)
        print(f"✅ 경조사 등록: {date} [{kind}] {name}"+(f" · 경조금 {amount}" if amount else "")+linked)
    # ── 문자 생성: 상황별 풍부한 3문구(pool) 우선, 없으면 경조사 기본 메시지 ──
    key=next((k for k in pool if kind==k), None) or next((k for k in pool if (kind in k or k.startswith(kind))), None)
    if key:
        msgs=pool[key]; start=sum(ord(c) for c in name)%len(msgs)
        picks=[msgs[(start+i)%len(msgs)] for i in range(min(3,len(msgs)))]
        L=[f"════ {key} 축하·위로 · 카톡에 붙여넣기 ({name} {jik}) ════",
           "※ 아래 【1】【2】【3】 중 마음에 드는 문구를 하나 골라 복사해서 카톡으로 보내세요(무료).",""]
        for i,it in enumerate(picks,1):
            L.append(f"【{i}】 (성구 {it.get('성구','')})"); L.append(it['문구'].format(church=CHURCH,name=name,jik=jik,pastor=PASTOR)); L.append("")
        out=os.path.join(CAT("07"),f"[축하위로_{key}] {sanit(name)}_{today()}.txt")
        open(out,'w',encoding='utf-8').write("\n".join(L)); _openfile(out)
        print(f"✅ {key} 축하·위로 문구 3가지: {out}  (【1】【2】【3】 중 골라 카톡 붙여넣기)")
        for i,it in enumerate(picks,1): print(f"  【{i}】 {it['문구'].format(church=CHURCH,name=name,jik=jik,pastor=PASTOR)[:56]}…")
    else:
        base=CAREMSG.get(kind,"")
        if any(x in kind for x in ("장례","입원","소천","별세","위독")):
            msg=f"[{CHURCH}] {name}{jik+' ' if jik else ''}님 가정에 {base or '위로를 전합니다'}. {PASTOR}와 온 성도가 함께 기도하며 위로합니다. 주님의 위로와 평강이 함께하시길 빕니다."
        else:
            msg=f"[{CHURCH}] {name}{jik+' ' if jik else ''}님, {base or '축하합니다'}! {PASTOR}와 온 성도가 함께 기뻐하며 축복합니다. 하나님의 은혜가 늘 함께하시길 기도합니다."
        out=os.path.join(CAT("07"),f"[축하위로_{sanit(kind)}] {sanit(name)}_{today()}.txt")
        open(out,'w',encoding='utf-8').write(msg); _openfile(out)
        print(f"✅ {kind} 문자 초안(카톡 붙여넣기): {out}\n    → {msg}")
    if "출생" in kind: print("  ※ 아기를 교적에 올리려면 '교우 등록'에서 등록해 주세요(생일 자동 관리).")
def careevent_list(a):
    """경조사 목록 — 종류(출생/장례/결혼 등) 필터."""
    db=load(); es=db.get("경조사",[])
    if a.kind: es=[e for e in es if a.kind in e.get("종류","")]
    print(f"■ 경조사 {len(es)}건"+(f" · {a.kind}" if a.kind else ""))
    for e in es[:40]:
        print(f"  · {e['날짜']} [{e.get('종류','')}] {e.get('대상','')} — {e.get('내용','')}"+(f" · 경조금 {e['경조금']}" if e.get('경조금') else ""))
    if not es: print("  (없음 — '통합 경조사'에서 기록해 주세요)")
def calendar_cmd(a):
    """교회력·절기 — 그 해 부활절 기준 절기 날짜 자동 산출(설교·주보 준비 기준)."""
    import datetime as _dt
    y=int(a.year or today()[:4])
    A=y%19; B=y//100; Cc=y%100; D=B//4; E=B%4; F=(B+8)//25; G=(B-F+1)//3
    H=(19*A+B-D-G+15)%30; I=Cc//4; K=Cc%4; L=(32+2*E+2*I-H-K)%7; M=(A+11*H+22*L)//451
    mo=(H+L-7*M+114)//31; da=((H+L-7*M+114)%31)+1
    easter=_dt.date(y,mo,da)
    def fs(mon):
        d=_dt.date(y,mon,1); return d+_dt.timedelta(days=(6-d.weekday())%7)
    xmas=_dt.date(y,12,25); sun4=xmas-_dt.timedelta(days=(xmas.weekday()+1)%7); adv1=sun4-_dt.timedelta(days=21)
    items=[("재의 수요일(사순절 시작)",easter-_dt.timedelta(days=46)),("종려주일",easter-_dt.timedelta(days=7)),
      ("고난주간 시작",easter-_dt.timedelta(days=7)),("부활절",easter),("성령강림절(오순절)",easter+_dt.timedelta(days=49)),
      ("맥추감사주일(7월 첫주)",fs(7)),("추수감사주일(11월 셋째주)",fs(11)+_dt.timedelta(days=14)),
      ("대림절 첫 주일",adv1),("성탄절",xmas)]
    print(f"■ {y}년 교회력·절기 (설교·주보 준비 기준)")
    for name,dt in items: print(f"  · {dt.isoformat()}  {name}")
def giving_ledger(a):
    """교인별·항목별 헌금대장 — 헌금 기록을 (교인×종류)로 집계(심방·영수증 근거)."""
    from collections import defaultdict
    db=load(); rec=[r for r in db["재정"] if r["구분"]=="수입"]
    if a.year: rec=[r for r in rec if r["날짜"].startswith(a.year)]
    by=defaultdict(lambda: defaultdict(int)); tot=defaultdict(int)
    for r in rec:
        who=r.get("교인") or "무명"; it=r.get("항목") or "헌금"; by[who][it]+=r["금액"]; tot[who]+=r["금액"]
    if a.name:
        print(f"■ {a.name} 헌금 내역{f' ({a.year})' if a.year else ''}")
        for it,amt in by.get(a.name,{}).items(): print(f"  · {it}: {amt:,}원")
        print(f"  합계: {tot.get(a.name,0):,}원")
    else:
        print(f"■ 교인별 헌금 합계{f' ({a.year})' if a.year else ''} ({len(tot)}명)")
        for who,amt in sorted(tot.items(),key=lambda x:-x[1]): print(f"  · {who}: {amt:,}원")
        if not tot: print("  (헌금 기록 없음 — '재정 기록(수입/지출)'에서 입력해 주세요)")
def donation_receipt(a):
    """연말정산 기부금영수증 — 교인의 연간 헌금 합계로 영수증(docx) 발급."""
    db=load(); y=a.year or today()[:4]
    rec=[r for r in db["재정"] if r["구분"]=="수입" and r.get("교인")==a.name and r["날짜"].startswith(y)]
    total=sum(r["금액"] for r in rec); m=find(db,a.name); addr=m[0].get("주소","") if m else ""
    d=newdoc()
    para(d,"기부금 영수증",22,NAVY,True,AL.CENTER,after=2,font=SANS); hr(d,"A67C1E",12)
    para(d,f"제 {y}-{sanit(a.name or '')} 호",10,LG,after=12,font=SANS)
    for k,v in [("성    명",a.name or ""),("주    소",addr),("기부기간",f"{y}. 1. 1 ~ {y}. 12. 31"),("기부금액",f"{total:,} 원"),("기 부 처",CHURCH)]:
        para(d,f"{k} :   {v}",12,GRAY,after=6)
    para(d,"위와 같이 기부하였음을 증명합니다. (소득세법에 따른 기부금 세액공제용)",11,NAVY,True,before=10,after=16)
    para(d,f"{y}년      월      일",11,GRAY,AL.CENTER,after=14)
    para(d,f"{CHURCH}   담임 {PASTOR}   (인)",13,NAVY,True,AL.CENTER,font=SANS)
    out=os.path.join(CAT("06"),f"[기부금영수증] {sanit(a.name or '')}_{y}.docx"); _savedoc(d,out)
    print(f"✅ 기부금영수증 발급: {out} · {a.name} {y}년 {total:,}원")
def vip_add(a):
    """태신자(전도 대상) 등록 — 담당 성도·중보기도(전도 관리)."""
    db=load(); db.setdefault("태신자",[])
    db["태신자"].append({"이름":a.name or "","담당":a.sponsor or "","연락처":a.tel or "","상태":"기도중","접촉":[],"등록일":today(),"메모":a.memo or ""})
    save(db); print(f"✅ 태신자 등록: {a.name or ''}(담당 {a.sponsor or '-'}) — 매주 중보기도·편지·연락으로 품기")
def vip_contact(a):
    db=load(); v=next((x for x in db.get("태신자",[]) if x["이름"]==a.name),None)
    if not v: print(f"✗ 태신자 '{a.name}' 없음"); return
    v.setdefault("접촉",[]).append({"날짜":a.date or today(),"방법":a.how or "기도","내용":a.note or ""}); save(db)
    print(f"✅ [{a.name}] 접촉: {a.how or '기도'} {a.note or ''} (총 {len(v['접촉'])}회)")
def vip_list(a):
    db=load(); vs=db.get("태신자",[])
    print(f"■ 태신자 {len(vs)}명")
    for v in vs: print(f"  · {v['이름']} (담당 {v.get('담당','')}) · 접촉 {len(v.get('접촉',[]))}회 [{v.get('상태','')}]")
    if not vs: print("  (없음 — '태신자(전도대상) 등록'에서 등록해 주세요)")
def vip_convert(a):
    """태신자 결신 → 새가족 자동 등록."""
    import types
    db=load(); v=next((x for x in db.get("태신자",[]) if x["이름"]==a.name),None)
    if not v: print(f"✗ 태신자 '{a.name}' 없음"); return
    v["상태"]="결신"; save(db)
    newfamily_add(types.SimpleNamespace(name=v["이름"],tel=v.get("연락처",""),leader=v.get("담당",""),cell="",sex="",birth="",addr="",memo="태신자→결신",date=""))
    print(f"✅ {a.name} 결신 → 새가족 등록 완료 (정착 관리 시작)")
def presbytery_add(a):
    """노회 관련 — 회의·직무·제출서류·노회비·시찰 등 기록(기한 관리)."""
    db=load(); db.setdefault("노회",[])
    db["노회"].append({"날짜":a.date or today(),"구분":a.kind or "회의","내용":a.note or "","담당직무":a.role or "","기한":a.due or "","상태":"진행"})
    db["노회"].sort(key=lambda e:e["날짜"],reverse=True); save(db)
    print(f"✅ 노회 항목: {a.date or today()} [{a.kind or '회의'}] {a.note or ''}"+(f" · 직무 {a.role}" if a.role else "")+(f" · 기한 {a.due}" if a.due else ""))
def presbytery_list(a):
    """노회 관련 목록(회의·직무·서류·기한)."""
    db=load(); es=db.get("노회",[])
    print(f"■ 노회 관련 {len(es)}건")
    for e in es[:40]:
        print(f"  · {e['날짜']} [{e.get('구분','')}] {e.get('내용','')}"+(f" · 직무 {e['담당직무']}" if e.get('담당직무') else "")+(f" · 기한 {e['기한']}" if e.get('기한') else ""))
    if not es: print("  (없음 — '노회 관련'에서 기록해 주세요. 예: 시찰회·노회수련회·노회비·서류제출·총대)")
def event_plan(a):
    """대외·교회 행사 준비 기획서 — 개요·준비물·담당배정·당일진행·사후정리(docx)."""
    if not a.title: print("⚠ 행사명을 입력해 주세요. (예: event-plan --title 추수감사절 --date 2026-11-15)"); return
    d=newdoc()
    para(d,f"{a.title or '행사'} 준비 기획서",19,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,f"{a.date or ''} · {CHURCH}",10,GRAY,AL.CENTER,after=6,font=SANS); hr(d,"C25A2E",12)
    for l,v in [("행사명",a.title or ""),("일시",a.date or ""),("장소",a.place or ""),("대상·예상인원",""),("주제·목적",a.theme or ""),("총괄 담당",a.host or "")]:
        para(d,f"{l} :   {v}",11,GRAY,after=4)
    groups=[("준비 항목",["홍보(현수막·포스터·SNS)","초청·섭외(강사/게스트)","예배·프로그램 순서","찬양·음향·영상","등록·안내·주차","다과·식사","재정·예산 집행","사진·기록"]),
      ("담당 배정",["총괄","프로그램","섭외","홍보","음향·영상","식사·다과","안내·주차","재정"]),
      ("당일 진행",["리허설·세팅","인원 점검","순서 진행","정리·마무리"]),
      ("행사 후",["결과 보고","감사 인사","사진·영상 정리","재정 정산"])]
    for g,items in groups:
        para(d,g,13,RGBColor(0xC2,0x5A,0x2E),True,before=8,after=3,font=SANS)
        for it in items: para(d,f"  ☐  {it}",11,after=2)
    hr(d,"CCCCCC",6)
    para(d,"★ 준비 3원칙",12,MAROON,True,before=6,after=3,font=SANS)
    for t in ["① 목표를 한 문장으로 먼저 정하기 — 이 행사로 무엇이 남기를 바라는지가 모든 결정의 기준입니다.",
              "② 장소·프로그램·예산을 초기에 함께 결정 — 셋은 서로 맞물립니다(장소가 프로그램을, 예산이 장소를 정합니다).",
              "③ 담당을 일찍 세워 위임 — 목표만 공유하고 맡기십시오. 함께 준비할 때 더 멀리 갑니다."]:
        para(d,t,10,after=2)
    para(d,"※ 수련회·여름사역이라면 '여름성경학교·수련회 준비' 참고자료 카드를 꼭 함께 보세요 — 장소 고르는 법(유형별 장단점·선택 기준·비교평가표·답사 체크리스트·예약계약), 프로그램 설계(6요소·시간배분·하루 리듬), D-60 준비 로드맵이 자세히 담겨 있습니다.",9,LG,before=8,font=SANS)
    out=os.path.join(CAT("07"),f"[행사기획] {sanit(a.title or '행사')}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 행사 준비 기획서: {out} (준비 3원칙·참고자료 안내 포함)")
_GREET_CACHE={}
def _greet_pool(kind):
    """축하 문구 풀 로드(100편). kind='생일'|'결혼기념일'. _data/<파일>.json. 없으면 []."""
    import json as _j
    fn={"생일":"생일축하문구.json","결혼기념일":"결혼기념문구.json"}.get(kind)
    if not fn: return []
    if fn in _GREET_CACHE: return _GREET_CACHE[fn]
    try: pool=_j.load(open(os.path.join(BASE,"_data",fn),encoding='utf-8'))
    except Exception: pool=[]
    _GREET_CACHE[fn]=pool; return pool
def _pick_greet(m, kind, year, pool):
    """성도의 그해 축하문구를 고른다. 연도별 이력(m['축하이력'][kind]={연도:번호})으로 매년 다른 문구.
       반환 (문구dict, 이전이력[(연도,번호)], 새로기록됨?). 이미 올해 고른 게 있으면 그대로(재생성 일관)."""
    if not pool: return (None,[],False)
    hist=m.setdefault('축하이력',{}).setdefault(kind,{})
    by_num={it.get('번호'):it for it in pool}; ys=str(year)
    prev=[(y,hist[y]) for y in sorted(hist) if y!=ys and hist[y] in by_num]
    if ys in hist and hist[ys] in by_num: return (by_num[hist[ys]], prev, False)
    used=set(hist.values())
    avail=[it for it in pool if it.get('번호') not in used] or pool   # 다 썼으면(100년!) 다시 전체
    chosen=avail[(int(m.get('id',0))+int(year))%len(avail)]
    hist[ys]=chosen.get('번호')
    return (chosen, prev, True)
def birthday(a):
    """🎂 성도 생일·결혼기념일 축하 — 다가오는 것(향후 N일·기본30) 조회 + 축하 문자 자동생성(카톡 붙여넣기·100문구 연도순환) + 셀 명단.
       ★매년 다른 문구가 자동 선택되고, 지난해 보낸 문구도 함께 보여드립니다(중복 방지)."""
    import datetime as _dt
    from collections import defaultdict
    db=load(); days=_intarg(a.days, 30, "일수"); td=_dt.date.today()
    def nextd(mmdd):
        try:
            mo,da=int(mmdd[:2]),int(mmdd[3:5]); d=_dt.date(td.year,mo,da)
            if d<td: d=_dt.date(td.year+1,mo,da)
            return d
        except Exception: return None
    ups=[]
    for m in db["교인"]:
        for kind,field,ic in [("생일","생년월일","🎂"),("결혼기념일","결혼기념일","💍")]:
            v=(m.get(field) or "")[5:]
            if not v: continue
            nd=nextd(v)
            if nd and 0<=(nd-td).days<=days: ups.append(((nd-td).days,nd,ic,kind,m))
    ups.sort(key=lambda x:x[0])
    print(f"■ 다가오는 {days}일 이내 생일·결혼기념일 {len(ups)}건")
    bycell=defaultdict(list)
    for row in ups: bycell[row[4].get('소속셀') or '미배정'].append(row)
    for cell,items in bycell.items():
        print(f"  [{cell}]")
        for dd,nd,ic,kind,m in items:
            tag="오늘" if dd==0 else f"D-{dd}"
            print(f"    {ic} {tag} {nd.isoformat()} {m['이름']} {m.get('직분','')} {kind} · {m.get('연락처','')}")
    if not ups: print("  (해당 없음 — 교인 생년월일·결혼기념일을 입력하세요)"); return
    print("  ─ 축하 문자 초안(카톡 붙여넣기 · 매년 다른 문구) ─")
    kitcell=defaultdict(list); _dirty=False
    for dd,nd,ic,kind,m in ups:
        nm=m['이름']; jik=m.get('직분','') or ''; tel=m.get('연락처','') or ''; year=nd.year
        pool=_greet_pool(kind); by_num={it.get('번호'):it for it in pool}
        chosen,prev,newly=_pick_greet(m,kind,year,pool)
        if newly: _dirty=True
        if chosen:
            msg=chosen['문구'].format(church=CHURCH,name=nm,jik=jik,pastor=PASTOR)
        elif kind=="생일":   # 풀이 없을 때 안전망 기본문구
            msg=(f"{CHURCH} {nm} {jik}님, 생일을 진심으로 축하드립니다! 🎂 주님의 긍휼이 아침마다 새롭듯(애 3:22-23), "
                 f"새 한 해에 은혜가 날마다 새롭기를 {PASTOR}와 온 성도가 축복합니다. 샬롬 🙏")
        else:
            msg=(f"{CHURCH} {nm} {jik}님 부부, 결혼기념일을 진심으로 축하드립니다! 💍 세 겹줄은 쉬 끊어지지 않듯(전 4:12) "
                 f"두 분의 언약에 사랑과 은혜가 늘 새롭기를 {PASTOR}와 온 성도가 축복합니다. 샬롬 🙏")
        # 지난해 보낸 문구 요약(예전 것을 보며 다른 문구를 보내도록)
        prevnote=""
        if prev:
            ps=[]
            for (y,num) in prev[-3:]:
                sn=by_num.get(num,{}).get('문구','')
                if sn: sn=sn.format(church=CHURCH,name=nm,jik=jik,pastor=PASTOR)
                ps.append(f"{y}년: «{sn[:26]}…»" if sn else f"{y}년: (#{num})")
            prevnote="지난 축하 → "+"  /  ".join(ps)
        tag="오늘" if dd==0 else f"D-{dd}"
        print(f"    → {nm} ({tag}·{kind}): {msg.splitlines()[0][:40]}…"+(f"   [{prevnote}]" if prevnote else ""))
        kitcell[m.get('소속셀') or '미배정'].append((tag,nd,nm,jik,tel,kind,msg,prevnote))
    if _dirty:
        try: save(db)   # 올해 선택한 문구를 성도별 이력에 기록(내년엔 다른 문구가 나옴)
        except Exception: pass
    # ── 카톡에 붙여넣기 목록 파일 자동 생성·자동 열기 (아침에 복사→카톡 전송) ──
    L=[f"════ 생일·결혼기념일 축하 · 카톡에 붙여넣기 목록 ({td.isoformat()} 기준 · 향후 {days}일) ════",
       "※ 이 메모장을 열어, 아래 각 사람 블록을 복사해서 카톡에 붙여넣어 보내세요(무료·비용0).",
       "※ 생일 '당일 아침'에 보내면 성도들이 더 감격합니다. 이름·직분을 살짝 다듬어 더 개인적으로 보내셔도 좋습니다.",
       "※ 문구는 100가지 중 매년 다른 것이 자동 선택됩니다. '지난 축하'로 예전에 보낸 문구를 보며 겹치지 않게 하실 수 있습니다.",""]
    for cell,rows in kitcell.items():
        L.append(f"───── [{cell}] {len(rows)}명 ─────")
        for tag,nd,nm,jik,tel,kind,msg,prevnote in rows:
            L.append(f"▼ {tag} {nd.isoformat()} · {nm} {jik} ({tel}) · {kind}")
            if prevnote: L.append(f"  ({prevnote})")
            L.append(msg); L.append("")
    out=os.path.join(CAT("07"),f"[생일축하_카톡에붙여넣기] {td.isoformat()}.txt")
    open(out,'w',encoding='utf-8').write("\n".join(L)); _openfile(out)
    print(f"  ✅ 카톡에 붙여넣기 목록 저장: {out}  (열어서 복사→카톡 전송)")
def hwp_convert(a):
    """DOCX → 한글(.hwp) 변환 — 생성된 문서를 네이티브 한글파일로. (Windows+한글+pywin32 필요)
       ※ docx는 한글에서 그대로 열립니다. .hwp가 꼭 필요할 때만 사용."""
    import subprocess, time
    try: import win32com.client as w
    except Exception: print("✗ .hwp 변환은 Windows + 한글(한컴오피스) + pywin32 설치 필요. (docx는 한글에서 바로 열립니다)"); return
    targets=[]
    if a.file and os.path.exists(a.file): targets=[a.file]
    else:
        n=int(a.recent or 10); allx=[]
        for dp,_,fs in os.walk(ROOT):
            if _skip_scan(dp): continue
            for f in fs:
                if f.endswith(".docx"): allx.append(os.path.join(dp,f))
        allx.sort(key=lambda p:os.path.getmtime(p),reverse=True); targets=allx[:n]
    if not targets: print("변환할 docx 없음"); return
    try: subprocess.run(["taskkill","/F","/IM","Hwp.exe"],capture_output=True)
    except Exception: pass
    hwp=w.gencache.EnsureDispatch("HWPFrame.HwpObject")
    try: hwp.RegisterModule("FilePathCheckDLL","FilePathCheckerModule")
    except Exception: pass
    hwp.SetMessageBoxMode(0x20000); done=0
    for p in targets:
        out=os.path.splitext(p)[0]+".hwp"
        try: hwp.Open(p); hwp.SaveAs(out,"HWP"); hwp.Clear(1); done+=1
        except Exception: pass
    try: hwp.Quit()
    except Exception: pass
    print(f"✅ 한글(.hwp) 변환 {done}개 완료 (docx와 같은 폴더에 저장)")
def print_file(a):
    """프린터 직접 출력 — 파일(docx/xlsx/pdf/hwp/pptx)을 기본 프린터로 전송. (Windows)"""
    p=a.file
    if not p or not os.path.exists(p): print("사용법: print-file --file \"파일경로\"  (또는 파일 열어 Ctrl+P)"); return
    try: os.startfile(p,"print"); print(f"🖨️ 인쇄 전송: {os.path.basename(p)} → 기본 프린터")
    except Exception as e: print(f"✗ 인쇄 실패: {e}. 파일을 열어 Ctrl+P로 인쇄하세요.")
def read_file(a):
    """파일 읽기 — TXT/DOCX 내용을 읽어 표시(옛 설교·명단·자료 불러오기)."""
    p=a.file
    if not p or not os.path.exists(p): print("사용법: read-file --file \"파일경로\"  (.txt/.docx)"); return
    ext=os.path.splitext(p)[1].lower()
    if ext==".txt":
        raw=open(p,'rb').read()
        if raw[:2] in (b'\xff\xfe',b'\xfe\xff'): t=raw.decode('utf-16',errors='replace')
        else:
            try: t=raw.decode('utf-8')
            except UnicodeDecodeError: t=raw.decode('cp949',errors='replace')
    elif ext==".docx":
        t="\n".join(x.text for x in Document(p).paragraphs)
    else: print(f"✗ 지원 형식: .txt / .docx (현재 {ext}). 엑셀은 export-excel, PDF는 뷰어로."); return
    print(f"■ {os.path.basename(p)} ({len(t)}자)")
    print(t[:3000]+("\n... (이하 생략, 전체는 파일에서)" if len(t)>3000 else ""))
def ppt_lyrics(a):
    """찬양 가사 PPT — 예배 프로젝터용 파워포인트 슬라이드(등록곡 가사, 한 소절씩)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt as PPt
    from pptx.dml.color import RGBColor as PColor
    from pptx.enum.text import PP_ALIGN
    import json as _j
    lib=_j.load(open(SONGLIB,encoding='utf-8')) if os.path.exists(SONGLIB) else []
    byn={re.sub(r'\s','',s['제목']):s for s in lib}
    songs=[s.strip() for s in (a.songs or "").split(";") if s.strip()]
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); blank=prs.slide_layouts[6]
    def slide(text,size=40):
        s=prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb=PColor(0x0f,0x12,0x26)
        tb=s.shapes.add_textbox(Inches(0.6),Inches(2.3),Inches(12.1),Inches(3.0)); tf=tb.text_frame; tf.word_wrap=True
        for i,ln in enumerate(text.split("\n")):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER
            r=p.add_run(); r.text=ln; r.font.size=PPt(size); r.font.bold=True; r.font.color.rgb=PColor(0xff,0xff,0xff); r.font.name="맑은 고딕"
    found=0
    for q in songs:
        nk=re.sub(r'\s','',q); sng=None
        for k,v in byn.items():
            if nk in k or k in nk: sng=v; break
        title=sng['제목'] if sng else q; lyr=(sng.get('가사') if sng else '') or ''
        slide(title,54)
        for block in re.split(r'\n\s*\n',lyr):
            if block.strip(): slide(block.strip(),40); found+=1
    out=os.path.join(CAT("07"),f"[찬양PPT] {sanit(a.title or '찬양')}_{today()}.pptx"); prs.save(out); _ext_copy(out); _openfile(out)
    print(f"✅ 찬양 가사 PPT: {out} ({found}장)")
def agenda(a):
    """일정 캘린더 — 다가오는 모든 일정(집회·외부설교·행사·생일·기념일·경조사)을 날짜순 한눈에."""
    import datetime as _dt
    db=load(); days=_intarg(a.days, 30, "일수"); td=_dt.date.today(); items=[]
    for e in db.get("일정",[]):
        try:
            d=_dt.date.fromisoformat(e["날짜"])
            if 0<=(d-td).days<=days: items.append((d,f"📅 [{e.get('유형','일정')}] {e.get('장소','')} {e.get('주제','')}"))
        except Exception: pass
    def nextd(mmdd):
        try:
            mo,da=int(mmdd[:2]),int(mmdd[3:5]); d=_dt.date(td.year,mo,da)
            return d if d>=td else _dt.date(td.year+1,mo,da)
        except Exception: return None
    for m in db["교인"]:
        for kind,field,ic in [("생일","생년월일","🎂"),("결혼기념일","결혼기념일","💍")]:
            v=(m.get(field) or "")[5:]
            if v:
                nd=nextd(v)
                if nd and 0<=(nd-td).days<=days: items.append((nd,f"{ic} {m['이름']} {kind}"))
    for e in db.get("경조사",[]):
        try:
            d=_dt.date.fromisoformat(e["날짜"])
            if 0<=(d-td).days<=days: items.append((d,f"🕯️ [{e.get('종류','')}] {e.get('대상','')}"))
        except Exception: pass
    # 교회력 절기(올해·내년)
    for y in (td.year, td.year+1):
        A=y%19;B=y//100;Cc=y%100;D=B//4;E=B%4;Fx=(B+8)//25;G=(B-Fx+1)//3
        H=(19*A+B-D-G+15)%30;I=Cc//4;K=Cc%4;L=(32+2*E+2*I-H-K)%7;Mx=(A+11*H+22*L)//451
        ea=_dt.date(y,(H+L-7*Mx+114)//31,((H+L-7*Mx+114)%31)+1)
        def _fs(mon,yy=y):
            dd=_dt.date(yy,mon,1); return dd+_dt.timedelta(days=(6-dd.weekday())%7)
        xm=_dt.date(y,12,25); s4=xm-_dt.timedelta(days=(xm.weekday()+1)%7); a1=s4-_dt.timedelta(days=21)
        for name,dt in [("사순절 시작",ea-_dt.timedelta(days=46)),("종려주일",ea-_dt.timedelta(days=7)),("부활절",ea),
                        ("성령강림절",ea+_dt.timedelta(days=49)),("맥추감사주일",_fs(7)),("추수감사주일",_fs(11)+_dt.timedelta(days=14)),
                        ("대림절 첫 주일",a1),("성탄절",xm)]:
            if 0<=(dt-td).days<=days: items.append((dt,f"✝️ {name} (교회력)"))
    items.sort(key=lambda x:x[0])
    print(f"■ 다가오는 {days}일 일정 캘린더 ({len(items)}건)")
    for d,label in items:
        dd=(d-td).days; tag="오늘" if dd==0 else ("내일" if dd==1 else f"D-{dd}")
        print(f"  {d.isoformat()} ({tag})  {label}")
    if not items: print("  (예정 일정 없음 — '집회·외부설교 일정'·생일 입력 시 여기 모입니다)")
def sermon_files(a):
    """지난 설교 (목록·열기·재활용) 통합 — 비우면 목록, '열기'에 제목=그 파일 열기, '재활용'에 검색어=지난 설교로 새 초안 작성."""
    qy=(getattr(a,'query','') or "").strip()
    if qy:   # 이력 검색: 설교 기록(본문·제목)에서 찾기
        import types
        return sermon_list(types.SimpleNamespace(query=qy,service=getattr(a,'service','') or ""))
    bases=[CAT("09")]
    bd=_C.get("백업폴더")
    if bd and os.path.isdir(os.path.join(bd,"09 설교")): bases.append(os.path.join(bd,"09 설교"))
    raw=[]
    for base in bases:
        if not os.path.isdir(base): continue
        for dp,_,fs in os.walk(base):
            for f in fs:
                if f.lower().endswith((".docx",".hwp")) and "오늘의묵상" not in f:
                    p=os.path.join(dp,f)
                    try: raw.append((os.path.getmtime(p),p,f,os.path.relpath(dp,base)))
                    except Exception: pass
    raw.sort(reverse=True)
    seen=set(); files=[]   # 파일명 중복 제거(프로그램폴더·D백업 겹침)
    for t in raw:
        if t[2].lower() in seen: continue
        seen.add(t[2].lower()); files.append(t)
    if a.open:
        hit=[t for t in files if a.open in t[2]]
        if hit:
            try: os.startfile(hit[0][1]); print(f"▶ 여는 중: {hit[0][2]}  (Word/한글에서 편집하세요)")
            except Exception as e: print(f"✗ 열기 실패: {e} — 파일을 직접 더블클릭하세요")
        else: print(f"✗ '{a.open}' 관련 설교 파일 없음. 목록에서 확인하세요.")
        return
    print(f"■ 작성했던 설교 파일 {len(files)}개 — 아래를 눌러 여세요")
    for _,p,f,rel in files[:50]:
        print(f"▶열기|{p}|{f}")
    if not files: print("  (이 폴더엔 설교 파일이 없습니다 — 아래 ‘설교 폴더 열기’를 눌러 예전 설교가 있는 C·D·USB 폴더로 이동해 찾으세요)")
    _sermon_folders()   # ★ 폴더 열기 버튼(주보 카드처럼) — C·D·USB의 예전 설교로 직접 이동
def _find_docs(kw, limit=20):
    """이름·키워드로 관련 문서 검색 — 프로그램·D백업·문서폴더(설정 '문서폴더')."""
    if not kw: return []
    roots=[ROOT]
    for r in (_C.get("백업폴더"), _C.get("문서폴더")):
        if r and os.path.isdir(r) and os.path.abspath(r) not in [os.path.abspath(x) for x in roots]: roots.append(r)
    exts=(".docx",".hwp",".pdf",".txt",".xlsx",".pptx"); seen=set(); hits=[]
    for r in roots:
        for dp,_,fs in os.walk(r):
            if _skip_scan(dp): continue
            for f in fs:
                if kw in f and f.lower().endswith(exts) and f.lower() not in seen:
                    seen.add(f.lower()); hits.append(os.path.join(dp,f))
                    if len(hits)>=limit: return hits
    return hits
def open_file(a):
    """파일 불러오기 — 프로그램·D백업·문서폴더에 저장된 파일을 이름으로 찾아 엽니다."""
    if a.file and (os.path.exists(a.file) or str(a.file).startswith("shell:")):
        try: os.startfile(a.file); print(f"▶ 여는 중: {os.path.basename(a.file) or a.file}")
        except Exception as e: print(f"✗ 열기 실패: {e} — 파일을 직접 더블클릭하세요")
        return
    kw=(getattr(a,'name','') or "").strip()
    if not kw:
        return my_files(a)   # 이름 없으면 → 종류별 내 작업물 목록(눌러 열기) · 통합
    roots=[ROOT]
    for r in (_C.get("백업폴더"), _C.get("문서폴더")):
        if r and os.path.isdir(r) and os.path.abspath(r) not in [os.path.abspath(x) for x in roots]: roots.append(r)
    exts=(".docx",".hwp",".hwpx",".xlsx",".pptx",".txt",".pdf",".jpg",".jpeg",".png"); seen=set(); hits=[]
    for r in roots:
        for dp,_,fs in os.walk(r):
            if _skip_scan(dp): continue
            for f in fs:
                if kw in f and f.lower().endswith(exts) and f.lower() not in seen:
                    seen.add(f.lower()); hits.append(os.path.join(dp,f))
    if not hits: print(f"✗ '{kw}' 관련 파일 없음"); return
    if len(hits)==1:
        try: os.startfile(hits[0]); print(f"▶ 여는 중: {os.path.basename(hits[0])}")
        except Exception: print(f"파일 위치: {hits[0]}")
        return
    print(f"■ '{kw}' 관련 {len(hits)}개 — 아래를 눌러 여세요")
    for h in hits[:20]: print(f"▶열기|{h}|{os.path.basename(h)}")
def open_folder(a):
    """폴더 열기 — 설교/심방/재정 등 해당 폴더를 탐색기(윈도우 파일창)로 엽니다. 그 안에서 파일을 더블클릭해 여세요."""
    NMAP={"설교":"09","묵상":"09","심방":"04","목양":"04","새가족":"04","셀":"04",
     "주보":"07","행사":"07","홍보":"07","찬양":"07","콘티":"07","악보":"07",
     "예배":"02","예식":"02","성례":"02","교육":"03","교안":"03","훈련":"03","공과":"03",
     "선교":"05","재정":"06","증명":"06","증명서":"06","헌금":"06","영수증":"06",
     "대외":"08","노회":"08","시설":"08","조직":"01","기본":"01"}
    kind=(getattr(a,"kind",None) or "설교").strip()
    n=NMAP.get(kind,"09")
    folder=CAT(n)   # 설정 '아카이브루트' 있으면 실제 교회 원본 아카이브를 엶
    try:
        os.startfile(folder); print(f"▶ '{os.path.basename(folder)}' 폴더를 탐색기로 열었습니다.\n   그 안에서 지난 파일들을 더블클릭해 여세요.")
    except Exception as e:
        print(f"폴더 위치:\n{folder}\n(자동 열기 실패: {e} — 위 경로를 복사해 탐색기 주소창에 붙여넣으세요)")
def my_files(a):
    """내 작업물 열기 — 종류별(설교/심방/주보/묵상/찬양/악보/증명서/재정/교안)로 D폴더·USB·백업의 작업물을 최근순으로. 눌러서 엽니다."""
    KMAP={"설교":["설교"],"설교문":["설교"],"심방":["심방"],"주보":["주보"],
     "묵상":["묵상","QT"],"찬양":["콘티","찬양"],"콘티":["콘티","찬양"],"악보":["악보"],
     "작곡":["악보","자작곡","작곡"],"증명서":["증명"],"증명":["증명"],
     "재정":["헌금","재정","영수증","기부금","결산"],"헌금":["헌금","영수증","기부금"],
     "교안":["교안","공과","교육"]}
    kind=(getattr(a,"kind",None) or "전체").strip()
    kws=KMAP.get(kind)
    all_drive=(kind!="전체")   # 전체=프로그램 작업물만, 종류지정=D폴더 과거자료까지
    roots=[ROOT]
    extra=[_C.get("백업폴더")]+([_C.get("문서폴더")] if all_drive else [])
    for r in extra:
        if r and os.path.isdir(r) and os.path.abspath(r) not in [os.path.abspath(x) for x in roots]: roots.append(r)
    exts=(".docx",".hwp",".hwpx",".pdf",".txt",".xlsx",".pptx"); seen=set(); raw=[]
    for r in roots:
        for dp,_,fs in os.walk(r):
            if _skip_scan(dp): continue
            for f in fs:
                if f.startswith("~") or not f.lower().endswith(exts): continue
                if kws and not any(k in (dp+os.sep+f) for k in kws): continue
                if f.lower() in seen: continue
                seen.add(f.lower()); p=os.path.join(dp,f)
                try: raw.append((os.path.getmtime(p),p,f))
                except Exception: pass
    raw.sort(reverse=True)
    if not raw: print(f"✗ '{kind}' 관련 작업물이 아직 없습니다 — 만들면 여기서 바로 열립니다"); return
    print(f"■ '{kind}' 작업물 {len(raw)}개 (최근 작업순) — 아래를 눌러 여세요")
    for _,p,f in raw[:40]: print(f"▶열기|{p}|{f}")
def song_sheet(a):
    """작곡한 곡 악보 열기 — 등록된 곡의 악보 파일을 열어 다듬기."""
    import json as _j
    if not os.path.exists(SONGLIB): print("등록 곡 없음 — '찬양곡·자작곡 등록'에서 등록해 주세요"); return
    lib=_j.load(open(SONGLIB,encoding='utf-8'))
    hit=[s for s in lib if a.title and a.title in s.get("제목","")]
    if not hit: print(f"✗ '{a.title}' 곡 없음"); return
    p=hit[0].get("악보","")
    if p and os.path.exists(p):
        try: os.startfile(p); print(f"▶ 악보 여는 중: {hit[0]['제목']}")
        except Exception as e: print(f"✗ {e}")
    else: print(f"✗ '{hit[0]['제목']}' 악보 파일 없음 ('찬양곡·자작곡 등록'에서 악보를 함께 등록해 주세요)")
def menu(a):
    """간편 메뉴 — 번호만 누르면 되는 대화형(명령어 몰라도 사용). 시작.bat로 실행."""
    import types
    def NS(**kw): return types.SimpleNamespace(**kw)
    def ask(p): return input(p).strip()
    while True:
        print(f"\n══════ {CHURCH} 교회 종합행정 ══════")
        upd=_check_update()
        if upd: print(f"  🔔 새 업데이트 v{upd[0]} 있습니다!  →  99 를 눌러 업데이트 (교인·설교 등 자료는 영구 보존)")
        print("  1 교인등록    2 심방브리핑   3 심방기록    4 증명서발급")
        print("  5 주보만들기  6 오늘의묵상   7 생일자확인  8 찬양콘티")
        print("  9 성경찾기   10 엑셀정리    11 돌봄필요   12 주간브리핑")
        print(" 13 집회일정   14 학생등록    15 시험응원   16 선교준비")
        print(" 17 설교이력   18 설교작성    19 지난설교 재생성")
        print(" 20 기능요청   77 피닉스복구  88 자료백업   99 업데이트")
        print("  0 종료")
        c=ask("번호 선택> ")
        try:
            if c=="0": print("종료합니다. 샬롬 🙏"); break
            elif c=="1": member_add(NS(name=ask("이름: "),role=ask("직분: "),cell=ask("소속셀: "),tel=ask("연락처: "),birth=ask("생년월일(YYYY-MM-DD): "),sex="",addr="",baptism="",leader="",memo="",date=""))
            elif c=="2": visit_brief(NS(name=ask("교인 이름: ")))
            elif c=="3": visit_add(NS(name=ask("이름: "),word=ask("전한 말씀: "),prayer=ask("기도제목(;로 구분): "),note=ask("나눈 내용: "),followup=ask("후속조치: "),kind="",by="",date=""))
            elif c=="4": cert(NS(name=ask("이름: "),kind=ask("종류(교인/세례/헌금): "),purpose=ask("용도: "),no=""))
            elif c=="5": bulletin(NS(sermon=ask("설교 제목: "),notice=ask("교회소식(;로 구분): "),week=ask("주간일정(;로 구분): "),order="",pray="",date=""))
            elif c=="6": devotion(NS(title=ask("묵상 제목: "),text=ask("본문(예 마11:28): "),verse="",body=ask("묵상 내용: "),apply="",pray="",theme="",date=""))
            elif c=="7": birthday(NS(days=ask("며칠 이내?(기본30): ") or "30"))
            elif c=="8": setlist(NS(songs=ask("곡목(;로 구분): "),size=(ask("용지(A4/A3): ") or "A4"),title=ask("콘티 제목: "),date=""))
            elif c=="9": bible(NS(ref=ask("찾을 곳(예: 요한복음 3:16): "),version=(ask("역본(kjv/web): ") or "web")))
            elif c=="10": export_excel(NS())
            elif c=="11": care(NS(days=(ask("기준 일수(기본 90): ") or "90")))
            elif c=="12": weekly_brief(NS())
            elif c=="13": schedule_add(NS(date=ask("날짜(YYYY-MM-DD): "),type=(ask("유형(집회/외부설교/행사/찬양사역): ") or "집회"),place=ask("장소/교회: "),host=ask("담당자: "),tel=ask("연락처: "),theme=ask("주제: "),fee=ask("사례: "),memo="")); schedule_list(NS(upcoming="1"))
            elif c=="14": student_add(NS(name=ask("이름: "),school=ask("학교: "),grade=ask("학년: "),dept=(ask("부서(중등부/고등부/청년부): ") or ""),tel=ask("연락처: "),guardian=ask("보호자: "),gtel=ask("보호자연락처: "),major="",sex="",birth=ask("생일(YYYY-MM-DD): "),addr="",leader="",memo="",date=""))
            elif c=="15": exam_cheer(NS(exam=(ask("시험명(예 중간고사): ") or "시험"),dept=(ask("부서(전체는 엔터): ") or "")))
            elif c=="16": mission_plan(NS(title=(ask("선교 팀명: ") or "단기선교"),dday=ask("출국일(YYYY-MM-DD, 모르면 엔터): ")))
            elif c=="17": sermon_list(NS(query=ask("설교 검색어(본문/제목, 전체는 엔터): "),text="",year=""))
            elif c=="18": sermon(NS(title=ask("설교 제목: "),text=ask("본문(예 요3:16): "),service=(ask("예배유형(새벽/수요/목요집회/금요/주일오전/주일오후/중고등부/초등부): ") or "주일오전예배"),theme=ask("주제(키워드): "),points=ask("대지(;로 구분, 모르면 엔터): "),series=ask("시리즈(없으면 엔터): "),date=""))
            elif c=="19": sermon_reuse(NS(query=ask("지난 설교 검색(본문/제목): "),title="",text="",theme="",service="",date=""))
            elif c=="20": request(NS(text=ask("원하는 기능을 적어주세요: ")))
            elif c=="88": backup(NS())
            elif c=="99":
                if upd: update(NS(file=upd[1])); print("  ✅ 업데이트 완료! 교인·설교 등 자료는 영구 보존되었습니다. 프로그램을 다시 실행하세요.")
                else: print("  현재 최신 버전입니다.")
            else: print("→ 없는 번호입니다.")
        except Exception as e: print("오류:",e)
VERSION="2026-07-21 (지난 설교·재활용에서 예전 설교 C·D·USB 바로 찾기 · 전 카드 '지난 자료 찾기' 버튼 · 안정화)"
# ★업데이트 발행 주소(깃허브 raw). 발행 스크립트가 목사님 계정으로 자동 채웁니다.
# 예) https://raw.githubusercontent.com/사용자명/저장소명/main/   ← 끝에 / 포함. 비어있으면 설정(업데이트기준URL) 또는 _업데이트 폴더 사용.
_UPDATE_BASE_DEFAULT="https://raw.githubusercontent.com/welikewon/church-admin/main/"
def _vt(v):
    try: return tuple(int(x) for x in str(v).split("."))
    except Exception: return (0,)
def _check_update():
    """_업데이트 폴더에 더 새로운 church.py가 있으면 (새버전, 경로) 반환."""
    p=os.path.join(BASE,"_업데이트","church.py")
    if not os.path.exists(p): return None
    try:
        m=re.search(r'VERSION\s*=\s*["\']([\d.]+)["\']',open(p,encoding='utf-8').read())
        if m and _vt(m.group(1))>_vt(VERSION): return (m.group(1),p)
    except Exception: pass
    return None
def version(a):
    print(f"■ {CHURCH} 교회 종합행정시스템 · 엔진 church.py v{VERSION}")
    up=_check_update()
    if up: print(f"  🔔 새 업데이트 v{up[0]} 있음! →  python church.py update --file \"{up[1]}\"")
    else:  print("  ✅ 현재 최신 버전입니다.")
    print("  🔄 업데이트 가능 — 화석이 아니라 계속 발전하는 프로그램입니다.")
    print("  🔒 자료는 영구 보존 — 업데이트해도 교인·심방·설교·재정·예화·주석·찬양곡 등")
    print("      모든 자료(church_db.json·church_config.json·_내자료)는 절대 삭제되지 않습니다.")
    ph=len([f for f in os.listdir(PHOENIX) if f.startswith("db_")]) if os.path.isdir(PHOENIX) else 0
    print(f"  🔥 피닉스 복구 — 실수로 자료가 지워져도 언제든 되살아납니다. (복구지점 {ph}개 보관 중)")
    print("      복구:  python church.py phoenix   (지점 확인) →  phoenix --last  (즉시 복구)")
def backup(a):
    """자료 백업 — 모든 데이터 + 내자료를 날짜 폴더에 복사(업데이트·오작동 대비)."""
    import shutil, datetime
    stamp=datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    bdir=os.path.join(BASE,"_백업",stamp); os.makedirs(bdir,exist_ok=True)
    for f in ("church_db.json","church_config.json"):
        p=os.path.join(BASE,f)
        if os.path.exists(p): shutil.copy2(p,os.path.join(bdir,f))
    nd=os.path.join(BASE,"_내자료")
    if os.path.isdir(nd): shutil.copytree(nd,os.path.join(bdir,"_내자료"),dirs_exist_ok=True)
    print(f"✅ 백업 완료: {bdir}")
    print("  (교인·심방·설교·재정·예화·주석·찬양곡 전부 · 업데이트해도 안전)")
    return bdir
def _upd_fetch(base, rel, timeout=25):
    """업데이트 소스(base=URL 또는 로컬폴더)에서 rel 파일을 바이트로 가져온다. http(s)면 다운로드, 아니면 로컬 읽기."""
    if str(base).startswith(("http://","https://")):
        import urllib.request, urllib.parse
        url=base.rstrip("/")+"/"+"/".join(urllib.parse.quote(p) for p in rel.split("/"))
        req=urllib.request.Request(url, headers={"User-Agent":"church-update"})
        with urllib.request.urlopen(req, timeout=timeout) as r: return r.read()
    with open(os.path.join(base, rel.replace("/",os.sep)),'rb') as f: return f.read()
def _ver_date(s):
    """버전 문자열에서 YYYY-MM-DD 추출(비교용). 없으면 (0,0,0)."""
    m=re.search(r'(\d{4})-(\d{2})-(\d{2})', str(s))
    return tuple(int(x) for x in m.groups()) if m else (0,0,0)
def update(a):
    """프로그램 업데이트 — 최신 코드·참고자료를 받아 교체(자료는 자동 백업 후 그대로 보존).
       ① --file 지정: 그 church.py로 로컬 교체(구방식).
       ② 아니면: 설정 '업데이트기준URL'(또는 _업데이트 폴더)의 manifest.json 기준으로 자동 다운로드·교체."""
    import shutil, datetime, json as _j
    # ── ① 구방식: 로컬 단일 church.py 교체 ──
    if a.file:
        src=a.file
        if not os.path.exists(src):
            print("사용법: update --file \"새 church.py 경로\""); return
        backup(a); stamp=datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        self_path=os.path.abspath(__file__)
        shutil.copy2(self_path, self_path+f".{stamp}.old"); shutil.copy2(src, self_path)
        print(f"✅ 업데이트 완료 — church.py 교체 (이전 버전은 church.py.{stamp}.old 로 보존)")
        print("  🔒 모든 자료는 그대로 보존되었습니다."); return
    # ── ② 신방식: 매니페스트 기반 온라인/폴더 업데이트 ──
    base=_C.get("업데이트기준URL") or _C.get("업데이트URL") or _UPDATE_BASE_DEFAULT
    localup=os.path.join(BASE,"_업데이트")
    if not base and os.path.isdir(localup) and os.path.exists(os.path.join(localup,"manifest.json")): base=localup
    if not base:
        print("ℹ 업데이트 위치가 아직 설정되지 않았습니다.")
        print("  설정(church_config.json)에  \"업데이트기준URL\": \"https://...\"  를 넣으면 버튼 한 번으로 최신 기능을 받습니다.")
        return
    try:
        man=_j.loads(_upd_fetch(base,"manifest.json").decode('utf-8'))
    except Exception as e:
        print(f"✗ 업데이트 정보를 가져오지 못했습니다: {str(e)[:100]}")
        print("  (인터넷 연결·업데이트 주소를 확인해 주세요. 자료는 안전합니다.)"); return
    remote_ver=str(man.get("version",""))
    if _ver_date(remote_ver)<=_ver_date(VERSION):
        print(f"✅ 이미 최신 버전입니다.  (현재 {VERSION})"); return
    print(f"🔔 새 버전: {remote_ver}   (현재 {VERSION})")
    if man.get("notes"): print(f"   ✨ 새 기능: {man['notes']}")
    backup(a)  # 자료 먼저 백업
    ROOT=os.path.dirname(BASE); ROOTN=os.path.normpath(ROOT)   # 설치 루트(_시스템의 상위) = manifest 경로 기준
    prot_files={"church_config.json","church_db.json"}
    prot_dirs={"_내자료","_백업","_업데이트백업",os.path.basename(PHOENIX)}
    stamp=datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    bdir=os.path.join(BASE,"_업데이트백업",stamp)
    ok=0; fail=[]
    for rel in man.get("files",[]):
        reln=str(rel).replace("\\","/").lstrip("/")
        segs=reln.split("/")
        if os.path.basename(reln) in prot_files or any(s in prot_dirs for s in segs):
            continue   # 자료·개인폴더(교인·설정·내자료·백업)는 절대 덮지 않음
        dst=os.path.normpath(os.path.join(ROOT, reln.replace("/",os.sep)))
        if not (dst==ROOTN or dst.startswith(ROOTN+os.sep)):
            fail.append(reln+"(경로벗어남)"); continue   # 설치 폴더 밖 쓰기 차단(안전)
        try:
            data=_upd_fetch(base, reln)
            if not data: fail.append(reln); continue
            if os.path.exists(dst):  # 이전본 백업
                bp=os.path.join(bdir, os.path.relpath(dst, ROOT)); os.makedirs(os.path.dirname(bp) or bdir,exist_ok=True); shutil.copy2(dst,bp)
            os.makedirs(os.path.dirname(dst) or ROOT, exist_ok=True)
            with open(dst,'wb') as f: f.write(data)
            ok+=1
        except Exception as e:
            fail.append(f"{reln}({str(e)[:30]})")
    print(f"✅ 업데이트 완료 — {ok}개 파일을 최신으로 교체했습니다. (이전본: _업데이트백업/{stamp})")
    if fail: print(f"   ⚠ 일부 실패 {len(fail)}건: {', '.join(fail[:6])}")
    print("  🔒 교인·심방·설교·재정·예화·주석·찬양곡 등 모든 자료는 그대로 보존되었습니다.")
    print("  🔄 프로그램을 껐다 다시 켜시면(시작.bat 더블클릭) 새 기능이 모두 적용됩니다.")
def update_check(a):
    """업데이트 있는지 확인만(적용 안 함) — 웹 상단 '새 업데이트' 표시용. __UPD__{json} 출력."""
    import json as _j
    base=_C.get("업데이트기준URL") or _C.get("업데이트URL") or _UPDATE_BASE_DEFAULT
    localup=os.path.join(BASE,"_업데이트")
    if not base and os.path.isdir(localup) and os.path.exists(os.path.join(localup,"manifest.json")): base=localup
    out={"new":False,"version":"","notes":""}
    if base:
        try:
            man=_j.loads(_upd_fetch(base,"manifest.json",timeout=8).decode('utf-8'))
            rv=str(man.get("version",""))
            if _ver_date(rv)>_ver_date(VERSION):
                out={"new":True,"version":rv,"notes":man.get("notes","")}
        except Exception: pass
    print("__UPD__"+_j.dumps(out,ensure_ascii=False))
def phoenix(a):
    """🔥 피닉스 복구 — 삭제·손상된 자료를 지난 시점으로 되살립니다. 자료는 결코 사라지지 않습니다."""
    import shutil
    os.makedirs(PHOENIX,exist_ok=True)
    snaps=sorted(f for f in os.listdir(PHOENIX) if f.startswith("db_") and f.endswith(".json"))
    if a.restore or a.last:
        if a.last or (a.restore or "").lower()=="last":
            target=snaps[-1] if snaps else None
        else:
            target=next((s for s in reversed(snaps) if a.restore in s), None)
        if not target: print("✗ 해당 복구시점이 없습니다. 'phoenix'로 목록을 확인하세요."); return
        if os.path.exists(DB): shutil.copy2(DB, DB+".before_restore")  # 현재 상태도 보존
        shutil.copy2(os.path.join(PHOENIX,target), DB)
        db=json.load(open(DB,encoding='utf-8'))
        print(f"🔥 피닉스 복구 완료 — {target[3:-5]} 시점으로 되살아났습니다.")
        print(f"   교인 {len(db.get('교인',[]))}명 · 설교 {len(db.get('설교',[]))}편 · 재정 {len(db.get('재정',[]))}건 복구")
        return
    print(f"🔥 피닉스 복구 지점 {len(snaps)}개 (자료가 지워져도 이 시점들로 되살릴 수 있습니다)")
    for s in list(reversed(snaps))[:15]:
        try:
            d=json.load(open(os.path.join(PHOENIX,s),encoding='utf-8'))
            print(f"  · {s[3:-5]}   교인 {len(d.get('교인',[]))}명 · 설교 {len(d.get('설교',[]))}편")
        except Exception: pass
    if not snaps: print("  (아직 스냅샷 없음 — 자료를 저장하면 자동으로 쌓입니다)")
    print("  복구:  phoenix --last  (가장 최근)   또는   phoenix --restore 20260717_093000")
def set_backup(a):
    """USB·D 백업 (설정 & 지금 저장) 통합 — 폴더를 넣으면 자동저장 폴더로 설정하고 즉시 전체 저장, 비우면 설정된 폴더로 지금 전체 저장."""
    import json as _j, types
    path=(getattr(a,'path','') or getattr(a,'to','') or "").strip()
    if path:   # 폴더 설정
        cfg=dict(_C); cfg["백업폴더"]=path
        _j.dump(cfg,open(CONFIG,'w',encoding='utf-8'),ensure_ascii=False,indent=2); _C["백업폴더"]=path
        try: os.makedirs(path,exist_ok=True)
        except Exception: pass
        print(f"✅ USB·D 자동저장 폴더 설정 완료: {path}")
        print("  이제부터 모든 자료가 저장될 때마다 이 폴더에도 자동 저장됩니다. 안전 이중보관!")
    elif not _C.get("백업폴더"):
        print("■ 자동저장 폴더가 아직 설정되지 않았습니다.")
        print("  저장 폴더 칸에  E:\\교회백업  같은 경로를 넣고 실행하세요 (USB=E:,  D폴더=D:\\교회백업).")
        print("  ※ 한 번 지정하면 이후 모든 자료가 그 폴더에도 자동 저장되고, 지금 즉시 전체 저장도 됩니다.")
        return
    # 설정된(또는 방금 설정한) 폴더로 즉시 전체 저장
    return sync_out(types.SimpleNamespace(to=(path or _C.get("백업폴더"))))
def sync_out(a):
    """USB·D드라이브로 전체 저장 — 자료(교인·설교 등)+만든 모든 문서를 지정 폴더에 복사(안전 이중보관)."""
    import shutil
    dest=a.to or _C.get("백업폴더")
    if not dest:
        print("사용법: sync --to \"D:\\교회백업\"  또는  \"E:\\\"(USB)")
        print("  또는 church_config.json 에 \"백업폴더\":\"E:\\\\교회백업\" 를 넣으면 설교 만들 때마다 자동 이중저장됩니다."); return
    os.makedirs(dest,exist_ok=True); n=0
    for f in ("church_db.json","church_config.json"):
        p=os.path.join(BASE,f)
        if os.path.exists(p):
            try: shutil.copy2(p,os.path.join(dest,f)); n+=1
            except Exception: pass
    nd=os.path.join(BASE,"_내자료")
    if os.path.isdir(nd):
        try: shutil.copytree(nd,os.path.join(dest,"_내자료"),dirs_exist_ok=True)
        except Exception: pass
    for dp,_,fs in os.walk(ROOT):
        if "_시스템" in dp: continue
        for f in fs:
            if f.lower().endswith((".docx",".hwp",".xlsx",".pptx",".txt")):
                rel=os.path.relpath(dp,ROOT); dd=os.path.join(dest,rel); os.makedirs(dd,exist_ok=True)
                try: shutil.copy2(os.path.join(dp,f),os.path.join(dd,f)); n+=1
                except Exception: pass
    print(f"✅ 외부 저장 완료: {dest}")
    print(f"   자료(교인·설교이력 등) + 만든 문서 {n}개 이중 보관 — USB·D드라이브 안전 백업")
def request(a):
    """기능 요청서 — 원하는 기능을 적으면 요청 파일 생성(담당자에게 전달→업데이트로 반영)."""
    import datetime
    if not a.text: print("사용법: request --text \"원하는 기능 설명\"  (예: 출석을 QR로 체크하고 싶어요)"); return
    rdir=os.path.join(BASE,"_요청"); os.makedirs(rdir,exist_ok=True)
    stamp=datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    p=os.path.join(rdir,f"기능요청_{sanit(CHURCH)}_{stamp}.txt")
    open(p,'w',encoding='utf-8').write(
        f"[기능 요청서]\n교회: {CHURCH}\n담임: {PASTOR}\n현재버전: v{VERSION}\n일시: {stamp}\n\n■ 요청 내용\n{a.text}\n\n"
        "※ 이 파일을 배포해 주신 담당자에게 카톡/메일로 보내주세요.\n"
        "  담당자가 이 기능을 만들어 업데이트로 보내드립니다.\n"
        "  업데이트해도 교인·심방·설교 등 모든 자료는 영구 보존됩니다.\n")
    print(f"✅ 기능 요청서 생성: {p}")
    print("  → 이 파일을 담당자에게 보내면, 다음 업데이트(🔔)에 원하는 기능이 추가됩니다.")
def dashboard(a):
    """대시보드 요약(JSON) — 웹 상단 통계용: 교인수·이번주 심방·이번달 헌금·다가오는 생일·돌봄 필요."""
    import datetime as _dt, json as _j
    db=load(); td=_dt.date.today(); mem=db.get("교인",[])
    def _d(s):
        try: return _dt.date.fromisoformat(str(s)[:10])
        except Exception: return None
    def _amt(x):
        try: return int(x)
        except Exception:
            try: return int(re.sub(r'[^0-9-]','',str(x)) or 0)
            except Exception: return 0
    n_mem=len([m for m in mem if (m.get("상태") or "재적") in ("재적","재적중","")]) or len(mem)
    wk0=td-_dt.timedelta(days=td.weekday())
    n_visit=sum(1 for m in mem for v in m.get("심방이력",[]) if (lambda d: d and wk0<=d<=wk0+_dt.timedelta(days=6))(_d(v.get("날짜",""))))
    ym=td.isoformat()[:7]
    offering=sum(_amt(r.get("금액",0)) for r in db.get("재정",[]) if r.get("구분")=="수입" and str(r.get("날짜",""))[:7]==ym)
    def nextd(mmdd):
        try:
            mo,da=int(mmdd[:2]),int(mmdd[3:5]); d=_dt.date(td.year,mo,da)
            return d if d>=td else _dt.date(td.year+1,mo,da)
        except Exception: return None
    bdays=[]
    for m in mem:
        v=(m.get("생년월일") or "")[5:]
        nd=nextd(v) if v else None
        if nd and 0<=(nd-td).days<=7: bdays.append(((nd-td).days,m.get("이름",""),m.get("직분","") or ""))
    bdays.sort()
    care_names=[]
    for m in mem:
        H=m.get("심방이력",[]); last=_d(H[-1].get("날짜")) if H else None
        if last is None or (td-last).days>=90: care_names.append(m.get("이름",""))
    newf=0
    for m in mem:
        rd=_d(m.get("등록일",""))
        if m.get("상태")=="새가족" or (rd and 0<=(td-rd).days<=90): newf+=1
    wk=["월","화","수","목","금","토","일"][td.weekday()]
    out={"교인":n_mem,"이번주심방":n_visit,"이번달헌금":offering,"돌봄필요":len(care_names),
         "돌봄명단":care_names[:4],"새가족":newf,"연월":ym,"오늘":td.isoformat(),"요일":wk,
         "생일수":len(bdays),"오늘생일":[nm for d,nm,jk in bdays if d==0],
         "생일":[{"d":d,"이름":nm,"직분":jk} for d,nm,jk in bdays[:6]]}
    print("__DASH__"+_j.dumps(out,ensure_ascii=False))
def _congrats_pool():
    import json as _j
    try: return _j.load(open(os.path.join(BASE,"_data","축하위로문구.json"),encoding='utf-8'))
    except Exception: return {}
def congrats(a):
    """🎉 상황별 축하·위로 문자 — 군입대·전역·진급·취업·합격·졸업·입학·출산·백일돌·개업·이사·퇴원·수술·은퇴·유학·수상·임직.
       상황과 성도 이름을 넣으면 골라 쓸 수 있는 카톡 문구 3가지를 만들어 드립니다(무료 붙여넣기)."""
    pool=_congrats_pool()
    if not pool: print("상황별 문구 자료(_data/축하위로문구.json)가 없습니다."); return
    keys=list(pool.keys())
    occ=(getattr(a,'occasion','') or getattr(a,'kind','') or "").strip()
    name=(getattr(a,'name','') or "").strip()
    jik=(getattr(a,'jik','') or getattr(a,'role','') or "").strip()
    if not occ:
        print("■ 쓸 수 있는 상황: "+" · ".join(keys)); print("  예) congrats --occasion 군입대 --name 홍길동 --jik 형제"); return
    key=next((k for k in keys if occ==k), None) or next((k for k in keys if occ in k or k.startswith(occ)), None)
    if not key:
        print(f"‘{occ}’에 맞는 상황을 못 찾았습니다. 쓸 수 있는 상황: "+" · ".join(keys)); return
    if not name:
        print("성도 이름을 입력해 주세요. 예) --name 홍길동 --jik 형제"); return
    msgs=pool[key]
    start=sum(ord(c) for c in name)%len(msgs)   # 이름별 시작점 → 늘 같은 문구만 나오지 않게
    picks=[msgs[(start+i)%len(msgs)] for i in range(min(3,len(msgs)))]
    L=[f"════ {key} 축하·위로 · 카톡에 붙여넣기 ({name} {jik}) ════",
       "※ 아래 【1】【2】【3】 중 마음에 드는 문구를 하나 골라 복사해서 카톡으로 보내세요(무료).",""]
    for i,it in enumerate(picks,1):
        txt=it['문구'].format(church=CHURCH,name=name,jik=jik,pastor=PASTOR)
        L.append(f"【{i}】 (성구 {it.get('성구','')})"); L.append(txt); L.append("")
    out=os.path.join(CAT("07"),f"[상황축하_{key}] {sanit(name)}_{today()}.txt")
    open(out,'w',encoding='utf-8').write("\n".join(L)); _openfile(out)
    print(f"✅ {key} 축하·위로 문구 생성: {out}  (【1】【2】【3】 중 골라 카톡 붙여넣기)")
    for i,it in enumerate(picks,1):
        print(f"  【{i}】 {it['문구'].format(church=CHURCH,name=name,jik=jik,pastor=PASTOR)[:56]}…")
def harvest_plan(a):
    """🌾 새생명축제(태신자 운동) 준비 타임라인 — 축제일 기준 D-56~D+28 월·일별 준비사항(날짜 자동계산·체크박스)."""
    import datetime as _dt
    d=newdoc()
    para(d,f"{a.title or '새생명 초청잔치'} 준비 타임라인",19,NAVY,True,AL.CENTER,after=2,font=SANS)
    dd=None
    try: dd=_dt.date.fromisoformat(a.dday) if a.dday else None
    except Exception: dd=None
    para(d,(f"축제일 {a.dday}  ·  " if a.dday else "")+f"{CHURCH} · 주제 요삼 1:2 · {today()}",10,GRAY,AL.CENTER,after=6,font=SANS); hr(d,"2F7D5B",12)
    stages=[(-56,"8주 전 · 태신자 작정·중보 시작",["태신자 명단 작성(성도별 1~3명 작정)","태신자 이름 정해 중보기도 시작","전 교회 축제 비전 선포(발대식)","준비위원회 구성 착수"]),
      (-49,"7주 전 · 조직·형식 결정",["준비위 조직표 확정(홍보·양육·초청·행사·다과·새가족등록)","1부/2부 운영 여부·시간 결정(예: 오전 11:30·오후 2:00)","예산안 편성"]),
      (-42,"6주 전 · 양육·초청 전략",["태신자 만남·양육팀 배치","초청 전략 수립(관계전도·초청장·SNS)","21일 중보기도표 배부"]),
      (-35,"5주 전 · 프로그램 구성",["예배·초청 순서 구성","문화행사·레크·특송 기획","초청 대상별 접촉 시작"]),
      (-28,"4주 전 · 홍보물 제작",["초청장·초청카드 제작(초청받은분+초청자 양면)","포스터·현수막·SNS 홍보물","선물·다과 품목 결정"]),
      (-21,"3주 전 · 초청 집중",["초청장 배부·개별 초청 시작","중보기도 집중(21일 작정)","순서자·봉사자 1차 배정"]),
      (-14,"2주 전 · 초청 확정·리허설",["참석 예정 태신자 확정·집계","프로그램 리허설","역할 분담표 확정(안내·주차·환영·상담·새가족등록)","무대·동선·좌석(새가족석·태신자석) 배치"]),
      (-7,"1주 전 · 최종 점검",["초청 최종 확인(재초청 문자)","다과·선물·물품 준비","당일 순서표(분단위) 확정","결신 초청·상담 흐름 점검"]),
      (-1,"전날 · 준비 완료",["장소 세팅·리허설 마무리","봉사자 집결·역할 최종 확인","새가족 등록카드·결신카드 비치"]),
      (0,"당일 · 새생명 초청잔치",["봉사자 조기 집결·기도","안내·주차·환영·다과 가동","예배·초청·결신 초청","결신카드·새가족 등록","기념촬영·배웅"]),
      (1,"익일 · 결신자 정리",["결신자·방문자 명단 정리","초청자별 태신자 결과 집계","감사 문자 발송"]),
      (8,"후속 4주 · 정착",["새가족 심방·환영회","새신자 교육 연결(4주)","초청자가 정착까지 동행","한 달 뒤 정착률 점검"])]
    for off,title,items in stages:
        ds=""
        if dd: ds="  ("+(dd+_dt.timedelta(days=off)).isoformat()+")"
        head=("D-DAY" if off==0 else (f"D{off}" if off<0 else f"D+{off}"))+ds
        para(d,f"{head} · {title}",13,RGBColor(0x2F,0x7D,0x5B),True,before=8,after=3,font=SANS)
        for it in items: para(d,f"  ☐  {it}",11,after=2)
    para(d,"※ 초청카드에 '초청자(이름·셀)'를 함께 적어, 초청자별로 태신자를 추적하고 정착까지 책임지게 하는 것이 태신자 운동의 핵심입니다.",9.5,LG,before=10,font=SANS)
    out=os.path.join(CAT("04"),f"[새생명축제 준비타임라인] {sanit(a.title or '새생명초청잔치')}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 새생명축제 준비 타임라인: {out}"+(f" (축제일 {a.dday} 기준 날짜 계산)" if dd else " (--dday 축제일 넣으면 날짜 자동계산)"))
def harvest_checklist(a):
    """🌾 새생명축제 당일 진행 체크리스트 — 시간대·역할별 당일 운영 점검표."""
    d=newdoc()
    para(d,f"{CHURCH} · {a.title or '새생명 초청잔치'}",11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,"당일 진행 체크리스트",20,NAVY,True,AL.CENTER,before=4,after=2,font=SANS); hr(d,"2F7D5B",12)
    para(d,f"{a.date or today()}",10,GRAY,align=AL.CENTER,after=8,font=SANS)
    blocks=[("행사 2시간 전 · 세팅",["강단·좌석(새가족석·태신자석) 배치","음향·영상·PPT 점검","안내·주차·환영 배치","다과·선물·등록대 세팅","결신카드·새가족 등록카드 비치"]),
      ("행사 30분 전 · 맞이",["봉사자 기도회","주차 안내 시작","환영·명찰·안내지 배부","태신자·초청자 좌석 안내"]),
      ("예배·초청 진행",["사회·순서 진행","특송·문화순서","말씀·초청 메시지","결신 초청·손들기·상담 연결","결신카드 작성 안내"]),
      ("마무리·배웅",["다과·교제","새가족 등록 유도","기념촬영","배웅·귀가 안내","초청자에게 후속 연결 부탁"]),
      ("행사 후 정리",["결신카드·등록카드 회수·집계","초청자별 결과 정리","봉사자 격려·정리","익일 감사문자 준비"])]
    for t,items in blocks:
        para(d,f"■ {t}",13,MAROON,True,before=8,after=3,font=SANS)
        for it in items: para(d,f"  ☐  {it}",11,after=2)
    out=os.path.join(CAT("04"),f"[새생명축제 당일체크리스트] {sanit(a.title or '새생명초청잔치')}_{a.date or today()}.docx"); _savedoc(d,out)
    print(f"✅ 새생명축제 당일 체크리스트: {out}")
def cell_worship(a):
    """🏠 구역·속회 예배 순서지 — 찬송·기도·말씀 나눔·삶 나눔·중보 (인도용 교안 겸용). 구역장 출력용."""
    d=newdoc()
    para(d,f"{CHURCH} · 구역·속회 예배",11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,a.title or "구역·속회 예배 순서",20,NAVY,True,AL.CENTER,before=4,after=2,font=SANS); hr(d,"1E6E6A",12)
    para(d,f"{a.date or today()}  ·  인도 {a.leader or '구역장'}  ·  본문 {a.text or ''}",10,GRAY,AL.CENTER,after=8,font=SANS)
    steps=[("묵도·부름","다 함께 마음을 모아 주님 앞에 나아갑니다"),("찬송","____장 (다 함께)"),("대표 기도","____ (한 분)"),
           ("성경 봉독",a.text or "____________ (다 함께)"),("말씀 나눔","오늘 본문에서 받은 은혜를 나눕니다"),
           ("삶 나눔·교제","한 주간의 삶·감사 제목을 나눕니다"),("중보 기도","나눈 기도제목으로 서로 위해 기도"),
           ("광고·교제","구역 소식·다음 모임 안내 · 다과"),("주기도문","다 함께")]
    for i,(t,desc) in enumerate(steps,1):
        para(d,f"{i}. {t}",13,NAVY,True,before=6,after=1,font=SANS); para(d,f"   {desc}",10.5,LG,after=2)
    para(d,"■ 말씀 나눔 질문",12,MAROON,True,before=10,after=3,font=SANS)
    for q in ["이 말씀에서 하나님은 어떤 분으로 나타나시나요?","내 삶에 적용할 한 가지는 무엇인가요?","이번 주 함께 기도·실천할 제목은?"]:
        para(d,f"· {q}",10.5,after=2)
    para(d,"■ 우리 구역 기도제목",12,MAROON,True,before=8,after=3,font=SANS)
    para(d,"________________________________________________________",10.5,LG,after=2)
    para(d,"________________________________________________________",10.5,LG,after=2)
    out=os.path.join(CAT("04"),f"[구역예배] {sanit(a.title or '순서')}_{a.date or today()}.docx"); _savedoc(d,out)
    print(f"✅ 구역·속회 예배 순서지: {out}")
def member_transfer(a):
    """↔ 교적 이동(전입·전출·이명) — 이름을 넣으면 기록+상태 갱신, 그리고 항상 '교적 이동 대장'을 함께 출력. (이름 비우면 대장만)"""
    db=load(); db.setdefault("교적이동",[])
    name=(getattr(a,'name','') or "").strip()
    if name:
        kind=(getattr(a,'kind','') or "전출").strip()
        rec={"종류":kind,"이름":name,"날짜":a.date or today(),"상대교회":getattr(a,'church','') or "","사유":a.memo or "","등록일":today()}
        db["교적이동"].append(rec)
        hit=find(db,name)
        if hit:
            if kind in ("전출","이명"): hit[0]["상태"]="이명"
            elif kind=="전입": hit[0]["상태"]="재적"
        save(db)
        print(f"✅ 교적 이동 기록: {kind} · {name} · {rec['날짜']}"+(f" ({rec['상대교회']})" if rec['상대교회'] else ""))
        if kind in ("전출","이명"): print("   → 이명증서가 필요하면 '증명서 발급'에서 종류를 '이명'으로 발급하세요.")
    else:
        print("ℹ 이름을 비우고 실행 → 교적 이동 대장만 출력합니다.")
    transfer_ledger(a)   # 기록 후 항상 전체 대장(전입·전출·이명 내역)을 함께 만들어 드림
def transfer_ledger(a):
    """교적 이동 대장 — 전입·전출·이명 기록을 표로 출력."""
    db=load(); recs=db.get("교적이동",[])
    if not recs: print("교적 이동 기록이 없습니다."); return
    d=newdoc()
    para(d,f"{CHURCH}",11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,"교적 이동 대장",20,NAVY,True,AL.CENTER,before=2,after=2,font=SANS); hr(d,"1B2A4A",12)
    para(d,f"발행일 {today()} · 총 {len(recs)}건",10,GRAY,AL.CENTER,after=8,font=SANS)
    rows=[["날짜","종류","이름","상대 교회","사유"]]
    for r in sorted(recs,key=lambda x:str(x.get("날짜",""))):
        rows.append([r.get("날짜",""),r.get("종류",""),r.get("이름",""),r.get("상대교회",""),r.get("사유","")])
    _ref_table(d,rows,colored=True)
    out=os.path.join(CAT("01"),f"[교적이동대장]_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 교적 이동 대장: {out} ({len(recs)}건)")
def annual_plan(a):
    """📅 연간 사역 계획표 — 월별 절기·주요 사역 한 장(우리 교회 계획 기입란). 절기는 그 해 교단력 확인."""
    year=(getattr(a,'year','') or str(datetime.date.today().year)).strip()
    d=newdoc()
    para(d,f"{CHURCH}",11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,f"{year}년 연간 사역 계획",20,NAVY,True,AL.CENTER,before=4,after=2,font=SANS); hr(d,"1B2A4A",12)
    para(d,"※ 부활절·맥추·추수감사 등 변동 절기는 그 해 달력·교단력을 확인해 채우세요.",9,LG,align=AL.CENTER,after=8,font=SANS)
    seasons={1:"신년감사·신년특별새벽기도회",2:"설 명절·전교인 기도",3:"사순절 시작·상반기 사역 출발",4:"고난주간·부활절",
             5:"가정의 달(어린이·어버이·스승주일)",6:"맥추감사절·상반기 결산",7:"여름 수련회·단기선교",8:"여름성경학교·비전트립",
             9:"가을 전도·새생명축제 준비",10:"추수감사절·종교개혁주일",11:"추수감사·임직·성탄 준비(대강절)",12:"성탄절·송구영신"}
    rows=[["월","주요 절기·사역","우리 교회 계획 (직접 기입)"]]
    for m in range(1,13): rows.append([f"{m}월",seasons.get(m,""),""])
    _ref_table(d,rows,colored=True)
    para(d,"■ 연간 반복 사역: 주일예배·수요예배·금요기도회·새벽기도 / 심방·구역모임·제직회·당회·공동의회",10,GRAY,before=8,font=SANS)
    out=os.path.join(CAT("01"),f"[연간사역계획] {sanit(year)}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 연간 사역 계획표: {out}")
_SAC_KINDS=("유아세례","세례","학습","입교")
def sacrament_add(a):
    """💧 성례 등록 & 대장 — 이름을 넣으면 유아세례·세례·학습·입교를 성례대장에 기록(증명서 근거), 이름을 비우면 성례대장을 본다."""
    name=(getattr(a,'name','') or '').strip()
    if not name:
        return sacrament_ledger(a)   # 이름 비우면 대장 보기(통합)
    db=load(); db.setdefault("성례",[])
    kind=(getattr(a,'kind','') or "세례").strip()
    if kind not in _SAC_KINDS:
        print(f"종류는 {'·'.join(_SAC_KINDS)} 중 하나로 입력해 주세요."); return
    seq=sum(1 for r in db["성례"] if r.get("종류")==kind)+1
    rec={"번호":seq,"종류":kind,"이름":a.name or "","생년월일":a.birth or "","주소":a.addr or "",
         "날짜":a.date or today(),"집례자":a.by or PASTOR,"본문":a.text or "","비고":a.memo or "","등록일":today()}
    db["성례"].append(rec)
    hit=find(db,a.name) if a.name else []
    if hit and kind in ("유아세례","세례"):
        hit[0]["세례"]=kind; hit[0]["세례일"]=rec["날짜"]
    save(db)
    print(f"✅ 성례 등록: {kind} 제{seq}호 · {a.name} · {rec['날짜']} (집례 {rec['집례자']})")
    print("   → 성례대장에 기록됐습니다. 나중에 '증명서 발급(세례)'으로 세례 증명서를 만들 수 있습니다.")
def sacrament_ledger(a):
    """📖 성례대장 — 유아세례·세례·학습·입교 기록을 대장(표)으로 출력. 종류 필터 가능. 진학용 세례증명의 근거."""
    db=load(); recs=db.get("성례",[])
    kf=(getattr(a,'kind','') or "").strip()
    if kf: recs=[r for r in recs if r.get("종류")==kf]
    if not recs:
        print("성례 기록이 없습니다. 먼저 '성례 등록'으로 세례·학습·입교를 기록해 주세요."); return
    d=newdoc()
    para(d,f"{CHURCH}",11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,f"성례대장{(' · '+kf) if kf else ''}",20,NAVY,True,AL.CENTER,before=2,after=2,font=SANS); hr(d,"1B2A4A",12)
    para(d,f"발행일 {today()} · 총 {len(recs)}건",10,GRAY,AL.CENTER,after=8,font=SANS)
    rows=[["번호","종류","이름","생년월일","성례일","집례자","비고"]]
    for r in sorted(recs,key=lambda x:(x.get("종류",""),str(x.get("날짜","")))):
        rows.append([str(r.get("번호","")),r.get("종류",""),r.get("이름",""),r.get("생년월일",""),
                     r.get("날짜",""),r.get("집례자",""),r.get("비고","")])
    _ref_table(d,rows,colored=True)
    para(d,"※ 이 대장은 세례·학습·입교의 공식 기록이며, 세례교인 증명서 발급의 근거가 됩니다. 원본은 교회가 영구 보존합니다.",9,LG,before=8,font=SANS)
    out=os.path.join(CAT("02"),f"[성례대장]{('_'+sanit(kf)) if kf else ''}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ 성례대장 생성: {out} ({len(recs)}건)")
def sacrament_apply(a):
    """📝 성례 서식 — '무엇을'이 증명서면 세례증명서 발급, 아니면 신청서 서식(당회 제출·심의용)."""
    if '증명' in (getattr(a,'mode','') or '').strip():   # 세례교인 증명서 발급으로 라우팅
        import types
        return cert(types.SimpleNamespace(name=(getattr(a,'name','') or '').strip(),kind='세례',
                    purpose=getattr(a,'purpose','') or '',role='',term='',no='',date=''))
    kind=(getattr(a,'kind','') or "세례").strip()
    d=newdoc()
    para(d,f"{CHURCH}",11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,f"{kind} 신청서",22,NAVY,True,AL.CENTER,before=4,after=6,font=SANS); hr(d,"1B2A4A",12)
    if kind=="유아세례":
        fields=[("아기 이름",a.name),("생년월일",a.birth),("아버지",""),("어머니",""),("주소",a.addr),("연락처",a.tel),("소속 셀·부서",a.cell)]
    else:
        fields=[("성명",a.name),("생년월일",a.birth),("주소",a.addr),("연락처",a.tel),("소속(셀·부서)",a.cell),
                ("신앙 경력",""),("현재 신앙 상태",""),("추천 인도자",a.leader)]
    rows=[[k,(v if v else "____________________________")] for k,v in fields]
    _ref_table(d,rows,colored=False)
    para(d,f"위와 같이 {kind}을(를) 신청합니다.",11.5,NAVY,True,AL.CENTER,before=14,after=10,font=SANS)
    para(d,today(),10,GRAY,align=AL.CENTER,after=4,font=SANS)
    para(d,"신청인 : ______________________ (서명)",11,align=AL.CENTER,after=2)
    para(d,"※ 본 신청서는 당회 심의를 거쳐 성례대장에 등재됩니다.",9,LG,before=12,font=SANS)
    out=os.path.join(CAT("02"),f"[{sanit(kind)}신청서] {sanit(a.name or '양식')}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ {kind} 신청서 생성: {out}")
def event_add(a):
    """📅 일정 등록 — 날짜·제목·내용·중요표시. 캘린더와 대시보드에 표시됩니다."""
    db=load(); db.setdefault("일정",[])
    dt=(getattr(a,'date','') or "").strip(); ti=(getattr(a,'title','') or "").strip()
    if not (dt and ti):
        print("날짜와 제목을 입력해 주세요. 예) --date 2026-08-15 --title 맥추감사주일"); return
    imp=1 if str(getattr(a,'important','') or '').strip() in ('1','y','예','중요','o','O') else 0
    rec={"id":db.get("_evseq",0)+1,"날짜":dt,"제목":ti,"내용":getattr(a,'note','') or "","중요":imp}
    db["_evseq"]=rec["id"]; db["일정"].append(rec); save(db)
    print(f"✅ 일정 등록: {dt} · {ti}"+(" ★중요" if imp else ""))
def events_json(a):
    """일정 목록(JSON) — 월(YYYY-MM) 필터 가능. 캘린더 표시용. __EVT__ 출력."""
    import json as _j
    db=load(); evs=db.get("일정",[])
    ym=(getattr(a,'month','') or "").strip()
    if ym: evs=[e for e in evs if str(e.get("날짜","")).startswith(ym)]
    print("__EVT__"+_j.dumps(sorted(evs,key=lambda e:str(e.get("날짜",""))),ensure_ascii=False))
def event_del(a):
    """일정 삭제 — id로."""
    db=load(); eid=str(getattr(a,'id','') or "")
    n=len(db.get("일정",[])); db["일정"]=[e for e in db.get("일정",[]) if str(e.get("id"))!=eid]
    save(db); print(f"✅ 일정 삭제 완료" if len(db["일정"])<n else "삭제할 일정을 못 찾았습니다.")
def cal_print(a):
    """🗓️ 월간 일정 달력(탁상 달력형) 인쇄 — 그 달 일정을 날짜칸에 넣은 표(docx)."""
    import datetime as _dt, calendar as _cal
    ym=(getattr(a,'month','') or _dt.date.today().isoformat()[:7]).strip()
    try: y,mo=int(ym[:4]),int(ym[5:7])
    except Exception: y,mo=_dt.date.today().year,_dt.date.today().month
    db=load(); byday={}
    for e in db.get("일정",[]):
        s=str(e.get("날짜",""))
        if s.startswith(f"{y:04d}-{mo:02d}"):
            try: byday.setdefault(int(s[8:10]),[]).append(e)
            except Exception: pass
    d=newdoc()
    para(d,f"{CHURCH}",11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,f"{y}년 {mo}월 일정 달력",20,NAVY,True,AL.CENTER,before=4,after=6,font=SANS); hr(d,"1B2A4A",12)
    weeks=_cal.Calendar(firstweekday=6).monthdayscalendar(y,mo)   # 일요일 시작
    t=d.add_table(rows=1+len(weeks),cols=7)
    for j,h in enumerate(["일","월","화","수","목","금","토"]):
        c=t.rows[0].cells[j]; ctext(c,h,WHITE,True,AL.CENTER); shade(c,"1B2A4A")
    for i,wk in enumerate(weeks,1):
        for j,day in enumerate(wk):
            c=t.rows[i].cells[j]
            if day==0: continue
            lines=str(day)
            for e in byday.get(day,[]): lines+="\n"+("★" if e.get("중요") else "·")+str(e.get("제목",""))[:11]
            ctext(c,lines,NAVY if byday.get(day) else GRAY,False,AL.LEFT)
    tborders(t,"BBBBBB",4)
    out=os.path.join(CAT("01"),f"[일정달력] {sanit(ym)}_{today()}.docx"); _savedoc(d,out)
    print(f"✅ {y}년 {mo}월 일정 달력: {out}")
def church_setup(a):
    """⛪ 우리 교회 이름 설정 — 교회명·담임명을 입력하면 설정파일에 자동 저장(직접 편집 불필요).
       이후 모든 문서·주보·증명서·축하 문자에 이 이름이 자동으로 들어갑니다."""
    import json as _j
    p=os.path.join(BASE,"church_config.json")
    try: c=_j.load(open(p,encoding='utf-8')) if os.path.exists(p) else {}
    except Exception: c={}
    ch=(getattr(a,'church','') or "").strip(); pa=(getattr(a,'pastor','') or "").strip()
    if not ch and not pa:
        print("ℹ 우리 교회 이름과 담임 목사님 성함을 입력해 주세요.")
        print("   예) 교회명: 은혜교회 · 담임: 홍길동 목사")
        print(f"   현재 설정 — 교회명: {c.get('교회명','(미설정)')} · 담임: {c.get('담임','(미설정)')}")
        return
    if ch: c["교회명"]=ch
    if pa: c["담임"]=pa
    _j.dump(c,open(p,'w',encoding='utf-8'),ensure_ascii=False,indent=2)
    print(f"✅ 설정 완료 — 교회명: {c.get('교회명','')} · 담임: {c.get('담임','')}")
    print("   이제부터 만드는 모든 문서·주보·증명서·축하 문자에 이 이름이 자동으로 들어갑니다.")
    print("   🔄 화면 위쪽 이름까지 바꾸려면 프로그램을 껐다 다시 켜 주세요(시작.bat 더블클릭).")
def manual(a):
    """사용 설명서 — 모든 기능 사용법을 담은 친절한 매뉴얼(docx) 생성. 처음 쓰는 분께."""
    d=newdoc()
    para(d,"",after=44); hr(d,"A67C1E",12)
    para(d,CHURCH,13,GOLD,True,AL.CENTER,after=2,font=SANS)
    para(d,"교회 종합행정 프로그램",26,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,"사용 설명서",18,MAROON,True,AL.CENTER,after=6,font=SANS); hr(d,"A67C1E",12)
    para(d,f"v{VERSION} · {CHURCH} 교회 종합행정 · 무료",10,LG,AL.CENTER,before=8,font=SANS)
    d.add_page_break()
    secs=[
     ("시작하기",[
      "‘★ 교회행정 시작 (여기를 더블클릭).bat’ 을 더블클릭하면 브라우저에 예쁜 화면이 열립니다 — 카드를 눌러 사용하세요.",
      "‘(참고) 번호 메뉴로 시작.bat’ 은 번호를 눌러 쓰는 간단 메뉴입니다(컴퓨터가 익숙지 않으셔도 OK).",
      "처음 한 번: _시스템\\church_config.json 을 메모장으로 열어 우리 교회명·담임명을 넣고 저장.",
      "준비물: 파이썬과 필요한 라이브러리(python-docx·openpyxl·python-pptx)는 배포판에 이미 포함되어 있어 따로 설치하지 않으셔도 됩니다. (한글 .hwp 변환만 한컴오피스가 설치된 PC에서 가능)",
      "만든 문서는 01~09 번호 폴더에 예배유형·종류별로 자동 정리됩니다."]),
     ("★ 맨 처음 하실 일 — 우리 교회 이름 넣기 (아주 쉽습니다)",[
      "이 프로그램은 ‘우리 교회 이름’과 ‘담임 목사님 성함’을 한 번만 넣어두면, 이후 만드는 모든 문서·주보·증명서·축하 문자에 그 이름이 자동으로 들어갑니다.",
      "◆ 가장 쉬운 방법 — 이대로만 하세요:",
      "  1) 프로그램을 켭니다(‘★ 교회행정 시작’ 파일 더블클릭).",
      "  2) 왼쪽 메뉴에서 ‘시스템’을 누르고, 맨 위의 ‘⛪ 우리 교회 이름 설정’ 카드를 누릅니다.",
      "  3) ‘우리 교회 이름’ 칸에 교회 이름을(예: 은혜교회), ‘담임 목사님 성함’ 칸에 성함을(예: 홍길동 목사) 적고 아래 ‘실행’ 버튼을 누릅니다.",
      "  4) 끝입니다! 프로그램을 껐다 다시 켜면 화면과 모든 문서에 우리 교회 이름이 나옵니다.",
      "※ 위쪽 검색창에 ‘교회 이름’이라고 쳐도 이 카드를 바로 찾을 수 있습니다.",
      "※ 어려운 설정 파일을 직접 여실 필요가 전혀 없습니다. 이 카드가 알아서 넣어드립니다.",
      "※ (컴퓨터에 익숙하신 분만) _시스템\\church_config.json 을 메모장으로 열어 \"교회명\"·\"담임\" 값을 직접 바꾸셔도 되지만, 위 카드를 쓰시는 것이 가장 안전하고 쉽습니다."]),
     ("🎂 생일·결혼기념일 축하 문자 — 매년 다른 문구로",[
      "성도 등록 때 ‘생년월일’과 ‘결혼기념일’을 넣어두면, 그 날짜가 다가올 때 축하 문자를 자동으로 만들어 드립니다.",
      "쓰는 법: 목양 화면 맨 위 ‘🎂 생일·결혼기념일 축하’ 카드를 누르면 → 다가오는 생일·결혼기념일 성도(‘오늘’은 따로 표시)와 함께, 당일 보낼 축하 카톡 문구가 만들어집니다 → 각 사람 글을 복사해 카톡으로 보내시면 됩니다(무료).",
      "‘며칠 이내를 볼까요?’ 칸에 숫자를 넣으면 그 기간의 생일·기념일을 미리 봅니다(0을 넣으면 오늘만). 한 카드로 ‘미리 알림’과 ‘축하 문구’를 다 챙깁니다.",
      "★100가지 문구 자동 순환: 생일·결혼기념일 축하 문구가 각각 100가지 준비돼 있어, 같은 성도라도 매년 다른 문구가 자동으로 선택됩니다. 해가 바뀌어도 작년과 똑같은 인사가 가지 않습니다.",
      "★지난 축하 보기: 그 성도에게 지난해·재작년에 어떤 문구를 보냈는지 함께 보여드려, 겹치지 않게 새 축하를 보내실 수 있습니다.",
      "성구는 주소만(예: 애 3:22-23) 넣고 축복의 내용은 새로 쓴 글이라, 저작권 걱정 없이 그대로 보내셔도 됩니다.",
      "생일 ‘당일 아침’에 보내면 성도들이 특히 감격합니다. 문구의 이름·호칭을 살짝 다듬어 더 개인적으로 보내셔도 좋습니다."]),
     ("🔄 프로그램 업데이트 — 새 기능 받기 (자료는 그대로 보존)",[
      "이 프로그램은 계속 발전합니다. 새 기능이 나오면 아래 방법으로 받으실 수 있고, 어느 쪽이든 교인·심방·재정 등 입력하신 자료는 그대로 보존됩니다.",
      "◆ 방법 1 (가장 쉬움): 화면 오른쪽 위 ‘업데이트’ 버튼을 누르세요. 인터넷에서 최신 기능을 자동으로 받아옵니다.",
      "◆ 방법 2 (업데이트 ZIP을 받으셨으면):",
      "  1) 받으신 ‘★업데이트_v최신.zip’ 파일을 압축 풉니다(오른쪽 클릭 → 압축 풀기).",
      "  2) 풀린 파일들을 전체 선택·복사해서, 기존 프로그램 폴더에 붙여넣기 합니다.",
      "  3) ‘같은 파일이 있습니다’ 라고 물으면 → ‘파일 덮어쓰기(바꾸기)’를 선택합니다.  → 끝!",
      "★자료는 안 지워집니다: 교인·재정 자료(church_db.json)와 교회 이름(church_config.json)은 업데이트에 들어있지 않아, 덮어써도 그대로 유지됩니다. 몇 번을 덮어써도 안전합니다.",
      "※ ‘기존 폴더를 삭제하고 새로 설치’는 절대 하지 마세요 — 입력하신 자료가 사라집니다. 삭제하지 말고 ‘덮어쓰기’만 하세요.",
      "※ 만약을 위해 업데이트 전에 _시스템 폴더의 church_db.json 을 다른 곳에 한 번 복사해두시면 100% 안심입니다. 혹시 자료가 잘못돼도 ‘피닉스 복구’로 지난 시점으로 되살릴 수 있습니다."]),
     ("🎨 화면 꾸미기 — 테마·색상·배경 (오른쪽 위 ☀ 단추)",[
      "화면 오른쪽 위 ‘☀’(해 모양) 단추를 누르면 화면을 취향대로 꾸미는 창이 열립니다.",
      "밝게/어둡게: 낮에는 밝은 화면, 밤에는 눈이 편한 어두운 화면으로 바꿀 수 있습니다.",
      "색상: 파랑·장미·보라·주황·민트 등 강조 색을 고를 수 있습니다.",
      "계절 배경: 봄·여름·가을·겨울의 은은한 배경을 고르면 화면 분위기가 그 계절로 바뀝니다.",
      "★내 사진 배경: 가족·자녀 사진, 교회 전경, 좋아하시는 풍경을 화면 배경으로 쓸 수 있습니다.",
      "  가장 쉬운 법: ☀ 창의 ‘내 사진 배경’ 아래 ‘＋ 사진 추가’ 버튼을 눌러 사진을 고르면, 바로 올라가 배경으로 적용됩니다(폴더를 열 필요가 없습니다).",
      "  또는: _시스템\\_내자료\\배경 폴더에 사진(jpg·png)을 직접 넣어도 ‘내 사진 배경’에 나타납니다 → 누르면 적용, ✕를 누르면 원래대로.",
      "  글자가 잘 보이도록 사진 위에 은은한 막을 자동으로 씌워 드리니, 밝은 사진·어두운 사진 모두 괜찮습니다. (가로로 넓은 사진이 예쁘게 채워집니다.)",
      "고른 화면 설정은 이 컴퓨터에 저장되어, 프로그램을 껐다 켜도 그대로 유지됩니다."]),
     ("🎬 행사 홍보영상 만들기 — 영상이 처음이어도 괜찮습니다",[
      "‘영상·홍보’ 섹션은 영상 편집을 한 번도 안 해보신 목사님도 행사 홍보영상을 만들 수 있게 돕습니다. 준비(대본·자막·조립 순서)를 완벽하게 해드리니, 무료 편집 프로그램으로 순서만 따라 완성하시면 됩니다.",
      "◆ ‘행사 홍보영상 기획·대본 (만들 준비)’ 카드: 행사명·목적·길이만 넣고 실행하면 → ①장면 구성(스토리보드) ②나레이션 대본 초안 ③화면 자막 문구 ④영상 AI에 붙여넣을 프롬프트 ⑤무료 편집 프로그램으로 조립하는 순서까지 한 장으로 만들어 드립니다.",
      "◆ 만드는 순서(가장 쉬운 길): 1)작업지를 뽑는다 → 2)Clipchamp(윈도우11에 이미 있음)나 Canva에서 홍보영상 템플릿을 연다 → 3)작업지 스토리보드대로 사진·영상을 넣는다 → 4)‘자동 자막’·‘AI 음성’으로 자막과 나레이션을 넣는다 → 5)저작권 무료 음악을 깔고 내보낸다.",
      "◆ 영상 AI: 작업지에 있는 프롬프트를 Runway·Pika·Kling·Canva 등 영상 AI(무료 한도 있음)에 붙여넣으면 오프닝·배경 컷을 만들어 줍니다. (실제 사람·교회 건물은 직접 찍은 사진/영상이 더 자연스럽습니다.)",
      "◆ ‘무료 영상편집·영상 AI 완전 가이드’ 카드: 무료 프로그램 사용법, 영상 AI 사용법, 부활절·성탄·수련회 등 목적별 홍보영상 템플릿, 저작권 안전한 음악·이미지 출처까지 자세히 담겨 있습니다.",
      "◆ 배경음악은 ‘찬양’ 섹션의 ‘AI 찬양 작곡 (수노·작업지)’ 카드로 직접 만들어 넣을 수도 있습니다. 저작권 걱정 없이 우리 영상에 딱 맞는 음악을 쓸 수 있습니다.",
      "◆ 배포: 카카오톡·인스타 릴스는 세로(9:16), 유튜브는 가로(16:9)로 내보내면 화면에 꽉 차게 보입니다.",
      "※ 사진만 넣으면 자동으로 영상까지 만들어 주는 기능은 다음 업데이트에 정식 지원 예정입니다. 이번 판에서는 위 작업지 + 무료 편집기(Clipchamp·Canva)로 만드시는 방식입니다 — 오히려 목사님 손으로 다듬은 영상이 더 자연스럽고 좋습니다.",
      "◆ 설교·예배 영상: ‘설교 슬라이드(PPT) 만들기’ 카드에 설교 제목·본문·대지를 넣으면 예배 화면용 슬라이드(PPT)가 자동으로 만들어집니다. 그리고 ‘예배·설교 영상 — 녹화·유튜브 라이브 완전 가이드’ 카드에 녹화 요령·무료 자동자막·유튜브 업로드·실시간 방송(OBS·Prism Live)까지 쉽고 자세히 담겨 있습니다. (설교 본편은 목사님 목소리로 직접 녹화하시는 것이 가장 은혜롭습니다.)"]),
     ("① 목양 — 성도 관리",[
      "교인 등록: 이름·연락처·주소·생일·결혼기념일·세례일·직분·소속셀·심방주기",
      "가족 등록: 배우자·자녀 등 가족관계 / 직분 이력: 임직·취임 변동 기록",
      "★심방 브리핑: 심방 전에 ‘지난 말씀·받은 기도제목’을 보여줘 같은 말 반복·재질문을 막습니다",
      "심방 기록: 오늘 나눈 말씀·기도제목·후속 저장(교인별 누적)",
      "새가족: 등록 → 주차별 양육 → 정착 대시보드(빠진 사람 경고) → 정식 교인 전환",
      "소그룹(셀): 등록·조원명단·주간보고·셀보고서 인쇄양식",
      "태신자(전도대상): 담당 성도 배정·접촉 기록·결신 시 새가족 자동 전환",
      "생일·기념일: 다가오는 생일/결혼기념일 미리 알림 + 축하 문자 초안",
      "돌봄 필요: 오래 심방 못 한 성도 자동 감지 / 주간 목회 브리핑",
      "경조사: 출생·장례·결혼·회갑·입원 기록(경조금 포함)",
      "★상황별 축하·위로: 군입대·전역·진급·취업·합격·졸업·출산·개업·이사·퇴원·은퇴 등 상황을 고르고 이름을 넣으면, 골라 쓸 수 있는 축하·위로 카톡 문구를 만들어 드립니다(무료 붙여넣기).",
      "★교적 이동: 전입·전출·이명을 기록하고 교인 상태를 자동 갱신합니다. 다른 교회로 옮기는 성도에게는 ‘증명서 발급’에서 이명증서를 떼어 드릴 수 있습니다."]),
     ("★ 성례·전도·다음세대 기능",[
      "성례대장(세례·학습·입교): ‘성례 등록’으로 세례·학습·입교·유아세례를 기록하면 성례대장에 남고, 훗날 기독교 대학 진학 등에 필요한 세례 증명서 발급의 근거가 됩니다. 신청서 서식도 함께 있습니다.",
      "발급대장: 발급한 모든 증명서(교인·세례·재적·이명 등)가 문서번호·발급일·대상·용도와 함께 자동으로 기록됩니다.",
      "새생명축제·태신자운동: ‘준비 타임라인’에 축제일을 넣으면 D-56부터 후속까지 날짜별 준비사항을, ‘당일 진행 체크리스트’와 참고자료(초청카드 서식·부록 18종 포함)를 함께 제공합니다. 요즘은 주일 오전 예배를 축제일(D-DAY)로 삼습니다.",
      "여름공과: ‘여름공과’ 카드 하나에 7주제(1창조·2복음·3동행·4사랑섬김·5믿음·6기도·7제자도)가 모여 있어, 번호 1~7을 넣으면 해당 공과가 열립니다.",
      "구역·속회 예배 순서지: 구역장이 인도하기 좋게 순서·나눔 질문·기도제목 칸이 담긴 순서지를 만들어 드립니다."]),
     ("② 예배·설교",[
      "설교 작성: 예배유형(새벽·수요·금요·주일오전/오후·중고등부 등) 선택 → 유형별 폴더 자동정리",
      "  설교 틀(장면→진리→그리스도) + 내 예화·주석 + 성경본문(KJV·WEB) 준비 → 내용은 목사님이 채움",
      "설교 이력: ‘언제 무슨 본문 설교했나’ 검색 / 지난 설교 재생성(재묵상·재작성)",
      "주보 만들기 · 오늘의 묵상(성도 발송문) · 교안 · 증명서 발급(교인·세례·헌금)",
      "성경 찾기: KJV·WEB(영문) 즉시 조회 / 교회력·절기: 그 해 부활절 기준 절기 자동 산출"]),
     ("★ 설교문 작성 상세 가이드 (꼭 읽어보세요)",[
      "1) 성경 본문 넣기 — ‘성경 찾기’에서 요3:16 을 KJV/WEB로 조회해 복사하거나, 설교 작성 시 본문칸에 ‘요3:16’을 적으세요(영어설교는 본문이 자동 삽입). 개역개정은 저작권상 프로그램에 넣지 않았으니 각자 성경으로 확인하고, 설교문엔 성구 주소를 표기합니다.",
      "2) 내 성경주석 붙여 쓰기 — 가지고 계신 성경주석 PDF를 _시스템\\_내자료\\성경주석 폴더에 넣으면, 설교 초안에 ‘참고 성경주석’ 목록으로 자동 표시됩니다. commentary-add 명령으로 넣으면 NotebookLM(구글 무료 노트도구)에도 올려 ‘이 본문 주석이 어디 있지?’를 빠르게 검색할 수 있습니다. → 그 주석을 펴서 참고하며 목사님이 직접 설교를 완성하십니다.",
      "3) 예화 넣기 — 감동적인 예화를 ‘내 예화 추가’ 기능으로 등록하면, 설교 주제가 맞을 때 자동 추천됩니다.",
      "4) 설교 틀 — 각 대지를 ‘장면 → 진리 → 그리스도’로 전개하도록 틀이 잡혀 있습니다. 내용은 목사님이 채워 완성하십니다.",
      "5) 예배유형 선택 — 새벽·수요·금요·주일오전/오후·중고등부 등을 고르면 그 유형 폴더로 정리되어, ‘설교 이력’에서 예배별로 바로 찾습니다.",
      "6) 지난 설교 재생성 — 예전 설교를 검색해 재묵상·재작성용 새 초안을 만들 수 있고, 같은 본문을 전에 설교했으면 알려줍니다(반복 방지).",
      "7) ※ 이 프로그램은 설교를 ‘대신’ 써 드리지 않습니다. 성경본문·내 성경주석·예화·설교 틀을 한자리에 모아 정리해 드릴 뿐이고, 설교 내용은 목사님이 은혜받아 직접 작성하십니다."]),
     ("★ NotebookLM 무료 활용 — 내 자료로 검색·연구 (선택)",[
      "NotebookLM이란? 구글이 만든 ‘내 자료 전용’ 무료 연구·검색 도구입니다. 내가 올린 자료(성경주석·신학서적·설교/찬양 자료 등) 안에서만 찾아주고 정리해 줍니다.",
      "무료 가입 (1분): 인터넷 브라우저에서  notebooklm.google.com  접속 → 가지고 계신 구글(Gmail) 계정으로 로그인 → ‘새 노트북(New notebook)’ 만들기. 완전 무료입니다.",
      "자료 넣기: 노트북에서 ‘소스 추가(Add source)’로 PDF·문서를 올립니다. 성경주석·설교집·찬양/가사 자료 등 무엇이든 됩니다.",
      "설교에 쓰기: ‘이 본문 관련 주석이 어디 있지?’처럼 물으면 내 자료 안에서 찾아 줍니다 → 그 내용을 펴서 참고하며 목사님이 직접 설교를 작성하십니다.",
      "찬양에 쓰기 (작곡·작사 목사님): 예배·찬양·신학 자료를 올려두고, 작사하실 때 ‘이 주제 관련 내용?’을 검색해 참고합니다.",
      "이 프로그램과 연동 ① 성경주석: commentary-add --pdf \"경로\" --notebook 노트북ID → 주석이 NotebookLM에도 올라가 검색됩니다.",
      "이 프로그램과 연동 ② 찬양곡·가사·기타 자료: 'NotebookLM에 자료 올리기'로 악보·가사·예배자료 등 무엇이든 올려 검색할 수 있습니다.",
      "★성경책별 자동 연결: 성경 각 권마다 노트북을 따로 두셨다면, 노트북ID 대신 --book 로마서 (또는 --book 롬)처럼 책 이름만 대면 그 책 노트북으로 자동 연결됩니다. (church_config.json의 'nlm책노트북'에 책이름→노트북ID를 넣어두면 됩니다. 축약형·부분이름도 인식)",
      "연동 도구(nlm)가 없어도 괜찮습니다 — notebooklm.google.com 웹사이트에서 파일을 직접 ‘소스 추가’로 올리셔도 똑같이 됩니다(더 간단).",
      "※ NotebookLM은 ‘내 자료를 찾아·정리해 주는’ 연구 도구입니다. 설교나 찬양을 대신 만들지 않으며, 창작은 목사님이 하십니다."]),
     ("③ 찬양 사역",[
      "찬양 콘티(A4/A3): 곡 목록 → 악보·Key·BPM·가사·유튜브 자동 편집",
      "찬양곡·자작곡 등록: 작곡·작사·Key·BPM·가사·악보·유튜브·발표일·저작권",
      "찬양 작품집 / 찬양집(악보집) 출판 / AI 찬양 작곡(수노·작업지)",
      "작사·작곡 참고자료 — 가지고 계신 찬양·예배·신학 자료(가사집·묵상노트·주석 등)를 NotebookLM(구글 무료 노트도구)에 올려두면, 작사하실 때 ‘이 주제 관련 내용이 어디 있지?’를 검색해 참고할 수 있습니다. 창작은 목사님이 하시고, 자료를 찾아 드리는 역할입니다.",
      "가사 스크린(문서·PPT) · 찬양팀 배정표 · 유튜브 바로 열기"]),
     ("④ 사역 — 일정·선교·청소년·뮤지컬",[
      "일정 캘린더: 집회·외부설교·행사·생일·기념일·경조사·교회력을 날짜순 한눈에",
      "집회/외부설교 일정(D-day) · 노회 관련(회의·서류·기한) · 대외행사 기획서",
      "청소년: 학생 등록(학교·학년·보호자)·시험 응원 카톡",
      "뮤지컬·공연: 프로덕션 등록·캐스팅표",
      "단기선교: 프로젝트·팀명단(여권·비상연락)·D-day 준비 타임라인·현지정보(선교사·숙소·식당)·영어회화·영어설교·준비 체크리스트"]),
     ("⑤ 재정",[
      "재정 기록(수입/지출) · 월별 요약 · 헌금 케어신호(정기헌금 끊긴 교인)",
      "주일 재정 결산서(전주이월·차주이월 자동) · 재정 출납부(원장) · 헌금·지출 표준 항목표",
      "교인별 헌금대장(항목별) · 연말정산 기부금영수증 발급",
      "★예산 관리(전문 재정프로그램식): ‘예산안 편성 양식’으로 작년 실적을 참고칸에 자동으로 채운 표를 뽑아 제직회에서 정한 뒤 → ‘예산 편성’ 카드에 항목별 예산을 입력하면 → ‘예산 대비 집행 현황’이 연중 집행률·잔액·예산 초과를 자동 계산해 제직회 보고서로 만들어 드립니다.",
      "  예산 편성 팁: 항목명을 ‘헌금·지출 항목표’와 똑같이(예: 십일조, 인건비(사례비)) 쓰면 실제 기록과 정확히 대조됩니다."]),
     ("⑥ 출력·파일",[
      "한글(.hwp) 변환: 만든 docx를 네이티브 한글파일로 — 한컴오피스가 설치된 PC에서만 됩니다(안 되어 있어도 docx는 한글에서 그대로 열립니다)",
      "프린터 출력: 파일을 기본 프린터로 직접 인쇄 / 엑셀 관리대장 · 찬양 가사 PPT",
      "파일 읽기: TXT·DOCX 내용 불러오기"]),
     ("★ 저장 방식 — 어디에 어떻게 저장되나요? (꼭 읽어보세요)",[
      "자동 저장: 무엇을 만들면 곧바로 파일로 저장됩니다. 따로 ‘저장’ 버튼을 누를 필요가 없습니다.",
      "★파일 이름으로 한눈에 구분: 설교·묵상은 ‘[예배유형] 제목 (성경구절)_날짜’ 형식으로 저장돼, 파일명만 봐도 어떤 설교인지 바로 압니다 (예: [금요기도회] 정한 마음을 창조하소서 (시51·10-17)_2026-07-17.docx).",
      "정리 위치: 만든 문서는 01~09 번호 폴더에 종류·예배유형별로 자동 정리됩니다(예: 설교 → 09 설교\\금요기도회).",
      "★바로 열기: 문서를 만들면 결과 화면에 ‘📂 이 문서 열기’와 ‘🗂️ 폴더 열기’ 버튼이 떠서, 방금 만든 파일이나 그 폴더(지난 작업들)를 클릭 한 번에 엽니다.",
      "설교문은 DOCX로 저장됩니다(한글에서 그대로 열립니다). 네이티브 한글(.hwp)이 필요하면 ‘한글(.hwp) 변환’ 카드로 바꾸세요.",
      "모든 기록(교인·심방·재정·설교이력 등)은 한 곳(church_db.json)에 쌓이고, ‘엑셀 관리대장’으로 11개 시트(교인·심방·출석·재정·설교이력·경조사·일정·태신자·소그룹·양육·노회)로 뽑아 따로 관리하실 수 있습니다.",
      "★3중 안전보관: ①원본 폴더  ②‘USB·D 자동저장 설정’으로 지정한 USB·D 폴더(모든 문서·기록·작곡·악보·예화가 자동 복사)  ③🔥피닉스 스냅샷 — 실수로 지워도 지난 시점으로 되살아납니다."]),
     ("⑦ 시스템 — 안심하고 쓰세요",[
      "🔒 자료 영구 보존: 업데이트해도 교인·심방·설교 등 모든 자료는 삭제되지 않습니다.",
      "🔥 피닉스 복구: 실수로 지워져도 지난 시점으로 언제든 되살립니다(자동 스냅샷).",
      "💾 자료 백업: 언제든 전체 백업 / 🔔 업데이트: 새 기능이 오면 표시등이 켜지고 눌러서 업데이트",
      "💡 기능 요청: 원하는 기능을 적으면 담당자에게 전달되어 다음 업데이트에 반영됩니다."]),
     ("★ USB·D 자동 이중저장 (안전이 최고)",[
      "왜: 컴퓨터가 고장나거나 실수로 지워도, USB나 D드라이브에 똑같이 있으면 자료를 잃지 않습니다.",
      "설정(딱 한 번): ‘USB·D 자동저장 설정’ 카드(또는 set-backup)에 폴더를 넣으세요. 예) USB는 E:\\교회백업, D드라이브는 D:\\교회백업",
      "그다음부터: 교인·설교·주보·재정 등 ‘무엇을 만들거나 저장할 때마다’ 그 폴더(USB·D)에 자동으로 함께 저장됩니다. 목사님이 따로 할 일 없습니다.",
      "지금 전체 복사: ‘USB·D로 전체 저장(sync)’을 누르면 지금까지의 모든 자료·문서가 한 번에 그 폴더로 복사됩니다.",
      "TIP: USB를 꽂아두면 예배·설교 자료가 원본(내 PC)+USB 두 곳에 항상 보관됩니다."]),
    ]
    for title,lines in secs:
        para(d,title,15,NAVY,True,before=6,after=4,font=SANS); hr(d,"1B2A4A",6,0,8)
        for ln in lines: para(d,"· "+ln,10.5,after=4)
        para(d,"",after=4)
    para(d,"※ 이 프로그램은 무료입니다. 서버비·사용료가 없습니다(공개 라이브러리·퍼블릭도메인 성경 사용).",9.5,MAROON,italic=True,before=6,font=SANS)
    para(d,"※ 성경 개역개정·NIV, 남의 주석·묵상 전문을 복제 배포하지 마세요(저작권). 성구 주소+창작 묵상은 안전.",9,LG,before=2,font=SANS)
    hr(d,"C9BF9E",8); para(d,"【저작권 및 이용 안내】",10,NAVY,True,before=8,after=2,font=SANS)
    para(d,COPYRIGHT_NOTICE,9,GRAY,after=4,font=SANS)
    folder=CAT("01")
    out=os.path.join(folder,f"[사용설명서] {sanit(CHURCH)}_v{sanit(VERSION)}.docx"); _savedoc(d,out)
    print(f"✅ 사용 설명서(DOCX): {out}")
    print("   ※ 한글(.hwp)로 원하시면 '한글 변환' 카드를 눌러주세요. (DOCX도 한글에서 그대로 열립니다)")
def helpcmd(a):
    print(__doc__); print("""
명령 요약:
  member-add   --name 이름 [--sex --birth --tel --addr --role --cell --leader --baptism --memo]
  member-list  [--cell 셀]          member-show --name 이름
  visit-add    --name 이름 --word 말씀 --prayer "제목1;제목2" [--note --followup --kind --by --date]
  visit-brief  --name 이름          ← 다음 심방 전 반드시! (중복 말씀·재질문 방지)
  cert         --name 이름 --kind 교인|세례|헌금 [--purpose 용도 --no 번호]
  attend-add   --service 주일낮 --men N --women N [--new N --date]
  finance-add  --kind 수입|지출 --item 항목 --amount 금액 [--name --memo --date]
  finance-sum  [--month YYYY-MM]
  bulletin     --sermon 설교 --notice "소식1;소식2" --week "일정1;일정2" [--order --pray --date]
  lesson       --title 제목 --target 대상 --text 본문 --theme 주제 --goal 목표 [--date]
  setlist      --songs "곡1;곡2;곡3" [--size A4|A3 --title 제목 --date]  ← 찬양팀 콘티(악보·key·BPM·가사)
  song-add     --title 곡명 [--composer 작곡가 --key C --bpm 72 --theme 주제 --lyrics 가사 --sheet 악보경로 --youtube URL]  ← 찬양곡·자작곡 등록(유튜브 링크 포함)
  song-list    [--theme 주제]         ← 내 찬양곡 라이브러리
  song-open    [--title 곡명]         ← 작곡한 찬양곡 유튜브 바로 열기(목록만 보려면 인자 없이)
  song-catalog [--title]              ← 찬양 작품집(자작곡 발표순·유튜브·저작권) · 사역 유산 보존
  songbook     [--songs "곡1;곡2" --title]  ← 찬양집(악보집) 출판(표지·목차·곡별 악보/가사)
  compose      --title 곡명 --theme 주제 [--bible --mood --key --bpm --lyrics]  ← 찬양 작곡 작업지(가사틀+작곡 프롬프트)
  lyrics-screen --songs "곡1;곡2" [--title]  ← 예배 프로젝터용 큰글씨 가사 슬라이드
  worship-roster --date --leader --keys --guitar --bass --drums --singers [--songs]  ← 찬양팀 배정표+리마인드
  ── 집회·청소년·뮤지컬·선교 사역 ──
  schedule-add --date --type 집회|부흥회|찬양집회 --place 교회 --host 담당 --tel --theme --fee 사례  ← 집회·초청 일정
  schedule-list [--upcoming 1]        ← 일정 목록(D-day)
  student-add  --name --school 학교 --grade 학년 --dept 부서 --tel --guardian 보호자 --gtel 보호자연락처 [--major 전공]  ← 청소년/대안학교 학생
  student-list [--dept 부서 --grade 학년]   exam-cheer --exam 중간고사 [--dept]  ← 학생명단 / 시험응원 카톡
  production-add --title 작품 --date 공연일 --place --director 연출 --music 음악감독   production-list  ← 뮤지컬·공연
  casting      --title 작품 --roles "배역:배우:파트;..."  ← 캐스팅표
  mission-add  --title 팀명 --country 국가 --period 기간 --purpose 목적   ← 단기선교 프로젝트
  mission-member --title 팀 --name --role 역할 --expire 여권만료 --emergency 비상연락 --health 건강  ← 선교 팀원
  mission-roster --title 팀 / mission-checklist --title 팀  ← 선교 팀명단·준비 체크리스트
  mission-plan  --title 팀 --dday 출국일  ← 단기선교 D-day 단계별 준비 타임라인
  mission-place --title 팀 --kind 선교사|숙소|식당|사역지 --name --contact --address  ← 현지 정보 등록
  mission-field --title 팀            ← 현지 정보 시트(선교사·숙소·식당·사역지)+일자별 플랜
  mission-english / mission-sermon --title --text 본문 --points  ← 선교지 영어회화 / 영어설교(영어성경 자동)
  ── 설교 관리 · 새가족 정착 ──
  sermon-log   --date --title --text 본문 --service --series  ← 설교 이력 기록
  sermon-list  [--query 검색 --year]  ← "언제 무슨 설교 했나"(본문·제목 검색·반복확인)
  song-catalog / songbook             ← 찬양 작품집 / 찬양집(악보집) 출판
  newfamily-add --name --tel --leader 인도자 --cell 셀   ← 새가족 등록(정착관리 시작)
  newfamily-step --name --how 심방|전화|문자 --note      ← 새가족 양육 접촉 기록
  newfamily-board                     ← 새가족 정착 대시보드(미이행 경고·뒷문 방지)
  newfamily-graduate --name [--role]  ← 정착 완료→정식 교인 전환
  ── 시스템 ──
  version   ← 버전·업데이트 안전 안내    backup   ← 자료 전체 백업
  update    --file 새church.py경로     ← 프로그램 업데이트(자료 자동백업·보존)
  menu                                ← 번호만 누르는 간편 메뉴(명령어 몰라도 사용)
  sermon       --title 제목 --text 본문 --theme 주제 --points "대지1;대지2;대지3"  ← 설교초안(예화·주석 참고)
  illus-add    --name 예화명 --source 출처 --topic "키워드;키워드" --bible 성경연결 --text 전문  ← 내 예화 추가
  illus-list                          ← 내 예화은행 목록
  commentary-add  --pdf 경로 [--notebook 노트북ID]  ← 내 성경주석 PDF 추가(선택:NotebookLM 검색연동)·설교 참고
  commentary-list                     ← 내 성경주석 라이브러리 목록
  bible        --ref "요한복음 3:16" --version 개역개정|niv|kjv   ← 성경본문(데이터 필요)
  devotion     --title 제목 --text 본문 --verse 말씀 --body 묵상 ← 오늘의 묵상(성도 발송)
  birthday     [--days 30]            ← 성도 생일 축하(다가오는 생일·결혼기념일+축하문자·셀명단)
  care         [--days 90]            ← 돌봄 필요 성도(심방 오래 경과)
  newcomer     [--days 90]            ← 새가족 정착 파이프라인
  kakao        --type birthday|devotion|custom [--to 셀|all --text --ahead]  ← 카톡에 붙여넣기 목록(무료 붙여넣기형)
  weekly-brief                        ← 주간 목회 브리핑(출석·재정·생일·돌봄 한 장 요약)
  prayer-add   --name 교인 --text 내용 --theme 주제 / prayer-digest  ← 기도제목 접수·주제별 정리
  giving-insight [--days 60]          ← 헌금 케어 신호(정기헌금 끊긴 교인 감지)
  export-excel                       ← DB를 색깔 Excel 관리대장으로 동기화
""")

# ───────── 행정서식 (v2 batch A) ─────────
def _docno(db,kind="문서"):
    """문서번호 자동채번 — 교회명약자-연도-일련(연도별 누적). ★db 객체를 받아 증가만 하고 저장은 호출자가 한 번(stale-read/lost-update 방지)."""
    y=datetime.date.today().year
    seqs=db.setdefault("_docseq",{}); k=f"{kind}:{y}"
    seqs[k]=seqs.get(k,0)+1
    return f"{(CHURCH or '교회')[0]}-{y}-{seqs[k]:03d}"
def meeting_minutes(a):
    """🗒️ 회의록 — 종류(제직회/공동의회/당회/교사회)·안건→토의→의결 구획·참석/불참 집계·서명란."""
    kind=a.type or "제직회"; date=a.date or today()
    att=[x.strip() for x in re.split(r'[;,]', a.attend or "") if x.strip()]
    absn=[x.strip() for x in re.split(r'[;,]', a.absent or "") if x.strip()]
    agenda=[x.strip() for x in re.split(r'[;\n]', a.agenda or "") if x.strip()]
    d=newdoc()
    para(d,CHURCH,11,GOLD,True,AL.CENTER,after=0,font=SANS)
    para(d,f"{kind} 회의록",22,NAVY,True,AL.CENTER,before=4,after=2,font=SANS); hr(d,"1F5C9E",12)
    info=[("일시",f"{date} {a.time or ''}".strip()),("장소",a.place or ""),("사회",a.host or PASTOR),("서기",a.clerk or "")]
    t=d.add_table(rows=2,cols=4)
    for i,(k,v) in enumerate(info):
        rr=0 if i<2 else 1; cc=(i%2)*2
        ctext(t.rows[rr].cells[cc],k,10,WHITE,True,SANS,AL.CENTER); shade(t.rows[rr].cells[cc],"1F5C9E")
        ctext(t.rows[rr].cells[cc+1],v,10,GRAY,font=SERIF)
    tborders(t)
    para(d,f"참석 {len(att)}명 · 불참 {len(absn)}명",11,MAROON,True,before=8,after=3,font=SANS)
    if att: para(d,"참석: "+", ".join(att),10,GRAY,after=2,font=SERIF)
    if absn: para(d,"불참: "+", ".join(absn),10,LG,after=4,font=SERIF)
    hr(d,"CCCCCC",4)
    para(d,"■ 안건 및 의결사항",13,NAVY,True,before=6,after=4,font=SANS)
    for i,ag in enumerate(agenda or ["(안건 — 직접 입력)"],1):
        para(d,f"제{i}호 안건: {ag}",12,NAVY,True,before=6,after=2,font=SANS)
        para(d,"  · 토의: ",11,GRAY,after=2,font=SERIF)
        para(d,"  · 의결: ",11,MAROON,after=6,font=SERIF)
    hr(d,"CCCCCC",4)
    para(d,"위 회의록이 사실과 같음을 확인합니다.",11,GRAY,AL.CENTER,before=10,after=10,font=SERIF)
    para(d,date,11,GRAY,AL.CENTER,after=8,font=SANS)
    para(d,f"사회 {a.host or PASTOR}  (인)          서기 {a.clerk or ''}  (인)",12,NAVY,True,AL.CENTER,font=SANS)
    para(d,f"· {CHURCH}",9,LG,align=AL.RIGHT,before=8,font=SANS)
    out=os.path.join(CAT("01"),f"[회의록] {kind}_{date}.docx"); _savedoc(d,out)
    print(f"✅ {kind} 회의록 생성: {out}")
def official_doc(a):
    """📄 공문(협조문) — 문서번호 자동채번·수신/제목/본문/발신 표준 두문·본문·결문."""
    db=load(); no=_docno(db,"공문"); save(db); date=a.date or today()
    d=newdoc()
    para(d,CHURCH,16,NAVY,True,AL.CENTER,after=2,font=SANS); hr(d,"1B2A4A",16)
    para(d,f"문서번호: 제 {no} 호",10,GRAY,after=1,font=SANS)
    para(d,f"시행일자: {date}",10,GRAY,after=1,font=SANS)
    para(d,f"수    신: {a.to or '수신처'}",11,NAVY,True,after=1,font=SANS)
    if a.via: para(d,f"경    유: {a.via}",10,GRAY,after=1,font=SANS)
    para(d,f"제    목: {a.title or '(제목)'}",12,MAROON,True,before=4,after=8,font=SANS); hr(d,"CCCCCC",4)
    for pnum,line in enumerate([l.strip() for l in re.split(r'[\n;]', a.body or "") if l.strip()] or ["(본문 — 직접 입력)"],1):
        para(d,f"{pnum}. {line}",11,GRAY,after=4,font=SERIF)
    para(d,"위와 같이 협조하여 주시기 바랍니다.  끝.",11,GRAY,before=6,after=16,font=SERIF)
    para(d,date,13,NAVY,True,AL.CENTER,after=6,font=SANS)
    para(d,f"{CHURCH}  {a.sender or ('담임목사 '+PASTOR)}   (직인)",14,NAVY,True,AL.CENTER,font=SANS)
    out=os.path.join(CAT("07"),f"[공문] {sanit(a.title or '공문')}_{date}.docx"); _savedoc(d,out)
    print(f"✅ 공문 발급(제 {no} 호): {out}")
def asset_register(a):
    """📦 비품대장 — 품목·수량·구입일·위치·담당·상태(가로 대장, 누적). --name 지정 시 추가 후 대장 재생성."""
    from docx.enum.section import WD_ORIENT
    db=load(); db.setdefault("비품",[])
    if a.name:
        db["비품"].append({"품목":a.name,"수량":a.qty or "1","구입일":a.date or today(),"위치":a.place or "","담당":a.manager or "","상태":a.status or "사용"}); save(db)
        print(f"  + 비품 등록: {a.name}")
    items=db["비품"]
    d=newdoc()
    for s in d.sections: s.orientation=WD_ORIENT.LANDSCAPE; s.page_width,s.page_height=s.page_height,s.page_width
    para(d,f"{CHURCH} 비품대장",20,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,f"작성일 {today()} · 총 {len(items)}건",10,GRAY,AL.CENTER,after=6,font=SANS); hr(d,"1F5C9E",10)
    hdr=["번호","품목","수량","구입일","위치","담당","상태"]
    BLANK=18; nrows=max(len(items),BLANK)   # 기록 + 손기입 빈칸, 최소 18줄
    t=d.add_table(rows=nrows+1,cols=len(hdr))
    for j,h in enumerate(hdr): ctext(t.rows[0].cells[j],h,10,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],"1F5C9E")
    for i in range(1,nrows+1):
        if i<=len(items):
            it=items[i-1]; vals=[str(i),it.get("품목",""),str(it.get("수량","")),it.get("구입일",""),it.get("위치",""),it.get("담당",""),it.get("상태","")]
        else:
            vals=[str(i),"","","","","",""]
        for j,v in enumerate(vals):
            ctext(t.rows[i].cells[j],v,10,GRAY,font=SERIF)
    tborders(t)
    out=os.path.join(CAT("07"),f"[비품대장] {today()}.docx"); _savedoc(d,out)
    print(f"✅ 비품대장 생성({len(items)}건): {out}")
def vehicle_log(a):
    """🚗 차량 운행일지 — 일자·운전자·행선지·거리(km)·주유(누적). --driver/--dest 지정 시 추가 후 일지 재생성."""
    db=load(); db.setdefault("차량운행",[])
    if a.driver or a.dest:
        db["차량운행"].append({"일자":a.date or today(),"운전자":a.driver or "","행선지":a.dest or "","거리":a.km or "","주유":a.fuel or ""}); save(db)
        print(f"  + 운행 기록: {a.date or today()} {a.driver or ''} {a.dest or ''}")
    rows_=db["차량운행"]
    tot=sum(int(re.sub(r'\D','',str(r.get('거리','') or '0')) or 0) for r in rows_)
    d=newdoc()
    para(d,f"{CHURCH} 차량 운행일지",20,NAVY,True,AL.CENTER,after=2,font=SANS)
    para(d,f"작성일 {today()} · 총 {len(rows_)}건 · 누적 {tot}km",10,GRAY,AL.CENTER,after=6,font=SANS); hr(d,"1F5C9E",10)
    hdr=["번호","일자","운전자","행선지","거리(km)","주유"]
    BLANK=18   # 손으로 바로 기입할 빈 줄(기록이 적어도 양식이 되게)
    nrows=max(len(rows_),BLANK)   # 기록 + 빈칸, 최소 18줄
    t=d.add_table(rows=nrows+1,cols=len(hdr))
    for j,h in enumerate(hdr): ctext(t.rows[0].cells[j],h,10,WHITE,True,SANS,AL.CENTER); shade(t.rows[0].cells[j],"1F5C9E")
    for i in range(1,nrows+1):
        if i<=len(rows_):
            r=rows_[i-1]; vals=[str(i),r.get("일자",""),r.get("운전자",""),r.get("행선지",""),str(r.get("거리","")),str(r.get("주유",""))]
        else:
            vals=[str(i),"","","","",""]   # 번호만 매긴 빈 기입란
        for j,v in enumerate(vals):
            ctext(t.rows[i].cells[j],v,10,GRAY,font=SERIF)
    tborders(t)
    para(d,"※ 표에 직접 손으로 적으시거나, 프로그램 '차량운행일지' 카드에 운전자·행선지·거리를 넣어 실행하면 자동으로 채워집니다.",9,LG,before=6,font=SANS)
    out=os.path.join(CAT("07"),f"[차량운행일지] {today()}.docx"); _savedoc(d,out)
    print(f"✅ 차량 운행일지 생성({len(rows_)}건): {out}")

# ───────── 목회 참고자료(6대 도메인 + 주일학교) — 저작권 안전 지침서 → DOCX/HWP ─────────
_REF_DOCS=[
    ("주일학교","교회학교·주일학교","_v2참고자료_주일학교.md","ss"),
    ("성경","성경자료·성경교육","_v2참고자료_성경자료.md","bible"),
    ("찬양","찬양·예배음악","_v2참고자료_찬양.md","wor"),
    ("전도","전도·새신자","_v2참고자료_전도새신자.md","ev"),
    ("제자훈련","제직·제자훈련","_v2참고자료_제직제자훈련.md","dis"),
    ("가정","가정·결혼·상담","_v2참고자료_가정상담.md","fam"),
    ("행정","교회행정·재정","_v2참고자료_교회행정재정.md","adm"),
    ("선교영어","선교지 선교영어","_v2참고자료_선교영어.md","me"),
    ("갈등관리","목회자 갈등관리(성도·당회)","_v2참고자료_교회갈등.md","conflict"),
    ("여름수련회","여름성경학교·수련회 준비","_v2참고자료_여름수련회.md","summer"),
    ("여름공과창조","여름공과 ① 나를 지으신 하나님","_v2참고자료_여름공과_1창조.md","vbs1"),
    ("여름공과복음","여름공과 ② 나를 구원하신 예수님","_v2참고자료_여름공과_2복음.md","vbs2"),
    ("여름공과동행","여름공과 ③ 함께하시는 하나님","_v2참고자료_여름공과_3동행.md","vbs3"),
    ("여름공과사랑","여름공과 ④ 사랑하며 섬기는 우리","_v2참고자료_여름공과_4사랑섬김.md","vbs4"),
    ("여름공과믿음","여름공과 ⑤ 믿음의 사람들","_v2참고자료_여름공과_5믿음.md","vbs5"),
    ("여름공과기도","여름공과 ⑥ 기도하는 다음세대","_v2참고자료_여름공과_6기도.md","vbs6"),
    ("여름공과제자","여름공과 ⑦ 예수님을 따라가요","_v2참고자료_여름공과_7제자도.md","vbs7"),
    ("심방","심방 지침(상황별)","_v2참고자료_심방.md","visit"),
    ("예식","예식서(순서·기도문)","_v2참고자료_예식.md","ceremony"),
    ("기도문","대표기도문·상황별","_v2참고자료_기도문.md","prayer"),
    ("성경공부","성경공부 인도","_v2참고자료_성경공부.md","bstudy"),
    ("목회자영성","목회자 영성·자기돌봄","_v2참고자료_목회자영성.md","spirit"),
    ("선교","단기선교·비전트립 준비","_v2참고자료_선교.md","mission"),
    ("새신자교육","새신자 교육 4주 프로그램","_v2참고자료_새신자교육.md","newbeliever"),
    ("새생명축제","새생명 축제(태신자 운동) 준비","_v2참고자료_새생명축제.md","harvest"),
    ("절기사역","교회 절기 사역 준비","_v2참고자료_절기사역.md","season"),
    ("영상제작","무료 영상편집·영상 AI 완전 가이드","_v2참고자료_영상제작.md","video"),
    ("예배영상방송","예배·설교 영상 — 녹화·유튜브 라이브 완전 가이드","_v2참고자료_예배영상방송.md","worshipvid"),
    ("이단분별","이단 분별·대처 (성도 보호)","_v2참고자료_이단분별.md","cult"),
    ("비전코칭","청소년 비전코칭(진로·소명)","_v2참고자료_비전코칭.md","vision"),
    ("기독교세계관","청소년 기독교 세계관","_v2참고자료_기독교세계관.md","worldview"),
    ("성경적상담","성경적 상담 실전 가이드","_v2참고자료_성경적상담.md","counsel"),
    ("평신도사역","평신도 사역·사역자 세우기","_v2참고자료_평신도사역.md","laity"),
]
def _ref_find(key):
    """도메인 키워드로 참고자료 md 파일 찾기. ★정확일치(rid·short) 먼저 — '선교'가 '선교영어'로 잘못 가는 오매칭 방지."""
    if not key: return None
    k=key.strip().lower()
    # 1) 정확 일치(rid 또는 short) 우선 — 카드(ref-*)는 이 경로로 정확히 매칭
    for short,title,fn,rid in _REF_DOCS:
        if k==rid or k==short.lower():
            p=os.path.join(ROOT,fn)
            if os.path.exists(p): return (short,title,p)
    # 2) 부분 일치(관대) — 사용자가 직접 도메인 일부를 입력한 경우
    for short,title,fn,rid in _REF_DOCS:
        s=short.lower(); t=title.lower()
        if s in k or k in s or k in t:
            p=os.path.join(ROOT,fn)
            if os.path.exists(p): return (short,title,p)
    return None
def _md_runs(p,text,size=11,font=SERIF,color=None):
    """**굵게** 를 반영해 한 문단에 여러 run 추가."""
    for seg in re.split(r'(\*\*.+?\*\*)', text):
        if not seg: continue
        b=seg.startswith('**') and seg.endswith('**')
        r=p.add_run(seg[2:-2] if b else seg); r.font.size=Pt(size); r.bold=b
        if color: r.font.color.rgb=color
        setf(r,font)
def reference_list(a):
    """목회 참고자료 목록 — 6대 도메인 + 주일학교(총 7편) 저작권 안전 지침서."""
    print(f"📚 {CHURCH} 목회 참고자료 (저작권 안전·바로 출력 가능)")
    n=0
    for short,title,fn,rid in _REF_DOCS:
        ok=os.path.exists(os.path.join(ROOT,fn))
        if ok: n+=1
        print(f"  {'✅' if ok else '—'} [{short}] {title}")
    print(f"\n  총 {n}편 준비됨. 참고자료 섹션에서 원하는 자료 카드를 누르면 문서가 만들어져 바로 열립니다.")
def _ref_chapters(md):
    """md를 (머리말줄들, [{num,title,lines}]) 로 분할 — '## ' 레벨2 헤더 기준."""
    head=[]; chaps=[]; cur=None
    for ln in md.split("\n"):
        st=ln.strip()
        if st.startswith('## ') and not st.startswith('### '):
            if cur: chaps.append(cur)
            mm=re.match(r'##\s+(\d+)', st)
            cur={"num":int(mm.group(1)) if mm else None,"title":st[3:].strip(),"lines":[ln]}
        elif cur is None: head.append(ln)
        else: cur["lines"].append(ln)
    if cur: chaps.append(cur)
    return head,chaps
def _ref_parts(spec):
    """'4' · '3-5' · '2,4,6' → {정수...}"""
    want=set()
    for tok in re.split(r'[,\s]+',(spec or "").strip()):
        if not tok: continue
        if '-' in tok:
            try:
                a2,b2=tok.split('-',1); want.update(range(int(a2),int(b2)+1))
            except Exception: pass
        else:
            try: want.add(int(tok))
            except Exception: pass
    return want
def _add_pagenum(d):
    """DOCX 하단 가운데 '- N -' 페이지 번호 — 한글·워드 모두 인식하도록 완전한 필드(begin·instr·separate·결과값·end).
       (begin/end만 있으면 한글이 0으로 표시 → separate와 결과값 '1'을 넣어 해결)"""
    for sec in d.sections:
        fp=sec.footer.paragraphs[0] if sec.footer.paragraphs else sec.footer.add_paragraph()
        fp.alignment=AL.CENTER
        r1=fp.add_run("- "); r1.font.size=Pt(9); r1.font.color.rgb=GRAY; setf(r1,SANS)
        r=fp.add_run(); r.font.size=Pt(9); r.font.color.rgb=GRAY; setf(r,SANS)
        b=OxmlElement('w:fldChar'); b.set(qn('w:fldCharType'),'begin')
        it=OxmlElement('w:instrText'); it.set(qn('xml:space'),'preserve'); it.text=" PAGE "
        sep=OxmlElement('w:fldChar'); sep.set(qn('w:fldCharType'),'separate')
        tv=OxmlElement('w:t'); tv.text="1"
        en=OxmlElement('w:fldChar'); en.set(qn('w:fldCharType'),'end')
        for el in (b,it,sep,tv,en): r._r.append(el)
        r2=fp.add_run(" -"); r2.font.size=Pt(9); r2.font.color.rgb=GRAY; setf(r2,SANS)
def _ref_table(d, rows, colored=False):
    """참고자료 표. colored=True: 고급 색표(남색 헤더+얼룩·정식 markdown 표) / False: 무채색 깔끔(ASCII·서식 표). 공통: 여백·세로가운데·테두리."""
    rows=[r for r in rows if r]
    if not rows: return
    cols=max(len(r) for r in rows); t=d.add_table(rows=len(rows),cols=cols)
    for ri,r in enumerate(rows):
        for j in range(cols):
            vv=re.sub(r'\*\*','',r[j]) if j<len(r) else ""
            if ri==0:
                if colored: ctext(t.rows[ri].cells[j],vv,9.5,WHITE,True,SANS,AL.CENTER); shade(t.rows[ri].cells[j],"1B2A4A")
                else: ctext(t.rows[ri].cells[j],vv,9.5,NAVY,True,SANS,AL.CENTER)
            else: ctext(t.rows[ri].cells[j],vv,9.5,GRAY,font=SERIF)
    tborders(t,"BBBBBB" if colored else "C9BF9E",4)
    mar=OxmlElement('w:tblCellMar')
    for side,val in (('top','40'),('bottom','40'),('left','100'),('right','100')):
        e=OxmlElement('w:'+side); e.set(qn('w:w'),val); e.set(qn('w:type'),'dxa'); mar.append(e)
    t._tbl.tblPr.append(mar)
    for ri,row in enumerate(t.rows):
        for c in row.cells:
            tcPr=c._tc.get_or_add_tcPr(); vA=OxmlElement('w:vAlign'); vA.set(qn('w:val'),'center'); tcPr.append(vA)
        if colored and ri>0 and ri%2==0:
            for c in row.cells: shade(c,"F5F2EB")
def vbs(a):
    """📖 여름공과 (7주제 모음) — 번호 1~7로 해당 공과를 엽니다. 1창조·2복음·3동행·4사랑섬김·5믿음·6기도·7제자도."""
    m={"1":"여름공과창조","2":"여름공과복음","3":"여름공과동행","4":"여름공과사랑","5":"여름공과믿음","6":"여름공과기도","7":"여름공과제자"}
    raw=(getattr(a,'no','') or "1").strip()   # 비우면 1과(창조)로
    dg=re.search(r'[1-7]',raw); key=m.get(dg.group() if dg else raw)
    if not key:
        print("여름공과 주제 번호를 1~7로 넣어 주세요:")
        print("  1 창조(나를 지으신 하나님) · 2 복음(구원하신 예수님) · 3 동행(함께하시는 하나님)")
        print("  4 사랑섬김 · 5 믿음의 사람들 · 6 기도하는 다음세대 · 7 제자도(예수님을 따라가요)")
        return
    setattr(a,'name',key); reference(a)
def reference(a):
    """목회 참고자료 문서화 — 지침서를 서식 DOCX(원하면 .hwp)로. --part 장선택·--toc 목차·페이지번호 포함."""
    hit=_ref_find(getattr(a,"name",None))
    if not hit:
        if (getattr(a,"name",None) or "").strip():
            print(f"✗ '{a.name}'에 맞는 자료가 없습니다. 아래 목록의 도메인 이름을 '도메인' 칸에 넣어 주세요.\n")
        else:
            print("👇 원하는 자료의 도메인 이름을 '도메인' 칸에 넣고 다시 실행하면 그 자료가 출력됩니다. (장 번호·목차·한글 옵션도 사용 가능)\n")
        reference_list(a); return
    short,title,path=hit
    try: md=open(path,encoding="utf-8").read()
    except Exception as e: print(f"✗ 자료를 읽지 못했습니다: {e}"); return
    head,chaps=_ref_chapters(md)
    if str(getattr(a,"toc","") or "").strip() in ("1","y","Y","예"):
        nnum=len([c for c in chaps if c['num']])
        print(f"📑 [{title}] 목차 (번호 있는 장 {nnum}개) — '장 번호' 칸에 번호(예: 4 또는 3-5)를 넣으면 그 부분만 출력됩니다.")
        for c in chaps: print(f"  {c['title']}" if c['num'] else f"  · {c['title']}")
        return
    suffix=""; part=(getattr(a,"part","") or "").strip()
    if part:
        want=_ref_parts(part); sel=[c for c in chaps if c['num'] in want]
        if not sel:
            print(f"✗ '{part}'에 해당하는 장이 없습니다. '목차 보기' 칸에 1을 넣어 장 번호를 먼저 확인하세요."); return
        md="\n".join(head+[ln for c in sel for ln in c['lines']])
        suffix=" — "+",".join(str(c['num']) for c in sel)+"장"
    d=newdoc()
    para(d,f"{CHURCH} · 목회 참고자료",9,GOLD,align=AL.CENTER,before=8,after=1,font=SANS)
    para(d,title+suffix,24,NAVY,True,AL.CENTER,before=1,after=1,font=SANS)
    para(d,today(),9,LG,align=AL.CENTER,after=9,font=SANS); hr(d,"A67C1E",14)
    lines=md.split("\n"); i=0; N=len(lines)
    while i<N:
        s=lines[i].strip()
        if s.startswith('```'):     # ``` 코드펜스 안 ASCII/파이프 표 → 진짜 DOCX 표로
            blk=[]; i+=1
            while i<N and not lines[i].strip().startswith('```'): blk.append(lines[i]); i+=1
            i+=1; tb=[]
            def _ft():
                if tb: _ref_table(d,tb,False); tb.clear()   # ``` 안 ASCII/서식 표 → 무채색 깔끔
            for ln in blk:
                st=ln.strip()
                if not st: continue
                if set(st)<=set('|-─—=+ :·┌┐└┘├┤┬┴┼│╭╮╯╰'): continue    # 구분선·박스 테두리 스킵
                if '│' in st or '|' in st:                              # │(박스) 또는 |(파이프) 열 구분
                    cs=[c.strip() for c in re.split(r'[│|]', st)]
                    while cs and cs[0]=='': cs.pop(0)
                    while cs and cs[-1]=='': cs.pop()
                    if cs: tb.append(cs)
                else: _ft(); para(d,st,10,GRAY,after=2,font=SANS)
            _ft(); para(d,"",after=4); continue
        if s.startswith('|'):     # 표 — 정식 markdown(|---| 구분선)=고급 색표 / 그 외 ASCII 파이프표=무채색
            is_md = (i+1<N and re.match(r'^\|[\s:\|\-]+\|?$', lines[i+1].strip()) is not None)
            prows=[]
            while i<N and lines[i].strip().startswith('|'):
                st=lines[i].strip()
                if set(st)<=set('|-─—=+ :·'): i+=1; continue
                prows.append([c.strip() for c in st.strip('|').split('|')]); i+=1
            if prows: _ref_table(d,prows,is_md); para(d,"",after=4)
            continue
        if s.startswith('#### '): para(d,re.sub(r'\*\*','',s[5:]),11,GRAY,True,before=4,after=2,font=SANS); i+=1; continue
        if s.startswith('### '):  para(d,re.sub(r'\*\*','',s[4:]),12.5,GOLD,True,before=6,after=2,font=SANS); i+=1; continue
        if s.startswith('## '):   para(d,re.sub(r'\*\*','',s[3:]),15,NAVY,True,before=10,after=3,font=SANS); i+=1; continue
        if s.startswith('# '):    i+=1; continue
        if s=='---':              hr(d,"D8CBB0",6); i+=1; continue
        if s.startswith('>'):
            buf=[]
            while i<N and lines[i].strip().startswith('>'):
                buf.append(lines[i].strip()[1:].strip()); i+=1
            p=d.add_paragraph(); p.paragraph_format.left_indent=Cm(0.5)
            p.paragraph_format.space_before=Pt(4); p.paragraph_format.space_after=Pt(6)
            _md_runs(p," ".join(x for x in buf if x),10.5,SERIF,GRAY); continue
        if re.match(r'^- \[[ xX]\] ', s):
            while i<N and re.match(r'^- \[[ xX]\] ', lines[i].strip()):
                p=d.add_paragraph(); p.paragraph_format.space_after=Pt(2); p.paragraph_format.left_indent=Cm(0.3)
                r0=p.add_run("☐  "); r0.font.size=Pt(11); setf(r0,SANS)
                _md_runs(p,re.sub(r'^- \[[ xX]\] ','',lines[i].strip()),11,SERIF); i+=1
            continue
        if s.startswith('- '):
            while i<N and lines[i].strip().startswith('- ') and not re.match(r'^- \[[ xX]\] ', lines[i].strip()):
                p=d.add_paragraph(); p.paragraph_format.space_after=Pt(2); p.paragraph_format.left_indent=Cm(0.5)
                r0=p.add_run("•  "); r0.font.size=Pt(11); r0.font.color.rgb=GOLD; setf(r0,SANS)
                _md_runs(p,lines[i].strip()[2:],11,SERIF); i+=1
            continue
        if re.match(r'^\d+\. ', s):
            while i<N and re.match(r'^\d+\. ', lines[i].strip()):
                m=re.match(r'^(\d+)\. (.*)', lines[i].strip())
                p=d.add_paragraph(); p.paragraph_format.space_after=Pt(2); p.paragraph_format.left_indent=Cm(0.5)
                r0=p.add_run(f"{m.group(1)}.  "); r0.font.size=Pt(11); r0.bold=True; r0.font.color.rgb=NAVY; setf(r0,SANS)
                _md_runs(p,m.group(2),11,SERIF); i+=1
            continue
        if len(s)>2 and s.startswith('*') and s.endswith('*') and not s.startswith('**'):
            para(d,s[1:-1],9,LG,before=8,after=2,font=SANS,italic=True); i+=1; continue
        if s=='': i+=1; continue
        p=d.add_paragraph(); p.paragraph_format.space_after=Pt(5); _md_runs(p,s,11,SERIF); i+=1
    _add_pagenum(d)
    out=os.path.join(CAT("03"),f"[참고자료] {sanit(title+suffix)}_{today()}.docx")
    if not _savedoc(d,out): return
    print(f"✅ 문서를 만들어 바로 열었습니다{suffix}: {os.path.basename(out)}")
    print("   📄 각 쪽 아래에 페이지 번호가 있습니다. 인쇄(Ctrl+P) 후 '인쇄 범위'에 원하는 쪽(예: 3-5)을 넣으세요.")
    print("   ※ 한글에서 번호가 안 보이면: 인쇄 미리보기를 열거나, 쪽 아무 곳 누르고 자판 F9(필드 새로고침) 한 번.")
    if str(getattr(a,"hwp","") or "").strip() in ("1","y","Y","예"):
        hp=_also_hwp(out)
        if hp: print(f"   ✅ 한글(.hwp)도 저장: {hp}")
    # (문서 열기는 _savedoc가 자동 처리)

def main():
    p=argparse.ArgumentParser(add_help=False); sub=p.add_subparsers(dest="cmd")
    def add(name,fn,*flags):
        sp=sub.add_parser(name);
        for fl in flags: sp.add_argument("--"+fl, dest=fl)
        sp.set_defaults(func=fn); return sp
    add("help",helpcmd)
    add("member-add",member_add,"name","sex","birth","tel","addr","role","cell","leader","baptism","baptismdate","wedding","cycle","memo","date","family")
    add("office-add",office_add,"name","role","memo","date")
    add("office-list",office_list,"name")
    ml=add("member-list",member_list,"cell")
    add("member-show",member_show,"name")
    add("family-add",family_add,"name","member","relation","birth","role","tel")
    add("family-list",family_list,"name")
    add("group-add",group_add,"name","leader","subleader","day","place","memo")
    add("group-list",group_list)
    add("group-assign",group_assign,"name","group")
    add("group-roster",group_roster,"name")
    add("group-report",group_report,"name","date","attend","note","prayer","absent")
    add("group-form",group_form,"name")
    add("train-add",train_add,"name","course","stage","status","memo","date")
    add("train-list",train_list,"name","course")
    add("visit-add",visit_add,"name","word","prayer","note","followup","kind","by","date")
    add("visit-daesim",visit_daesim,"year")
    add("visit-brief",visit_brief,"name")
    add("cert",cert,"name","kind","role","term","purpose","no")
    add("cert-ledger",cert_ledger,"name")
    a1=add("attend-add",attend_add,"service","date");
    for f in ("men","women","new"): a1.add_argument("--"+f,type=_pint,default=0)
    f1=add("finance-add",finance_add,"kind","item","name","memo","date","month","dept"); f1.add_argument("--amount",type=_pint,default=0)
    add("finance-sum",finance_sum,"month")
    add("finance-stats",finance_stats,"year")
    add("finance-report",finance_report,"date","dfrom","dto","year")
    add("finance-items",finance_items)
    add("finance-ledger",finance_ledger,"month","year","dfrom","dto")
    b1=add("budget-set",budget_set,"year","kind","item"); b1.add_argument("--amount",type=_pint,default=0)
    add("budget-status",budget_status,"year")
    add("budget-plan",budget_plan,"year")
    add("bulletin",bulletin,"sermon","notice","week","order","pray","date")
    add("lesson",lesson,"title","target","text","theme","goal","date")
    add("setlist",setlist,"songs","size","title","date")
    add("sermon",sermon,"title","text","theme","points","date","series","service")
    add("illus-add",illus_add,"name","source","topic","bible","text")
    add("illus-list",illus_list)
    add("commentary-add",commentary_add,"pdf","notebook","book")
    add("commentary-list",commentary_list)
    add("nlm-add",nlm_add,"file","notebook","book")
    add("song-add",song_add,"title","composer","lyricist","key","bpm","theme","bible","lyrics","sheet","type","youtube","release","album","copyright")
    add("song-list",song_list,"theme")
    add("song-open",song_open,"title")
    add("song-catalog",song_catalog,"title")
    add("songbook",songbook,"songs","title")
    add("compose",compose,"title","theme","bible","mood","key","bpm","lyrics")
    add("video-plan",video_plan,"title","purpose","length","tone","message","when","cta","aspect","mood")
    add("sermon-slides",sermon_slides,"title","text","theme","points")
    add("video-render",video_render,"title","message","when","cta","length","aspect","mood")
    add("lyrics-screen",lyrics_screen,"songs","title")
    add("worship-roster",worship_roster,"date","leader","keys","guitar","bass","drums","singers","songs")
    add("schedule-add",schedule_add,"date","type","place","host","tel","theme","fee","memo","upcoming")
    add("schedule-list",schedule_list,"upcoming")
    add("student-add",student_add,"name","school","grade","dept","tel","guardian","gtel","major","sex","birth","addr","leader","memo","date")
    add("student-list",student_list,"dept","grade")
    add("exam-cheer",exam_cheer,"exam","dept")
    add("production-add",production_add,"title","date","place","director","music","status","memo")
    add("production-list",production_list)
    add("casting",casting,"title","date","roles")
    add("mission-add",mission_add,"title","country","period","purpose")
    add("mission-member",mission_member,"title","name","role","expire","emergency","health")
    add("mission-roster",mission_roster,"title")
    add("mission-checklist",mission_checklist,"title")
    add("mission-plan",mission_plan,"title","dday")
    add("mission-place",mission_place,"title","kind","name","contact","address","memo")
    add("mission-field",mission_field,"title")
    add("mission-english",mission_english)
    add("mission-sermon",mission_sermon,"title","text","points")
    add("newfamily-add",newfamily_add,"name","tel","leader","cell","sex","birth","addr","memo","date")
    add("newfamily-step",newfamily_step,"name","how","by","note","date")
    add("newfamily-board",newfamily_board)
    add("newfamily-graduate",newfamily_graduate,"name","role")
    add("sermon-log",sermon_log,"date","title","text","service","series","points")
    add("sermon-list",sermon_list,"query","text","year","service")
    add("sermon-reuse",sermon_reuse,"query","text","title","theme","service","date")
    add("sermon-files",sermon_files,"open","service","query")
    add("open-file",open_file,"name","file","kind")
    add("my-files",my_files,"kind")
    add("open-folder",open_folder,"kind")
    add("song-sheet",song_sheet,"title")
    add("careevent-add",careevent_add,"kind","name","note","amount","action","date","jik","occasion","role")
    add("careevent-list",careevent_list,"kind")
    add("calendar",calendar_cmd,"year")
    add("birthday",birthday,"days")
    add("agenda",agenda,"days")
    add("hwp",hwp_convert,"file","recent")
    add("print-file",print_file,"file")
    add("ppt",ppt_lyrics,"songs","title")
    add("read-file",read_file,"file")
    add("giving-ledger",giving_ledger,"name","year")
    add("donation-receipt",donation_receipt,"name","year")
    add("vip-add",vip_add,"name","sponsor","tel","memo")
    add("vip-contact",vip_contact,"name","how","note","date")
    add("vip-list",vip_list)
    add("vip-convert",vip_convert,"name")
    add("presbytery-add",presbytery_add,"kind","note","role","due","date")
    add("presbytery-list",presbytery_list)
    add("event-plan",event_plan,"title","date","place","theme","host")
    add("version",version)
    add("backup",backup)
    ph=add("phoenix",phoenix,"restore"); ph.add_argument("--last",action="store_true")
    add("update",update,"file")
    add("update-check",update_check)
    add("setup",church_setup,"church","pastor")
    add("congrats",congrats,"occasion","kind","name","jik","role")
    add("sacrament-add",sacrament_add,"kind","name","birth","addr","date","by","text","memo")
    add("sacrament-ledger",sacrament_ledger,"kind")
    add("sacrament-apply",sacrament_apply,"kind","name","birth","addr","tel","cell","leader","mode","purpose")
    add("harvest-plan",harvest_plan,"title","dday")
    add("harvest-checklist",harvest_checklist,"title","date")
    add("cell-worship",cell_worship,"title","text","date","leader")
    add("member-transfer",member_transfer,"kind","name","date","church","memo")
    add("transfer-ledger",transfer_ledger)
    add("annual-plan",annual_plan,"year")
    add("event-add",event_add,"date","title","note","important")
    add("events-json",events_json,"month")
    add("event-del",event_del,"id")
    add("cal-print",cal_print,"month")
    add("dashboard",dashboard)
    add("request",request,"text")
    add("manual",manual)
    add("set-backup",set_backup,"path","to")
    add("sync",sync_out,"to")
    add("menu",menu)
    add("bible",bible,"ref","version")
    add("bible-plan",bible_plan,"type","date","dept")
    add("bible-quiz",bible_quiz,"ref","version","count","type","level","dept")
    add("memory-verse",memory_verse,"refs","set","version","dept")
    add("reference-list",reference_list)
    add("reference",reference,"name","hwp","part","toc")
    def _mk_ref(_short):
        def _f(a): setattr(a,"name",_short); reference(a)
        return _f
    for _s,_t,_fn,_rid in _REF_DOCS:
        add("ref-"+_rid,_mk_ref(_s),"part","toc","hwp")
    add("meeting-minutes",meeting_minutes,"type","date","time","place","host","clerk","attend","absent","agenda")
    add("official-doc",official_doc,"to","via","title","body","sender","date")
    add("asset-register",asset_register,"name","qty","place","manager","status","date")
    add("vehicle-log",vehicle_log,"driver","dest","km","fuel","date")
    add("vbs",vbs,"no","part","toc","hwp")
    add("devotion",devotion,"title","text","verse","body","apply","pray","theme","date")
    add("care",care,"days")
    add("newcomer",newcomer,"days")
    add("kakao",kakao,"type","to","title","text","ahead")
    add("weekly-brief",weekly_brief)
    add("prayer-add",prayer_add,"name","text","theme","date")
    add("prayer-digest",prayer_digest)
    add("giving-insight",giving_insight,"days")
    add("export-excel",export_excel)
    a=p.parse_args()
    if not getattr(a,"cmd",None): helpcmd(a); return
    try:
        a.func(a)
    except PermissionError as e:
        fn=os.path.basename(str(getattr(e,'filename','') or '')) or "해당 파일"
        print(f"⚠ '{fn}' 파일이 열려 있어 저장할 수 없습니다.\n   그 문서를 Word·한글에서 닫은 뒤 다시 실행해 주세요. (자료는 안전합니다)")
    except KeyboardInterrupt:
        print("\n중단되었습니다.")

if __name__=="__main__": main()
